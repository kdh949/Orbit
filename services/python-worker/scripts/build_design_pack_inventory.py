from __future__ import annotations

import argparse
import json
from collections import Counter
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Literal

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


LicenseStatus = Literal["approved", "pending", "rejected"]


@dataclass(frozen=True)
class ReferenceSource:
    source_id: str
    path: Path
    license_status: LicenseStatus = "pending"


@dataclass(frozen=True)
class SlideInventory:
    source: str
    slideNumber: int
    role: str
    silhouette: str
    koreanCapacity: str
    fontFamilies: list[str]
    licenseStatus: LicenseStatus
    eligibleForActivePack: bool


def build_inventory(sources: list[ReferenceSource]) -> dict[str, object]:
    slides = [
        inventory
        for source in sources
        for inventory in inspect_reference(source)
    ]
    source_counts = Counter(slide.source for slide in slides)
    return {
        "schemaVersion": 1,
        "sourceCount": len(sources),
        "slideCount": len(slides),
        "sources": [
            {
                "source": source.source_id,
                "slideCount": source_counts[source.source_id],
                "licenseStatus": source.license_status,
            }
            for source in sources
        ],
        "slides": [asdict(slide) for slide in slides],
    }


def inspect_reference(source: ReferenceSource) -> list[SlideInventory]:
    presentation = Presentation(source.path)
    last_slide = len(presentation.slides)
    return [
        inspect_slide(source, slide, index, last_slide)
        for index, slide in enumerate(presentation.slides, start=1)
    ]


def inspect_slide(
    source: ReferenceSource,
    slide: object,
    slide_number: int,
    last_slide: int,
) -> SlideInventory:
    shapes = list(slide.shapes)  # type: ignore[attr-defined]
    texts = [
        shape.text.strip()
        for shape in shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    text_chars = sum(len(text) for text in texts)
    picture_count = sum(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in shapes
    )
    table_count = sum(getattr(shape, "has_table", False) for shape in shapes)
    chart_count = sum(getattr(shape, "has_chart", False) for shape in shapes)
    role = infer_role(
        slide_number,
        last_slide,
        text_chars,
        picture_count,
        table_count,
        chart_count,
    )
    return SlideInventory(
        source=source.source_id,
        slideNumber=slide_number,
        role=role,
        silhouette=infer_silhouette(role, len(texts), picture_count),
        koreanCapacity=infer_korean_capacity(text_chars, len(texts)),
        fontFamilies=sorted(extract_fonts(shapes)),
        licenseStatus=source.license_status,
        eligibleForActivePack=source.license_status == "approved",
    )


def infer_role(
    slide_number: int,
    last_slide: int,
    text_chars: int,
    picture_count: int,
    table_count: int,
    chart_count: int,
) -> str:
    if slide_number == 1:
        return "cover"
    if slide_number == last_slide:
        return "closing"
    if chart_count or table_count:
        return "data"
    if text_chars <= 80:
        return "section"
    if picture_count:
        return "media-content"
    return "content"


def infer_silhouette(role: str, text_count: int, picture_count: int) -> str:
    if role in {"cover", "closing", "section"}:
        return f"{role}-focal"
    if picture_count:
        return "media-split" if text_count > 1 else "media-full"
    if role == "data":
        return "data-evidence"
    return "title-body-wide" if text_count <= 4 else "title-multi-item"


def infer_korean_capacity(text_chars: int, text_count: int) -> str:
    if text_chars <= 120 and text_count <= 3:
        return "high"
    if text_chars <= 360 and text_count <= 8:
        return "medium"
    return "low"


def extract_fonts(shapes: list[object]) -> set[str]:
    fonts: set[str] = set()
    for shape in shapes:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name:
                    fonts.add(run.font.name)
    return fonts


def parse_source(value: str) -> ReferenceSource:
    source_id, separator, path = value.partition("=")
    if not separator or not source_id or not path:
        raise argparse.ArgumentTypeError("source must use source-id=/path/file.pptx")
    return ReferenceSource(source_id=source_id, path=Path(path))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--source",
        action="append",
        required=True,
        type=parse_source,
        help="Stable source ID and local read-only PPTX path",
    )
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()
    report = build_inventory(args.source)
    args.output.write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
