from __future__ import annotations

import hashlib
import json
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation

from scripts import run_ooxml_reference_fidelity_artifacts as runner


TEMPLATE_IDS = [
    "simple-light",
    "simple-dark",
    "operating-review",
    "business-review",
    "project-kickoff",
    "team-alignment",
    "market-trends-report",
]


def _package() -> bytes:
    presentation = Presentation()
    for _ in range(8):
        presentation.slides.add_slide(presentation.slide_layouts[5])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _png(color: str) -> bytes:
    output = BytesIO()
    Image.new("RGB", (160, 90), color).save(output, format="PNG")
    return output.getvalue()


def _write_plan(tmp_path: Path) -> Path:
    package = _package()
    templates: list[dict[str, object]] = []
    for template_id in TEMPLATE_IDS:
        root = tmp_path / "private-input" / template_id
        root.mkdir(parents=True)
        source = root / "source.pptx"
        source.write_bytes(package)
        generated = root / "generated.pptx"
        generated.write_bytes(package)
        manifest = root / "manifest.json"
        manifest.write_text(
            json.dumps({"templateId": template_id, "version": 1}),
            encoding="utf-8",
        )
        templates.append(
            {
                "templateId": template_id,
                "templateVersion": 1,
                "sourcePath": str(source),
                "manifestPath": str(manifest),
                "generatedPptxPath": str(generated),
                "sourceSlideIds": [f"source-{index}" for index in range(1, 9)],
            }
        )
    plan = tmp_path / "private-plan.json"
    plan.write_text(
        json.dumps({"schemaVersion": 1, "templates": templates}),
        encoding="utf-8",
    )
    return plan


def _fake_manifest(_source: Path, value: dict[str, object]) -> SimpleNamespace:
    return SimpleNamespace(
        template_id=value["templateId"],
        version=value["version"],
        provenance=SimpleNamespace(authorization_status="approved"),
        canvas=SimpleNamespace(
            aspect_ratio="16:9",
            width_emu=12_192_000,
            height_emu=6_858_000,
        ),
        source_slides=[
            SimpleNamespace(
                source_slide_id=f"source-{index}",
                source_slide_part=f"ppt/slides/slide{index}.xml",
                slots=[],
            )
            for index in range(1, 9)
        ],
    )


def _fake_render(
    package: bytes,
    _aspect_ratio: str,
) -> runner.LibreOfficeRenderResult:
    digest = hashlib.sha256(package).hexdigest()[:2]
    return runner.LibreOfficeRenderResult(
        status="passed",
        version="26.8.0",
        pngs=tuple(_png(f"#{digest}{index:02x}66") for index in range(8)),
        issue_codes=(),
    )


def test_runner_generates_bounded_review_artifacts_without_approving_fidelity(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)

    summary = runner.run_fidelity_artifacts(
        plan,
        tmp_path / "private-output",
        render_deck=_fake_render,
    )

    assert summary["status"] == "generated"
    assert summary["approvalStatus"] == "pending"
    assert summary["generatedTemplateCount"] == 7
    assert summary["slideCount"] == 56
    report_path = (
        tmp_path
        / "private-output/simple-light/v1/manifests/fidelity-report.json"
    )
    report = json.loads(report_path.read_text(encoding="utf-8"))
    assert report["status"] == "generated"
    assert report["approvalStatus"] == "pending"
    assert report["structuralStatus"] == "passed"
    assert "FONT_AVAILABILITY_VALIDATION_PENDING" in report["issueCodes"]
    assert "HUMAN_FIDELITY_REVIEW_PENDING" in report["issueCodes"]
    assert len(report["slides"]) == 8
    root = report_path.parents[1]
    assert len(list((root / "baseline").glob("source-slide-*.png"))) == 8
    assert len(list((root / "generated").glob("generated-slide-*.png"))) == 8
    assert len(list((root / "diff").glob("intended-slot-mask-slide-*.png"))) == 8
    assert len(list((root / "diff").glob("locked-overlay-slide-*.png"))) == 8
    assert len(list((root / "montage").glob("*.png"))) == 3
    assert (root / "manifests/package.json").is_file()
    assert (root / "manifests/font.json").is_file()
    serialized = "\n".join(
        path.read_text(encoding="utf-8")
        for path in (root / "manifests").glob("*.json")
    )
    assert str(tmp_path) not in serialized
    assert "source.pptx" not in serialized
    assert report["slides"][0]["intendedSlotMaskPixelCount"] == 0


def test_runner_fails_only_the_fixture_with_an_unknown_source_slide(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    value = json.loads(plan.read_text(encoding="utf-8"))
    value["templates"][0]["sourceSlideIds"][0] = "unknown-source"
    plan.write_text(json.dumps(value), encoding="utf-8")
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)

    summary = runner.run_fidelity_artifacts(
        plan,
        tmp_path / "private-output",
        render_deck=_fake_render,
    )

    assert summary["status"] == "failed"
    assert summary["generatedTemplateCount"] == 6
    assert "SOURCE_SLIDE_SEQUENCE_INVALID" in summary["issueCodes"]


def test_runner_rejects_private_artifact_output_inside_repository(
    tmp_path: Path,
) -> None:
    plan = _write_plan(tmp_path)

    with pytest.raises(
        runner.FidelityArtifactRunnerError,
        match="FIDELITY_ARTIFACT_OUTPUT_INVALID",
    ):
        runner.run_fidelity_artifacts(
            plan,
            runner.REPOSITORY_ROOT / "private-output",
            render_deck=_fake_render,
        )


def test_runner_fails_summary_on_structural_drift(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)
    monkeypatch.setattr(runner, "_locked_snapshot_drift", lambda *_args: (1, 0, 0))

    summary = runner.run_fidelity_artifacts(
        plan,
        tmp_path / "private-output",
        render_deck=_fake_render,
    )

    assert summary["status"] == "failed"
    assert summary["generatedTemplateCount"] == 0
    assert "OOXML_REFERENCE_FIDELITY_LOCKED_GEOMETRY_DRIFT" in summary["issueCodes"]


def test_runner_fails_package_warning_hard_gate(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)
    monkeypatch.setattr(
        runner,
        "validate_cloned_package",
        lambda _package: ["RELATIONSHIP_TARGET_MISSING"],
    )

    summary = runner.run_fidelity_artifacts(
        plan,
        tmp_path / "private-output",
        render_deck=_fake_render,
    )

    assert summary["status"] == "failed"
    assert "OOXML_REFERENCE_PACKAGE_VALIDATION_FAILED" in summary["issueCodes"]
    package_manifest = json.loads(
        (
            tmp_path
            / "private-output/simple-light/v1/manifests/package.json"
        ).read_text(encoding="utf-8")
    )
    assert package_manifest["generatedPackageWarningCodes"] == [
        "RELATIONSHIP_TARGET_MISSING"
    ]


def test_locked_diff_detects_single_blue_channel_step() -> None:
    source = _png("#000000")
    generated_image = Image.new("RGB", (160, 90), "#000000")
    generated_image.putpixel((0, 0), (0, 0, 1))
    output = BytesIO()
    generated_image.save(output, format="PNG")
    mask = BytesIO()
    Image.new("L", (160, 90), 0).save(mask, format="PNG")

    _overlay, pixel_count = runner._locked_diff(
        source,
        output.getvalue(),
        mask.getvalue(),
    )

    assert pixel_count == 1


def test_runner_rejects_nonempty_output_directory(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    output = tmp_path / "private-output"
    output.mkdir()
    (output / "stale.json").write_text("{}", encoding="utf-8")
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)

    with pytest.raises(
        runner.FidelityArtifactRunnerError,
        match="FIDELITY_ARTIFACT_OUTPUT_NOT_EMPTY",
    ):
        runner.run_fidelity_artifacts(plan, output, render_deck=_fake_render)
