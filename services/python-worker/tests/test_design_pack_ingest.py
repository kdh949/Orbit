from pathlib import Path

from pptx import Presentation

from scripts.build_design_pack_inventory import ReferenceSource, build_inventory


def test_inventory_records_bounded_metadata_without_source_paths(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "reference.pptx"
    presentation = Presentation()
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = "Orbit Native Cover"
    content = presentation.slides.add_slide(presentation.slide_layouts[1])
    content.shapes.title.text = "핵심 내용"
    content.placeholders[1].text = "한글 본문 용량을 확인하는 설명입니다."
    presentation.save(source_path)

    report = build_inventory(
        [ReferenceSource("simple-light", source_path, "pending")]
    )

    assert report["sourceCount"] == 1
    assert report["slideCount"] == 2
    slides = report["slides"]
    assert isinstance(slides, list)
    assert slides[0]["role"] == "cover"
    assert slides[0]["licenseStatus"] == "pending"
    assert slides[0]["eligibleForActivePack"] is False
    assert str(tmp_path) not in str(report)


def test_approved_source_is_the_only_active_pack_eligible_status(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "approved.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    presentation.save(source_path)

    report = build_inventory(
        [ReferenceSource("orbit-native", source_path, "approved")]
    )

    slides = report["slides"]
    assert isinstance(slides, list)
    assert slides[0]["eligibleForActivePack"] is True
