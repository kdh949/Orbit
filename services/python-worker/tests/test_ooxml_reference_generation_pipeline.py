from __future__ import annotations

import hashlib
import io
import json
from pathlib import Path
from typing import Any

import pytest
from PIL import Image
from pptx import Presentation

from app.ai.ooxml_reference_templates.content_adapter import (
    ReferenceContentPlan,
)
from app.ai.ooxml_reference_templates.generation_stage import (
    GENERATION_STAGE_ORDER,
    OoxmlReferenceGenerationStageRequest,
    OoxmlReferenceStageError,
    execute_ooxml_reference_generation_stage,
)
from app.ai.ooxml_reference_templates.generation_runtime import (
    GeneratedAsset,
    LoadedReferenceTemplate,
    ReferenceInput,
    RenderValidationInput,
)
from app.ai.ooxml_reference_templates.models import (
    OoxmlReferenceTemplateGenerationRequest,
    OoxmlReferenceTemplateManifest,
)


def test_all_stages_execute_real_clone_slot_fidelity_and_materialization_path(
    tmp_path: Path,
) -> None:
    runtime = _FakeRuntime(tmp_path)
    dependencies: list[dict[str, object]] = []
    responses = []

    for stage in GENERATION_STAGE_ORDER:
        request = OoxmlReferenceGenerationStageRequest.model_validate(
            {
                **_base_request(),
                "stage": stage,
                "dependencies": dependencies,
            }
        )
        response = execute_ooxml_reference_generation_stage(
            request,
            runtime=runtime,
        )
        responses.append(response)
        dependencies.append(
            {
                "stage": stage,
                "payload": {
                    "data": response.artifact,
                    "metrics": {
                        "sourceSlideCount": response.source_slide_count,
                        "slotCount": response.slot_count,
                    },
                    "issueCodes": response.issue_codes,
                },
            }
        )

    by_stage = {response.stage: response for response in responses}
    package_artifact = by_stage["package-generation"].artifact
    assert package_artifact["currentPackage"] == {
        "fileId": "file-current-job-1",
        "originalName": "operating-review-generated.pptx",
        "size": runtime.current_package_size,
    }
    assert package_artifact["packageSha256"] == runtime.current_package_sha256
    assert runtime.generated_slide_text() == [
        ["운영 리뷰", "핵심 의사결정"],
        ["다음 단계", "실행 항목 확정"],
    ]
    assert (
        by_stage["render-validation"].artifact["fidelityReport"]["status"] == "passed"
    )

    materialization = by_stage["materialization"].artifact
    assert materialization["deck"]["metadata"] == {
        "language": "ko",
        "locale": "ko-KR",
        "sourceType": "import",
        "generatedBy": "ai",
        "audience": "general",
        "purpose": "inform",
        "tone": "professional",
        "createdFrom": {
            "topic": "운영 리뷰",
            "references": [{"fileId": "reference-1"}],
            "designReferences": [],
        },
    }
    assert materialization["templateSnapshot"] == {
        "catalogTemplateId": "operating-review",
        "catalogTemplateVersion": 1,
        "sourceSha256": runtime.source_sha256,
        "sourceSlideIds": ["cover-01", "closing-02"],
        "slotAssignmentCount": 4,
    }
    assert materialization["jobResult"]["fidelityReport"]["status"] == "passed"
    assert materialization["baselinePackage"] == {
        "fileId": "file-baseline",
        "originalName": "operating-review-generated-baseline.pptx",
        "size": runtime.current_package_size,
    }
    assert runtime.baseline_package == runtime.current_package
    assert runtime.baseline_package != runtime.source_bytes
    assert by_stage["content-planning"].artifact["outline"] == [
        {"order": 1, "title": "운영 리뷰"},
        {"order": 2, "title": "다음 단계"},
    ]
    assert materialization["renderAssets"] == [
        {
            "fileId": "file-render-job-1-001",
            "originalName": "slide-001.png",
            "size": runtime.render_png_size,
        },
        {
            "fileId": "file-render-job-1-002",
            "originalName": "slide-002.png",
            "size": runtime.render_png_size,
        },
    ]

    encoded = json.dumps(
        [response.model_dump(by_alias=True, mode="json") for response in responses],
        ensure_ascii=False,
    )
    assert all(
        len(
            json.dumps(
                response.artifact,
                ensure_ascii=False,
                separators=(",", ":"),
            ).encode("utf-8")
        )
        <= 1_048_576
        for response in responses
    )
    assert not any(
        forbidden in encoded.casefold()
        for forbidden in (
            "storagekey",
            "signedurl",
            "rawpackage",
            "packagebase64",
        )
    )
    planning = json.dumps(
        [
            by_stage["content-planning"].artifact,
            by_stage["template-planning"].artifact,
        ],
        ensure_ascii=False,
    )
    assert not any(
        forbidden in planning for forbidden in ('"x"', '"y"', "zIndex", "geometry")
    )


def test_same_inputs_produce_the_same_planning_and_package_checksums(
    tmp_path: Path,
) -> None:
    source_bytes = _source_package(tmp_path / "source")
    first = _execute_through_package(
        _FakeRuntime(tmp_path / "first", source_bytes=source_bytes)
    )
    second = _execute_through_package(
        _FakeRuntime(tmp_path / "second", source_bytes=source_bytes)
    )

    assert first["content-planning"] == second["content-planning"]
    assert first["template-planning"] == second["template-planning"]
    assert (
        first["package-generation"]["packageSha256"]
        == second["package-generation"]["packageSha256"]
    )


def test_runtime_cannot_return_a_private_locator_in_a_stage_artifact(
    tmp_path: Path,
) -> None:
    runtime = _FakeRuntime(tmp_path)
    runtime.reference_inputs = [
        ReferenceInput(
            file_id="reference-1",
            content="검증된 근거",
            metadata={"storageKey": "private/reference-1"},
        )
    ]
    request = OoxmlReferenceGenerationStageRequest.model_validate(_base_request())

    with pytest.raises(OoxmlReferenceStageError) as caught:
        execute_ooxml_reference_generation_stage(request, runtime=runtime)

    assert caught.value.code == "OOXML_REFERENCE_STAGE_ARTIFACT_INVALID"
    assert caught.value.retryable is False


def test_package_generation_recomputes_and_rejects_a_tampered_template_plan(
    tmp_path: Path,
) -> None:
    runtime = _FakeRuntime(tmp_path)
    dependencies: list[dict[str, object]] = []
    for stage in GENERATION_STAGE_ORDER[:4]:
        response = execute_ooxml_reference_generation_stage(
            OoxmlReferenceGenerationStageRequest.model_validate(
                {
                    **_base_request(),
                    "stage": stage,
                    "dependencies": dependencies,
                }
            ),
            runtime=runtime,
        )
        dependencies.append(
            {
                "stage": stage,
                "payload": {
                    "data": response.artifact,
                    "metrics": {
                        "sourceSlideCount": response.source_slide_count,
                        "slotCount": response.slot_count,
                    },
                    "issueCodes": response.issue_codes,
                },
            }
        )

    template_data = dependencies[-1]["payload"]
    assert isinstance(template_data, dict)
    template_plan = template_data["data"]
    assert isinstance(template_plan, dict)
    plan = template_plan["templatePlan"]
    assert isinstance(plan, dict)
    slides = plan["slides"]
    assert isinstance(slides, list)
    slides[0]["sourceSlideId"] = "closing-02"

    request = OoxmlReferenceGenerationStageRequest.model_validate(
        {
            **_base_request(),
            "stage": "package-generation",
            "dependencies": dependencies,
        }
    )
    with pytest.raises(OoxmlReferenceStageError) as caught:
        execute_ooxml_reference_generation_stage(request, runtime=runtime)

    assert caught.value.code == "OOXML_REFERENCE_TEMPLATE_PLAN_DRIFT"
    assert runtime._packages == {}


def _execute_through_package(runtime: _FakeRuntime) -> dict[str, dict[str, Any]]:
    dependencies: list[dict[str, object]] = []
    artifacts: dict[str, dict[str, Any]] = {}
    for stage in GENERATION_STAGE_ORDER[:5]:
        response = execute_ooxml_reference_generation_stage(
            OoxmlReferenceGenerationStageRequest.model_validate(
                {
                    **_base_request(),
                    "stage": stage,
                    "dependencies": dependencies,
                }
            ),
            runtime=runtime,
        )
        artifacts[stage] = response.artifact
        dependencies.append(
            {
                "stage": stage,
                "payload": {
                    "data": response.artifact,
                    "metrics": {
                        "sourceSlideCount": response.source_slide_count,
                        "slotCount": response.slot_count,
                    },
                    "issueCodes": response.issue_codes,
                },
            }
        )
    return artifacts


def _base_request() -> dict[str, object]:
    return {
        "jobId": "job-1",
        "projectId": "project-1",
        "stage": "reference-extract-file",
        "templateId": "operating-review",
        "templateVersion": 1,
        "request": {
            "topic": "운영 리뷰",
            "slideCountRange": {"min": 2, "max": 2},
            "referencePolicy": "references-first",
            "referenceFileIds": ["reference-1"],
            "templateSelection": {
                "mode": "user",
                "templateId": "operating-review",
                "version": 1,
            },
        },
        "dependencies": [],
    }


class _FakeRuntime:
    def __init__(self, root: Path, *, source_bytes: bytes | None = None) -> None:
        root.mkdir(parents=True, exist_ok=True)
        self.source_bytes = source_bytes or _source_package(root)
        self.source_sha256 = hashlib.sha256(self.source_bytes).hexdigest()
        self.source_size = len(self.source_bytes)
        self.manifest = _manifest(self.source_sha256)
        self.reference_inputs = [
            ReferenceInput(
                file_id="reference-1",
                content="검증된 근거",
                metadata={},
            )
        ]
        self._packages: dict[str, bytes] = {}
        self._baseline_packages: dict[str, bytes] = {}
        self._render_assets: dict[str, tuple[GeneratedAsset, ...]] = {}
        self._render_png = _png()

    @property
    def current_package_size(self) -> int:
        return len(self._packages["file-current-job-1"])

    @property
    def current_package_sha256(self) -> str:
        return hashlib.sha256(self._packages["file-current-job-1"]).hexdigest()

    @property
    def current_package(self) -> bytes:
        return self._packages["file-current-job-1"]

    @property
    def baseline_package(self) -> bytes:
        return self._baseline_packages["file-baseline"]

    @property
    def render_png_size(self) -> int:
        return len(self._render_png)

    def load_template(
        self, template_id: str, template_version: int
    ) -> LoadedReferenceTemplate:
        assert (template_id, template_version) == ("operating-review", 1)
        return LoadedReferenceTemplate(
            manifest=self.manifest,
            catalog_version="catalog-v1",
            source_package=self.source_bytes,
            source_asset=GeneratedAsset(
                file_id="file-baseline",
                original_name="operating-review-source.pptx",
                size=len(self.source_bytes),
            ),
        )

    def extract_references(
        self,
        project_id: str,
        request: OoxmlReferenceTemplateGenerationRequest,
    ) -> list[ReferenceInput]:
        assert project_id == "project-1"
        assert request.reference_file_ids == ["reference-1"]
        return self.reference_inputs

    def plan_content(
        self,
        project_id: str,
        request: OoxmlReferenceTemplateGenerationRequest,
        references: list[ReferenceInput],
    ) -> ReferenceContentPlan:
        assert project_id == "project-1"
        assert request.topic == "운영 리뷰"
        assert [item.file_id for item in references] == ["reference-1"]
        return ReferenceContentPlan.model_validate(
            {
                "templateId": "operating-review",
                "title": "운영 리뷰",
                "slides": [
                    {
                        "order": 1,
                        "semanticRole": "cover",
                        "values": [
                            {
                                "contentItemId": "cover-title",
                                "semanticRole": "title",
                                "content": "운영 리뷰",
                            },
                            {
                                "contentItemId": "cover-subtitle",
                                "semanticRole": "subtitle",
                                "content": "핵심 의사결정",
                            },
                        ],
                        "sourceRefs": ["reference-1"],
                    },
                    {
                        "order": 2,
                        "semanticRole": "closing",
                        "values": [
                            {
                                "contentItemId": "closing-title",
                                "semanticRole": "title",
                                "content": "다음 단계",
                            },
                            {
                                "contentItemId": "closing-subtitle",
                                "semanticRole": "subtitle",
                                "content": "실행 항목 확정",
                            },
                        ],
                    },
                ],
            }
        )

    def available_fonts(self) -> set[str]:
        return set()

    def font_fallbacks(self) -> dict[str, str]:
        return {}

    def read_image_asset(self, project_id: str, file_id: str) -> tuple[bytes, str]:
        raise AssertionError((project_id, file_id))

    def store_current_package(
        self,
        job_id: str,
        project_id: str,
        template_id: str,
        content: bytes,
    ) -> GeneratedAsset:
        assert project_id == "project-1"
        file_id = f"file-current-{job_id}"
        self._packages[file_id] = content
        return GeneratedAsset(
            file_id=file_id,
            original_name=f"{template_id}-generated.pptx",
            size=len(content),
        )

    def stage_baseline_package(
        self,
        job_id: str,
        project_id: str,
        template_id: str,
        content: bytes,
    ) -> GeneratedAsset:
        assert (job_id, project_id, template_id) == (
            "job-1",
            "project-1",
            "operating-review",
        )
        self._baseline_packages["file-baseline"] = content
        return GeneratedAsset(
            file_id="file-baseline",
            original_name="operating-review-generated-baseline.pptx",
            size=len(content),
        )

    def read_current_package(
        self,
        job_id: str,
        project_id: str,
        template_id: str,
        file_id: str,
    ) -> bytes:
        assert (job_id, project_id) == ("job-1", "project-1")
        assert template_id == "operating-review"
        return self._packages[file_id]

    def render_and_prepare_fidelity(
        self,
        job_id: str,
        project_id: str,
        template_id: str,
        template_version: int,
        package_file_id: str,
        package_bytes: bytes,
        source_slide_ids: list[str],
    ) -> RenderValidationInput:
        assert (template_id, template_version) == ("operating-review", 1)
        assert package_bytes == self._packages[package_file_id]
        assets = tuple(
            GeneratedAsset(
                file_id=f"file-render-{job_id}-{order:03d}",
                original_name=f"slide-{order:03d}.png",
                size=len(self._render_png),
            )
            for order in range(1, len(source_slide_ids) + 1)
        )
        self._render_assets[package_file_id] = assets
        slides = [
            {
                "sourceSlideId": source_slide_id,
                "sourcePng": self._render_png,
                "generatedPng": self._render_png,
                "intendedSlotMaskPng": self._render_png,
                "sourceLockedSnapshot": {"shapes": []},
                "generatedLockedSnapshot": {"shapes": []},
            }
            for source_slide_id in source_slide_ids
        ]
        return RenderValidationInput(
            assets=assets,
            slides=slides,
            environment={
                "renderer": "fake-powerpoint",
                "rendererVersion": "1",
                "fontFiles": [{"family": "Arial", "role": "body", "sha256": "f" * 64}],
                "sourceSha256": self.source_sha256,
                "templateManifestSha256": "e" * 64,
                "artifactSha256": hashlib.sha256(package_bytes).hexdigest(),
            },
            calibration={
                "status": "calibrated",
                "lockedRegionSsimThreshold": 0.99,
                "geometryEdgeTolerancePx": 0,
                "rationale": "deterministic fixture",
                "identityBaselines": [
                    {
                        "templateId": template_id,
                        "version": 1,
                        "renderer": "fake-powerpoint",
                        "rendererVersion": "1",
                        "reportSha256": "d" * 64,
                    }
                    for template_id in (
                        "simple-light",
                        "simple-dark",
                        "operating-review",
                        "business-review",
                        "project-kickoff",
                        "team-alignment",
                        "market-trends-report",
                    )
                ],
            },
        )

    def render_assets(
        self,
        job_id: str,
        project_id: str,
        package_file_id: str,
        assets: list[GeneratedAsset],
    ) -> tuple[GeneratedAsset, ...]:
        assert (job_id, project_id) == ("job-1", "project-1")
        assert assets == list(self._render_assets[package_file_id])
        return self._render_assets[package_file_id]

    def generated_slide_text(self) -> list[list[str]]:
        presentation = Presentation(io.BytesIO(self._packages["file-current-job-1"]))
        return [
            [shape.text for shape in slide.shapes if hasattr(shape, "text_frame")]
            for slide in presentation.slides
        ]


def _source_package(root: Path) -> bytes:
    root.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = "Source cover"
    cover.placeholders[1].text = "Source subtitle"
    closing = presentation.slides.add_slide(presentation.slide_layouts[0])
    closing.shapes.title.text = "Source closing"
    closing.placeholders[1].text = "Source next step"
    path = root / "source.pptx"
    presentation.save(path)
    return path.read_bytes()


def _manifest(source_sha256: str) -> OoxmlReferenceTemplateManifest:
    def slide(order: int, role: str, layout: int) -> dict[str, Any]:
        part = f"ppt/slides/slide{order}.xml"
        return {
            "sourceSlideId": f"{role}-{order:02d}",
            "sourceSlidePart": part,
            "sourceOrder": order,
            "semanticRole": role,
            "relationships": {
                "layoutPart": f"ppt/slideLayouts/slideLayout{layout}.xml",
                "masterPart": "ppt/slideMasters/slideMaster1.xml",
                "themePart": "ppt/theme/theme1.xml",
            },
            "capacity": {
                "textSlotCount": 2,
                "imageSlotCount": 0,
                "tableSlotCount": 0,
                "chartSlotCount": 0,
            },
            "previewId": f"preview-{order}",
            "lockedInventorySha256": f"{order:064x}",
            "slots": [
                _text_slot(part, order, "title", "2"),
                _text_slot(part, order, "subtitle", "3"),
            ],
        }

    return OoxmlReferenceTemplateManifest.model_validate(
        {
            "templateId": "operating-review",
            "version": 1,
            "status": "active",
            "sourceFormat": "pptx",
            "sourceSha256": source_sha256,
            "slideCount": 2,
            "canvas": {
                "aspectRatio": "16:9",
                "widthEmu": 12_192_000,
                "heightEmu": 6_858_000,
            },
            "name": "Operating Review",
            "description": "운영 리뷰 fixture",
            "preview": {
                "coverPreviewId": "cover",
                "coverPreviewSha256": "a" * 64,
                "bodyPreviewId": "body",
                "bodyPreviewSha256": "b" * 64,
            },
            "sourceSlides": [slide(1, "cover", 1), slide(2, "closing", 2)],
            "provenance": {
                "authorizationStatus": "approved",
                "inventoryVersion": 1,
            },
        }
    )


def _text_slot(
    slide_part: str, order: int, semantic_role: str, shape_id: str
) -> dict[str, Any]:
    return {
        "slotId": f"operating-review-v1-{order:02d}-{semantic_role}",
        "semanticRole": semantic_role,
        "contentType": "text",
        "required": True,
        "locator": {
            "slidePart": slide_part,
            "shapeId": shape_id,
            "placeholderType": semantic_role,
            "relationshipId": None,
        },
        "capacity": {"maxChars": 120, "maxLines": 4},
        "mutationPolicy": ["text-content"],
        "replacementPolicy": {"overflow": "fail"},
    }


def _png() -> bytes:
    output = io.BytesIO()
    Image.new("L", (16, 16), 0).save(output, format="PNG")
    return output.getvalue()
