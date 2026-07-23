from __future__ import annotations

import copy
import hashlib
import json
import subprocess
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from app.ai.ooxml_reference_templates.annotation import (
    AnnotationValidationError,
    build_image_slot_candidate_report,
    build_spike_candidate,
    build_source_slide_catalog,
    locked_inventory_sha256,
    render_source_slide_montage,
    select_spike_template,
    validate_source_slide_annotations,
)
from app.ai.ooxml_reference_templates.inventory import (
    ReferenceSource,
    inspect_reference_package,
)


SHA256 = "a" * 64
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
ROLES = [
    "cover",
    "agenda",
    "statement",
    "metric",
    "comparison",
    "process",
    "timeline",
    "evidence",
    "summary",
    "closing",
]


def _fixture_pptx(tmp_path: Path, name: str, slide_count: int) -> tuple[Path, list[int]]:
    presentation = Presentation()
    decoration_ids: list[int] = []
    for index in range(slide_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Fixture {index + 1}"
        slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(7), Inches(2)
        ).text_frame.text = "Body"
        decoration = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5)
        )
        decoration_ids.append(decoration.shape_id)
    path = tmp_path / name
    presentation.save(path)
    return path, decoration_ids


def _manifest(path: Path, template_id: str, slide_count: int, annotated: int) -> dict:
    presentation = Presentation(path)
    slides: list[dict] = []
    for index, slide in enumerate(presentation.slides):
        part = str(slide.part.partname).lstrip("/")
        title = slide.shapes.title
        body = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.shape_id != title.shape_id
        )
        role = ROLES[index] if index < len(ROLES) else "statement"
        slots = []
        if index < annotated:
            slots.extend(
                [
                {
                    "slotId": f"{template_id}-v1-slide-{index + 1:02d}-title",
                    "semanticRole": "title" if index in {0, slide_count - 1} else "body",
                    "contentType": "text",
                    "required": True,
                    "locator": {
                        "slidePart": part,
                        "shapeId": str(title.shape_id),
                        "placeholderType": "title",
                        "relationshipId": None,
                    },
                    "capacity": {"maxChars": 160, "maxLines": 4},
                    "mutationPolicy": ["text-content"],
                    "replacementPolicy": {"overflow": "fail"},
                },
                {
                    "slotId": f"{template_id}-v1-slide-{index + 1:02d}-body",
                    "semanticRole": "body",
                    "contentType": "text",
                    "required": False,
                    "locator": {
                        "slidePart": part,
                        "shapeId": str(body.shape_id),
                        "placeholderType": None,
                        "relationshipId": None,
                    },
                    "capacity": {"maxChars": 500, "maxLines": 12},
                    "mutationPolicy": ["text-content"],
                    "replacementPolicy": {"overflow": "fail"},
                },
                ]
            )
        slides.append(
            {
                "sourceSlideId": f"{template_id}-slide-{index + 1:02d}",
                "sourceSlidePart": part,
                "sourceOrder": index + 1,
                "semanticRole": role,
                "relationships": {
                    "layoutPart": "ppt/slideLayouts/slideLayout6.xml",
                    "masterPart": "ppt/slideMasters/slideMaster1.xml",
                    "themePart": "ppt/theme/theme1.xml",
                },
                "capacity": {
                    "textSlotCount": len(slots),
                    "imageSlotCount": 0,
                    "tableSlotCount": 0,
                    "chartSlotCount": 0,
                },
                "previewId": f"slide-{index + 1:02d}",
                "lockedInventorySha256": SHA256,
                "slots": slots,
            }
        )
    for slide in slides:
        slide["lockedInventorySha256"] = locked_inventory_sha256(
            path,
            slide["sourceSlidePart"],
            (slot["locator"]["shapeId"] for slot in slide["slots"]),
            (slot["locator"]["relationshipId"] for slot in slide["slots"]),
        )
    return {
        "templateId": template_id,
        "version": 1,
        "status": "active",
        "sourceFormat": "pptx",
        "sourceSha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        "slideCount": slide_count,
        "canvas": {
            "aspectRatio": "16:9",
            "widthEmu": 12_192_000,
            "heightEmu": 6_858_000,
        },
        "name": template_id,
        "description": "annotation fixture",
        "preview": {
            "coverPreviewId": "slide-01",
            "coverPreviewSha256": SHA256,
            "bodyPreviewId": "slide-02",
            "bodyPreviewSha256": SHA256,
        },
        "sourceSlides": slides,
        "provenance": {"authorizationStatus": "approved", "inventoryVersion": 1},
    }


def _image_annotation_fixture(
    tmp_path: Path, *, shared: bool
) -> tuple[Path, dict]:
    path, _ = _fixture_pptx(tmp_path, "image-annotation.pptx", 10)
    image_path = tmp_path / "annotation-image.png"
    Image.new("RGBA", (400, 300), (30, 90, 150, 180)).save(image_path)
    presentation = Presentation(path)
    slide = presentation.slides[0]
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(4), Inches(3), Inches(2.25)
    )
    if shared:
        slide.shapes.add_picture(
            str(image_path), Inches(5), Inches(4), Inches(3), Inches(2.25)
        )
    shape_id = str(picture.shape_id)
    relationship_id = picture._pic.blipFill.blip.rEmbed
    presentation.save(path)

    value = _manifest(path, "image-annotation", 10, annotated=10)
    source_slide = value["sourceSlides"][0]
    source_slide["slots"].append(
        {
            "slotId": "image-annotation-v1-slide-01-image",
            "semanticRole": "image",
            "contentType": "image",
            "required": True,
            "locator": {
                "slidePart": source_slide["sourceSlidePart"],
                "shapeId": shape_id,
                "placeholderType": None,
                "relationshipId": relationship_id,
            },
            "capacity": {
                "minAspectRatio": 1.2,
                "maxAspectRatio": 1.5,
                "cropPolicy": "preserve-frame",
                "alphaRequired": True,
                "maskRequired": False,
            },
            "mutationPolicy": ["image-source"],
            "replacementPolicy": {"overflow": "fail"},
        }
    )
    source_slide["capacity"]["imageSlotCount"] = 1
    source_slide["lockedInventorySha256"] = locked_inventory_sha256(
        path,
        source_slide["sourceSlidePart"],
        (slot["locator"]["shapeId"] for slot in source_slide["slots"]),
        (slot["locator"]["relationshipId"] for slot in source_slide["slots"]),
    )
    return path, value


def _picture_candidate_fixture(
    tmp_path: Path,
    *,
    shared: bool = False,
    placeholder: bool = True,
    animated: bool = False,
    replacement_description: str | None = None,
    replacement_name: bool = False,
) -> tuple[Path, dict]:
    path, _ = _fixture_pptx(tmp_path, "candidate-source.pptx", 10)
    image_path = tmp_path / "candidate-image.png"
    Image.new("RGBA", (400, 300), (30, 90, 150, 180)).save(image_path)
    presentation = Presentation(path)
    slide = presentation.slides[0]
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(4), Inches(3), Inches(2.25)
    )
    if shared:
        slide.shapes.add_picture(
            str(image_path), Inches(5), Inches(4), Inches(3), Inches(2.25)
        )
    shape_id = str(picture.shape_id)
    presentation.save(path)

    with zipfile.ZipFile(path, "r") as package:
        entries = {
            item.filename: package.read(item.filename) for item in package.infolist()
        }
    root = ET.fromstring(entries["ppt/slides/slide1.xml"])
    source_picture = next(
        item
        for item in root.findall(f".//{{{PML_NS}}}pic")
        if item.find(f".//{{{PML_NS}}}cNvPr").get("id") == shape_id
    )
    if placeholder:
        non_visual = source_picture.find(
            f"{{{PML_NS}}}nvPicPr/{{{PML_NS}}}nvPr"
        )
        ET.SubElement(non_visual, f"{{{PML_NS}}}ph", {"type": "pic"})
    c_nv_pr = source_picture.find(f".//{{{PML_NS}}}cNvPr")
    if replacement_description is not None:
        c_nv_pr.set("descr", replacement_description)
    if replacement_name:
        c_nv_pr.set("name", "Replace with image.")
    if animated:
        timing = ET.SubElement(root, f"{{{PML_NS}}}timing")
        ET.SubElement(timing, f"{{{PML_NS}}}spTgt", {"spid": shape_id})
    entries["ppt/slides/slide1.xml"] = ET.tostring(root, xml_declaration=True)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as package:
        for filename, content in entries.items():
            package.writestr(filename, content)

    return path, _manifest(path, "candidate-report", 10, annotated=10)


def _mark_smartart_and_animation(path: Path, animation_shape_id: str) -> None:
    with zipfile.ZipFile(path, "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}
    slide = entries["ppt/slides/slide1.xml"]
    smartart = (
        b'<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="99" name="SmartArt 1"/>'
        b'<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr><p:xfrm/>'
        b'<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/'
        b'drawingml/2006/diagram"/></a:graphic></p:graphicFrame>'
    )
    timing = (
        b'<p:timing><p:tnLst><p:par><p:cTn id="1" dur="indefinite">'
        b'<p:childTnLst><p:anim><p:cBhvr><p:cTn id="2" dur="500"/>'
        b'<p:tgtEl><p:spTgt spid="'
        + animation_shape_id.encode()
        + b'"/></p:tgtEl></p:cBhvr></p:anim></p:childTnLst></p:cTn>'
        b'</p:par></p:tnLst></p:timing>'
    )
    entries["ppt/slides/slide1.xml"] = slide.replace(
        b"</p:spTree>", smartart + b"</p:spTree>"
    ).replace(b"</p:sld>", timing + b"</p:sld>")
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as package:
        for filename, content in entries.items():
            package.writestr(filename, content)


def _candidate(path: Path, manifest: dict) -> dict:
    source = ReferenceSource(manifest["templateId"], path)
    inventory = inspect_reference_package(source)
    return build_spike_candidate(source=source, inventory=inventory, manifest=manifest)


def test_spike_selection_is_deterministic_and_prefers_supported_locator_role_capacity(
    tmp_path: Path,
) -> None:
    supported_path, _ = _fixture_pptx(tmp_path, "supported.pptx", 10)
    larger_path, _ = _fixture_pptx(tmp_path, "larger.pptx", 12)
    supported = _candidate(
        supported_path, _manifest(supported_path, "supported", 10, annotated=10)
    )
    larger = _candidate(larger_path, _manifest(larger_path, "larger", 12, annotated=4))

    first = select_spike_template([larger, supported], target_slide_count_range=(8, 10))
    second = select_spike_template([supported, larger], target_slide_count_range=(8, 10))

    assert first == second
    assert first["templateId"] == "supported"
    assert first["rationale"]["supportedLocatorCoverage"] == 1.0
    assert first["rationale"]["hasCover"] is True
    assert first["rationale"]["hasClosing"] is True
    assert first["rationale"]["roleCoverage"] > larger["roleCoverage"]
    assert first["rationale"]["capacityEligibleSlideCount"] >= 10
    assert first["rationale"]["targetSlideCounts"] == [8, 9, 10]
    assert "slideCount" not in first["rationale"]["rankingCriteria"]


@pytest.mark.parametrize("location", ["source-slide", "slot"])
def test_annotation_rejects_unknown_fields(tmp_path: Path, location: str) -> None:
    path, _ = _fixture_pptx(tmp_path, "strict.pptx", 10)
    value = _manifest(path, "strict", 10, annotated=10)
    target = value["sourceSlides"][0]
    if location == "slot":
        target = target["slots"][0]
    target["rawXml"] = "<private/>"

    with pytest.raises(AnnotationValidationError, match="unknown_field"):
        validate_source_slide_annotations(path, value)


@pytest.mark.parametrize("duplicate", ["slot-id", "locator"])
def test_annotation_rejects_duplicate_slot_identity(tmp_path: Path, duplicate: str) -> None:
    path, _ = _fixture_pptx(tmp_path, "duplicate.pptx", 10)
    value = _manifest(path, "duplicate", 10, annotated=10)
    copied = copy.deepcopy(value["sourceSlides"][0]["slots"][0])
    if duplicate == "locator":
        copied["slotId"] += "-other"
    else:
        copied["locator"] = {
            **copied["locator"],
            "shapeId": str(int(copied["locator"]["shapeId"]) + 1),
        }
    value["sourceSlides"][0]["slots"].append(copied)

    code = "duplicate_locator" if duplicate == "locator" else "duplicate_slot_id"
    with pytest.raises(AnnotationValidationError, match=code):
        validate_source_slide_annotations(path, value)


def test_annotation_rejects_locator_not_present_in_source_ooxml(tmp_path: Path) -> None:
    path, _ = _fixture_pptx(tmp_path, "missing.pptx", 10)
    value = _manifest(path, "missing", 10, annotated=10)
    value["sourceSlides"][0]["slots"][0]["locator"]["shapeId"] = "999999"

    with pytest.raises(AnnotationValidationError, match="unsupported_locator"):
        validate_source_slide_annotations(path, value)


def test_annotation_accepts_package_wide_unique_image_media_target(
    tmp_path: Path,
) -> None:
    path, value = _image_annotation_fixture(tmp_path, shared=False)

    manifest = validate_source_slide_annotations(path, value)

    image_slots = [
        slot
        for slot in manifest.source_slides[0].slots
        if slot.content_type == "image"
    ]
    assert len(image_slots) == 1


def test_annotation_rejects_shared_image_media_target(tmp_path: Path) -> None:
    path, value = _image_annotation_fixture(tmp_path, shared=True)

    with pytest.raises(
        AnnotationValidationError, match="shared_image_media_target"
    ):
        validate_source_slide_annotations(path, value)


def test_image_slot_candidate_report_uses_conservative_replacement_intent(
    tmp_path: Path,
) -> None:
    path, value = _picture_candidate_fixture(tmp_path)
    manifest = validate_source_slide_annotations(path, value)

    report = build_image_slot_candidate_report(path, manifest)

    assert report["summary"] == {
        "directPictureCount": 1,
        "eligibleCandidateCount": 1,
        "highConfidenceCandidateCount": 1,
        "excludedPictureCount": 0,
        "exclusionReasonCounts": {},
    }
    candidate = report["candidates"][0]
    assert candidate["replacementIntent"] == {
        "sourceType": "placeholder",
        "usage": "media-slot",
        "replaceMode": "replace",
        "confidence": 0.95,
        "evidence": "direct-picture-placeholder",
    }
    assert candidate["highConfidence"] is True
    assert candidate["mediaTargetRelationshipCount"] == 1
    assert candidate["slideEmbedCount"] == 1
    serialized = json.dumps(report)
    assert str(path) not in serialized
    assert "Fixture 1" not in serialized
    assert "ppt/media/" not in serialized
    assert "<p:pic" not in serialized


def test_image_slot_candidate_report_records_shared_target_exclusion(
    tmp_path: Path,
) -> None:
    path, value = _picture_candidate_fixture(tmp_path, shared=True)
    manifest = validate_source_slide_annotations(path, value)

    report = build_image_slot_candidate_report(path, manifest)

    assert report["summary"]["directPictureCount"] == 2
    assert report["summary"]["eligibleCandidateCount"] == 0
    assert report["summary"]["exclusionReasonCounts"] == {
        "shared_media_target": 2
    }
    assert all(
        exclusion["mediaTargetRelationshipCount"] == 1
        and exclusion["slideEmbedCount"] == 2
        and exclusion["exclusionReasons"] == ["shared_media_target"]
        and exclusion["highConfidence"] is False
        for exclusion in report["exclusions"]
    )


def test_image_slot_candidate_report_does_not_infer_high_confidence_without_placeholder(
    tmp_path: Path,
) -> None:
    path, value = _picture_candidate_fixture(
        tmp_path,
        placeholder=False,
        replacement_name=True,
    )
    manifest = validate_source_slide_annotations(path, value)

    report = build_image_slot_candidate_report(path, manifest)

    assert report["summary"]["eligibleCandidateCount"] == 1
    assert report["summary"]["highConfidenceCandidateCount"] == 0
    candidate = report["candidates"][0]
    assert candidate["highConfidence"] is False
    assert candidate["replacementIntent"] == {
        "sourceType": "slide",
        "usage": "media-slot",
        "replaceMode": "preserve",
        "confidence": 0.55,
        "evidence": "no-explicit-source-replacement-intent",
    }


def test_image_slot_candidate_report_accepts_exact_normalized_authored_description_without_leaking_it(
    tmp_path: Path,
) -> None:
    path, value = _picture_candidate_fixture(
        tmp_path,
        placeholder=False,
        replacement_description="  Replace   with image. \n",
    )
    manifest = validate_source_slide_annotations(path, value)

    report = build_image_slot_candidate_report(path, manifest)

    assert report["summary"]["highConfidenceCandidateCount"] == 1
    candidate = report["candidates"][0]
    assert candidate["highConfidence"] is True
    assert candidate["replacementIntent"] == {
        "sourceType": "slide",
        "usage": "media-slot",
        "replaceMode": "replace",
        "confidence": 0.95,
        "evidence": "source-authored-image-replacement-description",
    }
    serialized = json.dumps(report)
    assert "Replace with image." not in serialized
    assert "Replace   with image." not in serialized
    assert "description" not in candidate["replacementIntent"]


def test_image_slot_candidate_report_excludes_animated_picture(
    tmp_path: Path,
) -> None:
    path, value = _picture_candidate_fixture(tmp_path, animated=True)
    manifest = validate_source_slide_annotations(path, value)

    report = build_image_slot_candidate_report(path, manifest)

    assert report["summary"]["eligibleCandidateCount"] == 0
    assert report["summary"]["exclusionReasonCounts"] == {"animated_picture": 1}
    assert report["exclusions"][0]["highConfidence"] is False


def test_annotation_cli_writes_read_only_image_slot_candidate_report(
    tmp_path: Path,
) -> None:
    path, value = _picture_candidate_fixture(tmp_path)
    manifest_path = tmp_path / "manifest.json"
    manifest_path.write_text(json.dumps(value), encoding="utf-8")
    output_path = tmp_path / "candidate-report.json"
    script = (
        Path(__file__).parents[1]
        / "scripts"
        / "annotate_ooxml_reference_template.py"
    )

    result = subprocess.run(
        [
            sys.executable,
            str(script),
            "--source",
            str(path),
            "--manifest",
            str(manifest_path),
            "--image-slot-candidate-output",
            str(output_path),
            "--target-slide-count",
            "99",
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    report = json.loads(output_path.read_text(encoding="utf-8"))
    assert report["summary"]["highConfidenceCandidateCount"] == 1
    assert value == json.loads(manifest_path.read_text(encoding="utf-8"))


@pytest.mark.parametrize(
    ("excluded", "code"),
    [
        ("decoration", "excluded_decoration"),
        ("master", "excluded_master_layout"),
        ("layout", "excluded_master_layout"),
        ("smartart", "excluded_smartart"),
        ("animation", "excluded_animation"),
    ],
)
def test_annotation_excludes_non_slot_ooxml_objects(
    tmp_path: Path, excluded: str, code: str
) -> None:
    path, decoration_ids = _fixture_pptx(tmp_path, "excluded.pptx", 10)
    value = _manifest(path, "excluded", 10, annotated=10)
    slot = value["sourceSlides"][0]["slots"][0]
    if excluded == "decoration":
        slot["locator"]["shapeId"] = str(decoration_ids[0])
    elif excluded in {"master", "layout"}:
        slot["locator"]["slidePart"] = (
            "ppt/slideMasters/slideMaster1.xml"
            if excluded == "master"
            else "ppt/slideLayouts/slideLayout6.xml"
        )
    else:
        animation_shape_id = slot["locator"]["shapeId"]
        _mark_smartart_and_animation(path, animation_shape_id)
        value["sourceSha256"] = hashlib.sha256(path.read_bytes()).hexdigest()
        if excluded == "smartart":
            slot["locator"]["shapeId"] = "99"

    with pytest.raises(AnnotationValidationError, match=code):
        validate_source_slide_annotations(path, value)


def test_builds_bounded_human_review_catalog_and_montage(tmp_path: Path) -> None:
    path, _ = _fixture_pptx(tmp_path, "review.pptx", 10)
    manifest = validate_source_slide_annotations(
        path, _manifest(path, "review", 10, annotated=10)
    )
    catalog = build_source_slide_catalog(manifest, target_slide_count=10)
    preview_directory = tmp_path / "previews"
    preview_directory.mkdir()
    for slide in catalog["slides"]:
        Image.new("RGB", (160, 90), "white").save(
            preview_directory / f"{slide['previewId']}.png"
        )
    montage_path = tmp_path / "review-output" / "montage.png"

    render_source_slide_montage(catalog, preview_directory, montage_path)

    assert catalog["targetSlideCount"] == 10
    assert catalog["slides"][0]["semanticRole"] == "cover"
    assert catalog["slides"][-1]["semanticRole"] == "closing"
    assert all("sourceSlidePart" not in slide for slide in catalog["slides"])
    with Image.open(montage_path) as montage:
        assert montage.size == (1000, 1550)
