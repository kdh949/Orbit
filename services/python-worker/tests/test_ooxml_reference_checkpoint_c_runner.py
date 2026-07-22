from __future__ import annotations

import hashlib
import json
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation

from scripts import run_ooxml_reference_checkpoint_c as runner


TEMPLATE_IDS = [
    "simple-light",
    "simple-dark",
    "operating-review",
    "business-review",
    "project-kickoff",
    "team-alignment",
    "market-trends-report",
]


def _generated_package(slide_count: int = 8) -> bytes:
    presentation = Presentation()
    for _ in range(slide_count):
        presentation.slides.add_slide(presentation.slide_layouts[5])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _png(color: str) -> bytes:
    output = BytesIO()
    Image.new("RGB", (160, 90), color).save(output, format="PNG")
    return output.getvalue()


def _write_plan(
    tmp_path: Path,
    *,
    power_point: bool,
    overflow_template_id: str | None = None,
) -> Path:
    package = _generated_package()
    package_sha256 = hashlib.sha256(package).hexdigest()
    templates: list[dict[str, object]] = []
    for template_id in TEMPLATE_IDS:
        root = tmp_path / "private-inputs" / template_id
        root.mkdir(parents=True)
        source_path = root / "source.pptx"
        source_path.write_bytes(package)
        manifest_path = root / "manifest.json"
        manifest_path.write_text(
            json.dumps({"templateId": template_id, "version": 1}),
            encoding="utf-8",
        )
        generated_path = root / "generated.pptx"
        generated_path.write_bytes(package)
        quality_path = root / "quality.json"
        quality_path.write_text(
            json.dumps(
                {
                    "schemaVersion": 1,
                    "templateId": template_id,
                    "templateVersion": 1,
                    "generatedPackageSha256": package_sha256,
                    "generatedSlideCount": 8,
                    "requiredUniqueSourceCount": 7,
                    "selectedUniqueSourceCount": 8,
                    "eligibleUniqueSourceCount": 8,
                    "requiredUniqueLayoutCount": 4,
                    "selectedUniqueLayoutCount": 8,
                    "eligibleUniqueLayoutCount": 8,
                    "adjacentRepeatCount": 0,
                    "slotOverflowCount": int(template_id == overflow_template_id),
                    "overlapCount": 0,
                    "cropErrorCount": 0,
                    "warningCodes": [],
                    "unsupportedObjectCodes": [],
                    "fontSubstitutionRiskCodes": [],
                    "supportedFixtureTypes": ["text"],
                }
            ),
            encoding="utf-8",
        )
        entry: dict[str, object] = {
            "templateId": template_id,
            "templateVersion": 1,
            "sourcePath": str(source_path),
            "manifestPath": str(manifest_path),
            "generatedPptxPath": str(generated_path),
            "qualityEvidencePath": str(quality_path),
        }
        if power_point:
            evidence_path = root / "powerpoint.json"
            evidence_path.write_text(
                json.dumps(
                    {
                        "schemaVersion": 1,
                        "application": "Microsoft PowerPoint",
                        "applicationVersion": "16.0",
                        "templateId": template_id,
                        "templateVersion": 1,
                        "generatedPackageSha256": package_sha256,
                        "evaluatedSlideCount": 8,
                        "reopenStatus": "passed",
                        "renderStatus": "passed",
                        "renderArtifactSha256": "a" * 64,
                        "warningCodes": [],
                    }
                ),
                encoding="utf-8",
            )
            entry["powerPointEvidencePath"] = str(evidence_path)
        templates.append(entry)
    plan_path = tmp_path / "fixture-plan.json"
    plan_path.write_text(
        json.dumps({"schemaVersion": 1, "templates": templates}),
        encoding="utf-8",
    )
    return plan_path


def _fake_annotation_validator(
    _source_path: Path, manifest: dict[str, object]
) -> SimpleNamespace:
    return SimpleNamespace(
        template_id=manifest["templateId"],
        version=manifest["version"],
        canvas=SimpleNamespace(aspect_ratio="16:9"),
        provenance=SimpleNamespace(authorization_status="approved"),
    )


def _fake_libreoffice_render(
    _package: bytes, _aspect_ratio: str
) -> runner.LibreOfficeRenderResult:
    return runner.LibreOfficeRenderResult(
        status="passed",
        version="26.8.0",
        pngs=tuple(_png(f"#{index:02x}3366") for index in range(8)),
        issue_codes=(),
    )


def test_runner_writes_bounded_reports_but_never_passes_without_powerpoint(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan_path = _write_plan(tmp_path, power_point=False)
    monkeypatch.setattr(
        runner,
        "validate_source_slide_annotations",
        _fake_annotation_validator,
    )

    summary = runner.run_checkpoint_c(
        plan_path,
        tmp_path / "private-output",
        render_deck=_fake_libreoffice_render,
    )

    assert summary["status"] == "not-run"
    assert summary["templateCount"] == 7
    assert summary["automatedPassedTemplateCount"] == 7
    assert summary["powerPointPassedTemplateCount"] == 0
    assert summary["qualityTotals"] == {
        "slotOverflowCount": 0,
        "overlapCount": 0,
        "cropErrorCount": 0,
        "warningCount": 0,
    }
    report_path = tmp_path / "private-output/simple-light/report.json"
    report = json.loads(report_path.read_text(encoding="utf-8"))
    assert report["status"] == "not-run"
    assert report["packageValidation"] == {
        "status": "passed",
        "warningCodes": [],
    }
    assert report["pythonPptxReopen"] == {
        "status": "passed",
        "slideCount": 8,
    }
    assert report["libreOffice"]["status"] == "passed"
    assert report["libreOffice"]["renderedSlideCount"] == 8
    assert report["powerPoint"]["status"] == "not-run"
    assert len(list((tmp_path / "private-output/simple-light/libreoffice").glob("*.png"))) == 9
    serialized = json.dumps(summary) + report_path.read_text(encoding="utf-8")
    assert str(tmp_path) not in serialized
    assert "source.pptx" not in serialized


def test_runner_passes_only_with_bound_powerpoint_evidence_and_zero_quality_counts(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan_path = _write_plan(tmp_path, power_point=True)
    monkeypatch.setattr(
        runner,
        "validate_source_slide_annotations",
        _fake_annotation_validator,
    )

    summary = runner.run_checkpoint_c(
        plan_path,
        tmp_path / "private-output",
        render_deck=_fake_libreoffice_render,
    )

    assert summary["status"] == "passed"
    assert summary["automatedPassedTemplateCount"] == 7
    assert summary["powerPointPassedTemplateCount"] == 7
    assert summary["passedTemplateCount"] == 7


def test_runner_fails_checkpoint_on_any_overflow_even_when_renderers_pass(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan_path = _write_plan(
        tmp_path,
        power_point=True,
        overflow_template_id="operating-review",
    )
    monkeypatch.setattr(
        runner,
        "validate_source_slide_annotations",
        _fake_annotation_validator,
    )

    summary = runner.run_checkpoint_c(
        plan_path,
        tmp_path / "private-output",
        render_deck=_fake_libreoffice_render,
    )

    assert summary["status"] == "failed"
    assert summary["qualityTotals"]["slotOverflowCount"] == 1
    report = json.loads(
        (
            tmp_path / "private-output/operating-review/report.json"
        ).read_text(encoding="utf-8")
    )
    assert report["automatedStatus"] == "failed"
    assert "OOXML_REFERENCE_SLOT_OVERFLOW" in report["issueCodes"]


def test_runner_rejects_a_fixture_plan_without_exactly_seven_templates(
    tmp_path: Path,
) -> None:
    plan_path = _write_plan(tmp_path, power_point=False)
    value = json.loads(plan_path.read_text(encoding="utf-8"))
    value["templates"].pop()
    plan_path.write_text(json.dumps(value), encoding="utf-8")

    with pytest.raises(runner.CheckpointRunnerError, match="FIXTURE_PLAN_INVALID"):
        runner.run_checkpoint_c(plan_path, tmp_path / "private-output")


def test_runner_reports_sequence_counts_instead_of_hiding_them_as_invalid_evidence(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan_path = _write_plan(tmp_path, power_point=True)
    value = json.loads(plan_path.read_text(encoding="utf-8"))
    operating_review = next(
        item
        for item in value["templates"]
        if item["templateId"] == "operating-review"
    )
    quality_path = Path(operating_review["qualityEvidencePath"])
    quality = json.loads(quality_path.read_text(encoding="utf-8"))
    quality["selectedUniqueSourceCount"] = 5
    quality["adjacentRepeatCount"] = 1
    quality_path.write_text(json.dumps(quality), encoding="utf-8")
    monkeypatch.setattr(
        runner,
        "validate_source_slide_annotations",
        _fake_annotation_validator,
    )

    runner.run_checkpoint_c(
        plan_path,
        tmp_path / "private-output",
        render_deck=_fake_libreoffice_render,
    )

    report = json.loads(
        (
            tmp_path / "private-output/operating-review/report.json"
        ).read_text(encoding="utf-8")
    )
    assert report["qualityMetrics"]["requiredUniqueSourceCount"] == 7
    assert report["qualityMetrics"]["selectedUniqueSourceCount"] == 5
    assert report["qualityMetrics"]["adjacentRepeatCount"] == 1
    assert "OOXML_REFERENCE_SOURCE_UNIQUENESS_FAILED" in report["issueCodes"]
    assert "OOXML_REFERENCE_ADJACENT_REPEAT" in report["issueCodes"]


def test_runner_records_a_malformed_generated_package_without_crashing(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan_path = _write_plan(tmp_path, power_point=False)
    value = json.loads(plan_path.read_text(encoding="utf-8"))
    target = next(
        item for item in value["templates"] if item["templateId"] == "simple-light"
    )
    Path(target["generatedPptxPath"]).write_bytes(b"not-an-ooxml-package")
    monkeypatch.setattr(
        runner,
        "validate_source_slide_annotations",
        _fake_annotation_validator,
    )

    summary = runner.run_checkpoint_c(
        plan_path,
        tmp_path / "private-output",
        render_deck=_fake_libreoffice_render,
    )

    assert summary["status"] == "failed"
    report = json.loads(
        (tmp_path / "private-output/simple-light/report.json").read_text(
            encoding="utf-8"
        )
    )
    assert report["status"] == "failed"
    assert report["pythonPptxReopen"]["status"] == "failed"
    assert "MALFORMED_OOXML_PACKAGE" in report["packageValidation"]["warningCodes"]


def test_cli_masks_private_paths_when_artifact_io_fails(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    private_detail = str(tmp_path / "private-output/report.json")

    def fail_runner(_plan: Path, _output: Path) -> dict[str, object]:
        raise OSError(private_detail)

    monkeypatch.setattr(runner, "run_checkpoint_c", fail_runner)

    exit_code = runner.main(
        [
            "--plan",
            str(tmp_path / "private-plan.json"),
            "--output-directory",
            str(tmp_path / "private-output"),
        ]
    )

    captured = capsys.readouterr()
    assert exit_code == 1
    assert captured.err == "CHECKPOINT_RUNNER_IO_FAILED\n"
    assert private_detail not in captured.err
