from __future__ import annotations

import base64
from io import BytesIO

import pytest
from pptx import Presentation

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.deck_generation.design_pack_registry import load_design_pack_catalog
from app.ai.deck_generation.design_pack_selector import (
    DESIGN_PACK_DIRECTORY,
    select_system_design_pack,
)
from app.ai.deck_generation.models import GenerateDeckRequest
from app.ai.deck_generation.pipeline import analyze_input
from app.ai.design_pack_layouts.kickoff_alignment import (
    compile_kickoff_alignment_layout,
)
from app.ai.design_program import DeckDesignProgram


def test_kickoff_intent_selects_pack_without_adjacent_process_silhouette() -> None:
    selection = select_system_design_pack(kickoff_input(), golden_slides())
    registry = load_design_pack_catalog(DESIGN_PACK_DIRECTORY)
    layouts = {layout.layout_id: layout for layout in registry.layouts}
    silhouettes = [layouts[layout_id].silhouette_id for layout_id in selection.layout_ids]

    assert selection.pack_id == "kickoff-alignment"
    assert selection.fallback_used is False
    assert len(set(silhouettes)) >= 4
    assert all(left != right for left, right in zip(silhouettes, silhouettes[1:]))


def test_role_and_schedule_support_three_to_six_native_items() -> None:
    program = design_program()
    direction = program.slides[0]
    for count in (3, 6):
        role = compile_kickoff_alignment_layout(
            "kickoff-roles-01", direction, slide("comparison", count, "역할과 담당"), program
        )
        schedule = compile_kickoff_alignment_layout(
            "kickoff-schedule-01", direction, slide("process", count, "실행 일정"), program
        )

        assert sum(element["type"] == "text" for element in role.elements) >= count
        assert sum(
            element["type"] == "rect"
            and "_kickoff_schedule_bar_" in element["elementId"]
            for element in schedule.elements
        ) == count
        response = export_deck_pptx(
            DeckPptxExportRequest(deck=export_deck([role.elements, schedule.elements]))
        )
        assert len(
            Presentation(BytesIO(base64.b64decode(response.content_base64))).slides
        ) == 2


def test_role_and_schedule_reject_overflow() -> None:
    program = design_program()
    direction = program.slides[0]
    for layout_id in ("kickoff-roles-01", "kickoff-schedule-01"):
        for count in (2, 7):
            with pytest.raises(ValueError, match="supports 3 to 6 items"):
                compile_kickoff_alignment_layout(
                    layout_id, direction, slide("process", count, "일정"), program
                )


def kickoff_input():
    request = GenerateDeckRequest.model_validate(
        {
            "projectId": "project_demo_1",
            "topic": "신규 프로젝트 킥오프와 팀 얼라인먼트",
            "prompt": "역할, 로드맵, 일정 계획을 합의합니다.",
            "metadata": {"purpose": "inform"},
            "design": {"mediaPolicy": "minimal"},
        }
    )
    return analyze_input(request).model_copy(
        update={"presentation_profile": "general-inform"}
    )


def golden_slides() -> list[dict]:
    return [
        slide("cover", 0, "프로젝트 킥오프"),
        slide("process", 4, "오늘의 아젠다"),
        slide("solution", 3, "공동 목표"),
        slide("comparison", 4, "역할과 담당"),
        slide("process", 5, "업무 프로세스"),
        slide("process", 4, "핵심 로드맵"),
        slide("process", 6, "실행 일정"),
        slide("summary", 2, "다음 단계"),
    ]


def slide(slide_type: str, item_count: int, title: str) -> dict:
    return {
        "title": title,
        "message": "담당자와 완료 기준을 명확하게 합의합니다",
        "slideType": slide_type,
        "contentItems": [
            {"contentItemId": f"item-{index}", "text": f"마일스톤 {index}"}
            for index in range(1, item_count + 1)
        ],
        "mediaIntent": {"kind": "none"},
    }


def design_program() -> DeckDesignProgram:
    return DeckDesignProgram.model_validate(
        {
            "visualConcept": "Kickoff alignment",
            "paletteRoles": {
                "dominant": "#FFFFFF",
                "surface": "#EEF2FF",
                "text": "#0F172A",
                "focal": "#4F46E5",
                "secondary": "#0F766E",
            },
            "typography": {
                "headingFont": "Pretendard",
                "bodyFont": "Pretendard",
                "typeScale": {"cover": 68, "title": 48, "body": 26},
            },
            "backgroundSequence": ["light"],
            "imageStyle": "No decorative imagery",
            "surfaceStyle": "Structured alignment fields",
            "slides": [
                {
                    "order": 1,
                    "compositionId": "editorial-split",
                    "variant": "light",
                    "backgroundMode": "light",
                    "focalType": "message",
                    "assetRole": "none",
                    "requiredAsset": False,
                }
            ],
        }
    )


def export_deck(slide_elements: list[list[dict]]) -> dict:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {
            "backgroundColor": "#FFFFFF",
            "textColor": "#0F172A",
            "fontFamily": "Pretendard",
        },
        "slides": [
            {
                "order": index,
                "style": {"backgroundColor": "#FFFFFF"},
                "speakerNotes": "프로젝트 합의 내용을 설명합니다.",
                "elements": elements,
            }
            for index, elements in enumerate(slide_elements, start=1)
        ],
    }
