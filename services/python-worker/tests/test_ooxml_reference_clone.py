from __future__ import annotations

import hashlib
import posixpath
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from app.ai.ooxml_reference_templates.clone import clone_source_slides
from app.ai.ooxml_reference_templates.inventory import (
    ReferenceSource,
    inspect_reference_package,
)
from app.ai.ooxml_reference_templates.package import validate_cloned_package


PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
EXTENDED_PROPERTIES_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
)


def _source_package(tmp_path: Path) -> Path:
    image_path = tmp_path / "source-image.png"
    Image.new("RGB", (32, 24), (20, 80, 140)).save(image_path)

    presentation = Presentation()
    for index in range(10):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Identity {index + 1}"
        slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(6), Inches(1.5)
        ).text_frame.text = f"Body {index + 1}"
        slide.notes_slide.notes_text_frame.text = f"Notes {index + 1}"
        if index == 0:
            slide.shapes.add_picture(
                str(image_path), Inches(7), Inches(1), Inches(2), Inches(1.5)
            )
            slide.shapes.add_table(2, 2, Inches(1), Inches(3.5), Inches(4), Inches(1.5))
            chart_data = ChartData()
            chart_data.categories = ["A", "B"]
            chart_data.add_series("Series", (1, 2))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(6),
                Inches(3),
                Inches(5),
                Inches(3),
                chart_data,
            )

    output = tmp_path / "source.pptx"
    presentation.save(output)
    _add_timing(output)
    return output


def _add_timing(path: Path) -> None:
    with zipfile.ZipFile(path, "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}
    timing = (
        b'<p:timing><p:tnLst><p:par><p:cTn id="1" dur="indefinite"/>'
        b"</p:par></p:tnLst></p:timing>"
    )
    entries["ppt/slides/slide1.xml"] = entries["ppt/slides/slide1.xml"].replace(
        b"</p:sld>", timing + b"</p:sld>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as package:
        for name, content in entries.items():
            package.writestr(name, content)


def _add_stale_document_properties(path: Path) -> None:
    with zipfile.ZipFile(path, "r") as package:
        entries = {
            item.filename: package.read(item.filename)
            for item in package.infolist()
        }
    entries["docProps/app.xml"] = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
 xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <TotalTime>999</TotalTime><Words>999</Words><Paragraphs>99</Paragraphs>
  <Application>Microsoft PowerPoint</Application><PresentationFormat>Widescreen</PresentationFormat>
  <Slides>10</Slides><Notes>10</Notes><HiddenSlides>4</HiddenSlides>
  <HeadingPairs><vt:vector size="2" baseType="variant"><vt:variant><vt:lpstr>Slide Titles</vt:lpstr></vt:variant><vt:variant><vt:i4>10</vt:i4></vt:variant></vt:vector></HeadingPairs>
  <TitlesOfParts><vt:vector size="1" baseType="lpstr"><vt:lpstr>PRIVATE UNSELECTED SOURCE TITLE</vt:lpstr></vt:vector></TitlesOfParts>
  <Company>PRIVATE SOURCE COMPANY</Company><Manager>PRIVATE SOURCE MANAGER</Manager>
  <Template>PRIVATE SOURCE TEMPLATE</Template><HyperlinkBase>PRIVATE SOURCE PATH</HyperlinkBase>
  <ScaleCrop>false</ScaleCrop><LinksUpToDate>false</LinksUpToDate><SharedDoc>false</SharedDoc>
  <HyperlinksChanged>false</HyperlinksChanged><AppVersion>16.0000</AppVersion>
</Properties>"""
    entries["docProps/core.xml"] = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
 xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>PRIVATE SOURCE TITLE</dc:title><dc:creator>PRIVATE SOURCE AUTHOR</dc:creator></cp:coreProperties>"""
    entries["docProps/custom.xml"] = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
 xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"><property fmtid="{D5CDD505-2E9C-101B-9397-08002B2CF9AE}" pid="2" name="private-source"><vt:lpwstr>PRIVATE SOURCE CUSTOM VALUE</vt:lpwstr></property></Properties>"""
    entries["docProps/thumbnail.jpeg"] = b"PRIVATE SOURCE THUMBNAIL"

    relationships = ET.fromstring(entries["_rels/.rels"])
    ET.SubElement(
        relationships,
        f"{{{PKG_REL_NS}}}Relationship",
        {
            "Id": "rIdPrivateThumbnail",
            "Type": f"{PKG_REL_NS}/metadata/thumbnail",
            "Target": "docProps/thumbnail.jpeg",
        },
    )
    ET.SubElement(
        relationships,
        f"{{{PKG_REL_NS}}}Relationship",
        {
            "Id": "rIdPrivateCustom",
            "Type": f"{REL_NS}/custom-properties",
            "Target": "docProps/custom.xml",
        },
    )
    entries["_rels/.rels"] = ET.tostring(
        relationships, encoding="utf-8", xml_declaration=True
    )

    content_types = ET.fromstring(entries["[Content_Types].xml"])
    ET.SubElement(
        content_types,
        "{http://schemas.openxmlformats.org/package/2006/content-types}Override",
        {
            "PartName": "/docProps/custom.xml",
            "ContentType": "application/vnd.openxmlformats-officedocument.custom-properties+xml",
        },
    )
    entries["[Content_Types].xml"] = ET.tostring(
        content_types, encoding="utf-8", xml_declaration=True
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as package:
        for name, content in entries.items():
            package.writestr(name, content)


def _parts(package_bytes: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return {item.filename: package.read(item.filename) for item in package.infolist()}


def _relationship_types(parts: dict[str, bytes], source_part: str) -> set[str]:
    rels_part = posixpath.join(
        posixpath.dirname(source_part),
        "_rels",
        f"{posixpath.basename(source_part)}.rels",
    )
    root = ET.fromstring(parts[rels_part])
    return {
        relationship.attrib["Type"].rsplit("/", 1)[-1]
        for relationship in root.findall(f"{{{PKG_REL_NS}}}Relationship")
    }


def _presentation_ids(parts: dict[str, bytes]) -> tuple[list[str], list[str]]:
    presentation = ET.fromstring(parts["ppt/presentation.xml"])
    slide_ids = presentation.findall(f".//{{{PML_NS}}}sldId")
    return (
        [item.attrib["id"] for item in slide_ids],
        [item.attrib[f"{{{REL_NS}}}id"] for item in slide_ids],
    )


def _assert_existing_inventory_validator_accepts(
    tmp_path: Path, package_bytes: bytes, slide_count: int
) -> None:
    output = tmp_path / f"cloned-{slide_count}.pptx"
    output.write_bytes(package_bytes)
    inventory = inspect_reference_package(ReferenceSource("clone", output))
    assert inventory["securityPreflight"] == "passed"
    assert inventory["slideCount"] == slide_count


def test_clone_preserves_slide_transitive_relationship_graph(tmp_path: Path) -> None:
    source = _source_package(tmp_path)
    result = clone_source_slides(
        source.read_bytes(), source_slide_parts=["ppt/slides/slide1.xml"]
    )
    parts = _parts(result.package_bytes)
    cloned_part = result.clones[0].cloned_slide_part
    relation_types = _relationship_types(parts, cloned_part)

    assert {"slideLayout", "notesSlide", "image", "chart"}.issubset(relation_types)
    assert b"<p:timing" in parts[cloned_part]
    assert b"<a:tbl>" in parts[cloned_part]
    assert result.clones[0].layout_part in parts
    assert result.clones[0].master_part in parts
    assert result.clones[0].theme_part in parts
    assert result.clones[0].notes_parts
    assert all(part in parts for part in result.clones[0].notes_parts)
    assert result.clones[0].media_parts
    assert result.clones[0].chart_parts
    assert b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"' in parts[
        "ppt/_rels/presentation.xml.rels"
    ]
    assert b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"' in parts[
        "[Content_Types].xml"
    ]
    assert validate_cloned_package(result.package_bytes) == []
    _assert_existing_inventory_validator_accepts(tmp_path, result.package_bytes, 1)


def test_repeated_clone_allocates_collision_free_package_identities(tmp_path: Path) -> None:
    source = _source_package(tmp_path)
    result = clone_source_slides(
        source.read_bytes(),
        source_slide_parts=["ppt/slides/slide1.xml"] * 3,
    )
    parts = _parts(result.package_bytes)
    slide_ids, presentation_rids = _presentation_ids(parts)

    assert len({clone.cloned_slide_part for clone in result.clones}) == 3
    assert len(slide_ids) == len(set(slide_ids)) == 3
    assert len(presentation_rids) == len(set(presentation_rids)) == 3
    assert len({clone.presentation_slide_id for clone in result.clones}) == 3
    assert len({clone.presentation_relationship_id for clone in result.clones}) == 3
    assert len({part for clone in result.clones for part in clone.notes_parts}) == 3
    assert len({part for clone in result.clones for part in clone.chart_parts}) == 3
    assert len({part for clone in result.clones for part in clone.workbook_parts}) == 3

    for clone in result.clones:
        rels_part = posixpath.join(
            posixpath.dirname(clone.cloned_slide_part),
            "_rels",
            f"{posixpath.basename(clone.cloned_slide_part)}.rels",
        )
        root = ET.fromstring(parts[rels_part])
        rel_ids = [
            item.attrib["Id"]
            for item in root.findall(f"{{{PKG_REL_NS}}}Relationship")
        ]
        assert len(rel_ids) == len(set(rel_ids))

    assert validate_cloned_package(result.package_bytes) == []
    _assert_existing_inventory_validator_accepts(tmp_path, result.package_bytes, 3)


def test_clone_sanitizes_stale_document_properties_and_thumbnail(
    tmp_path: Path,
) -> None:
    source = _source_package(tmp_path)
    _add_stale_document_properties(source)

    result = clone_source_slides(
        source.read_bytes(),
        source_slide_parts=["ppt/slides/slide1.xml", "ppt/slides/slide2.xml"],
    )
    parts = _parts(result.package_bytes)
    app_properties = ET.fromstring(parts["docProps/app.xml"])

    assert app_properties.findtext(f"{{{EXTENDED_PROPERTIES_NS}}}Slides") == "2"
    assert app_properties.findtext(f"{{{EXTENDED_PROPERTIES_NS}}}HiddenSlides") == "0"
    for local_name in (
        "TotalTime",
        "Words",
        "Paragraphs",
        "Notes",
        "AppVersion",
        "HeadingPairs",
        "TitlesOfParts",
        "Company",
        "Manager",
        "Template",
        "HyperlinkBase",
    ):
        assert app_properties.find(f"{{{EXTENDED_PROPERTIES_NS}}}{local_name}") is None
    assert b"PRIVATE SOURCE" not in parts["docProps/app.xml"]
    assert b"PRIVATE SOURCE" not in parts["docProps/core.xml"]
    assert "docProps/custom.xml" not in parts
    assert "docProps/thumbnail.jpeg" not in parts

    root_relationships = ET.fromstring(parts["_rels/.rels"])
    root_relationship_types = {
        item.attrib["Type"].rsplit("/", 1)[-1]
        for item in root_relationships.findall(f"{{{PKG_REL_NS}}}Relationship")
    }
    assert "custom-properties" not in root_relationship_types
    assert "thumbnail" not in root_relationship_types
    assert validate_cloned_package(result.package_bytes) == []


@pytest.mark.parametrize("slide_count", [8, 9, 10])
def test_identity_control_clones_eight_to_ten_unique_source_slides(
    tmp_path: Path, slide_count: int
) -> None:
    source = _source_package(tmp_path)
    source_parts = _parts(source.read_bytes())
    selected = [f"ppt/slides/slide{index}.xml" for index in range(1, slide_count + 1)]

    result = clone_source_slides(
        source.read_bytes(), source_slide_parts=selected, identity_control=True
    )
    cloned_parts = _parts(result.package_bytes)

    assert len(result.clones) == slide_count
    assert result.identity_control_slide_count == slide_count
    for clone in result.clones:
        assert clone.source_slide_part in selected
        source_hash = hashlib.sha256(source_parts[clone.source_slide_part]).digest()
        cloned_hash = hashlib.sha256(cloned_parts[clone.cloned_slide_part]).digest()
        assert source_hash == cloned_hash
    assert validate_cloned_package(result.package_bytes) == []
    _assert_existing_inventory_validator_accepts(
        tmp_path, result.package_bytes, slide_count
    )
