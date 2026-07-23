from __future__ import annotations

import hashlib
import base64
import json
from io import BytesIO
from types import SimpleNamespace
from typing import Any

import pytest
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation

from app.ai.ooxml_reference_templates.catalog_transport import (
    VerifiedPrivateCatalogSource,
)
from app.ai.ooxml_reference_templates.calibration import (
    CALIBRATION_CONTENT_TYPE,
    CALIBRATION_OBJECT_KEY,
    PrivateFidelityCalibrationError,
)
from app.ai.ooxml_reference_templates.fidelity import (
    EXPECTED_TEMPLATE_IDS,
    evaluate_ooxml_reference_fidelity,
)
from app.ai.ooxml_reference_templates.models import (
    OoxmlReferenceTemplateGenerationRequest,
    OoxmlReferenceTemplateManifest,
)
from app.ai.ooxml_reference_templates.private_generation_runtime import (
    PrivateGenerationRuntimeError,
    PrivateOoxmlReferenceGenerationRuntime,
    ProjectImageAsset,
    _locked_snapshot,
    _slot_mask_png,
    generated_storage_key,
)
from app.ai.deck_generation.models import (
    ContentPlan,
    DeckOutline,
    GenerateDeckRequest,
    PresentationTimingPlan,
    SlidePlan,
)
from app.config import load_config
from app.main import app, configure_ooxml_reference_template_catalog
from app.ai.pptx_design_importer import ImportedDesignAsset
from app.references import (
    EMBEDDING_DIMENSION,
    ReferenceSearchResult,
)
from test_config import VALID_ENV


PPTX_MIME_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
SOURCE = b"PK\x03\x04private-reference-source"


class MissingObjectError(Exception):
    def __init__(self) -> None:
        self.response = {"Error": {"Code": "NoSuchKey"}}


class PreconditionFailedError(Exception):
    def __init__(self) -> None:
        self.response = {"Error": {"Code": "PreconditionFailed"}}


class FakeS3Client:
    def __init__(self) -> None:
        self.objects: dict[str, dict[str, Any]] = {}
        self.put_count = 0

    def head_object(self, **kwargs: object) -> dict[str, object]:
        value = self.objects.get(str(kwargs["Key"]))
        if value is None:
            raise MissingObjectError
        return {
            "ContentLength": len(value["content"]),
            "ContentType": value["content_type"],
            "Metadata": {"sha256": value["sha256"]},
        }

    def get_object(self, **kwargs: object) -> dict[str, object]:
        value = self.objects.get(str(kwargs["Key"]))
        if value is None:
            raise MissingObjectError
        return {"Body": BytesIO(value["content"])}

    def put_object(self, **kwargs: object) -> dict[str, object]:
        key = str(kwargs["Key"])
        if kwargs.get("IfNoneMatch") == "*" and key in self.objects:
            raise PreconditionFailedError
        content = bytes(kwargs["Body"])
        metadata = dict(kwargs["Metadata"])
        self.objects[key] = {
            "content": content,
            "content_type": str(kwargs["ContentType"]),
            "sha256": str(metadata["sha256"]),
        }
        self.put_count += 1
        return {}


class CatalogRuntime:
    def __init__(self, source: bytes = SOURCE) -> None:
        self.source = source
        self.manifest = _manifest(source)

    def read_source(
        self,
        template_id: str,
        version: int,
    ) -> VerifiedPrivateCatalogSource:
        assert (template_id, version) == ("operating-review", 1)
        return VerifiedPrivateCatalogSource(
            manifest=self.manifest,
            content=self.source,
        )


class FakeEmbeddingClient:
    class Embeddings:
        def create(self, *, model: str, input: list[str]) -> object:
            del model
            data = [
                type(
                    "Embedding",
                    (),
                    {"embedding": [0.01] * EMBEDDING_DIMENSION},
                )()
                for _ in input
            ]
            return type("EmbeddingResponse", (), {"data": data})()

    embeddings = Embeddings()


class OrderedReferenceRepository:
    def __init__(self) -> None:
        self.calls: list[tuple[str, list[str] | None]] = []

    def search_chunks(
        self,
        project_id: str,
        query_embedding: list[float],
        *,
        limit: int = 6,
        file_ids: list[str] | None = None,
    ) -> list[ReferenceSearchResult]:
        del query_embedding, limit
        self.calls.append((project_id, file_ids))
        file_id = (file_ids or [""])[0]
        return [
            ReferenceSearchResult(
                chunk_id=f"{file_id}-chunk-2",
                project_id=project_id,
                file_id=file_id,
                chunk_index=2,
                content=f"{file_id} 두 번째 근거",
                metadata={"fileName": f"{file_id}.pdf"},
                score=0.95,
            ),
            ReferenceSearchResult(
                chunk_id=f"{file_id}-chunk-1",
                project_id=project_id,
                file_id=file_id,
                chunk_index=1,
                content=f"{file_id} 첫 번째 근거",
                metadata={"fileName": f"{file_id}.pdf"},
                score=0.9,
            ),
        ]


class ProjectAssetReader:
    def __init__(self, asset: ProjectImageAsset) -> None:
        self.asset = asset
        self.calls: list[tuple[str, str]] = []

    def get_image_asset(self, project_id: str, file_id: str) -> ProjectImageAsset:
        self.calls.append((project_id, file_id))
        return self.asset


def test_generated_storage_key_matches_typescript_parity_vectors() -> None:
    assert generated_storage_key(
        "project_1",
        "generation-1",
        "file_1",
        "operating-review-generated.pptx",
    ) == (
        "projects/project_1/ooxml-reference-generations/generation-1/"
        "file_1/operating-review-generated.pptx"
    )
    assert generated_storage_key(
        "project_1",
        "generation-1",
        "file_2",
        "운영 리뷰 (최종)!?.pptx",
    ) == (
        "projects/project_1/ooxml-reference-generations/generation-1/file_2/"
        "%EC%9A%B4%EC%98%81%20%EB%A6%AC%EB%B7%B0%20"
        "(%EC%B5%9C%EC%A2%85)!%3F.pptx"
    )


def test_generated_package_put_is_idempotent_and_conflict_fails_closed() -> None:
    client = FakeS3Client()
    runtime = _runtime(client)

    first = runtime.store_current_package(
        "generation-1",
        "project_1",
        "operating-review",
        b"same-package",
    )
    second = runtime.store_current_package(
        "generation-1",
        "project_1",
        "operating-review",
        b"same-package",
    )

    assert first == second
    assert client.put_count == 1
    assert list(client.objects) == [
        generated_storage_key(
            "project_1",
            "generation-1",
            first.file_id,
            first.original_name,
        )
    ]
    with pytest.raises(PrivateGenerationRuntimeError) as caught:
        runtime.store_current_package(
            "generation-1",
            "project_1",
            "operating-review",
            b"different-package",
        )
    assert caught.value.code == "OOXML_REFERENCE_GENERATED_ASSET_CONFLICT"
    assert "projects/" not in str(caught.value)


def test_generated_package_is_readable_by_a_fresh_runtime_instance() -> None:
    client = FakeS3Client()
    stored = _runtime(client).store_current_package(
        "generation-1",
        "project_1",
        "operating-review",
        b"restart-safe-package",
    )

    content = _runtime(client).read_current_package(
        "generation-1",
        "project_1",
        "operating-review",
        stored.file_id,
    )

    assert content == b"restart-safe-package"


def test_reference_chunks_are_loaded_with_project_and_requested_file_order() -> None:
    repository = OrderedReferenceRepository()
    runtime = _runtime(
        FakeS3Client(),
        reference_repository=repository,
        embedding_client=FakeEmbeddingClient(),
    )
    request = OoxmlReferenceTemplateGenerationRequest.model_validate(
        {
            "topic": "운영 리뷰",
            "referencePolicy": "references-only",
            "referenceFileIds": ["file-b", "file-a"],
            "templateSelection": {
                "mode": "user",
                "templateId": "operating-review",
                "version": 1,
            },
        }
    )

    references = runtime.extract_references("project_1", request)

    assert [item.file_id for item in references] == ["file-b", "file-a"]
    assert references[0].content == ("file-b 첫 번째 근거\n\nfile-b 두 번째 근거")
    assert references[0].metadata == {
        "chunkIds": ["file-b-chunk-1", "file-b-chunk-2"],
        "chunkCount": 2,
    }
    assert repository.calls == [
        ("project_1", ["file-b"]),
        ("project_1", ["file-a"]),
    ]


def test_topic_only_does_not_load_selected_reference_chunks() -> None:
    repository = OrderedReferenceRepository()
    runtime = _runtime(
        FakeS3Client(),
        reference_repository=repository,
        embedding_client=FakeEmbeddingClient(),
    )
    request = OoxmlReferenceTemplateGenerationRequest.model_validate(
        {
            "topic": "운영 리뷰",
            "referencePolicy": "topic-only",
            "referenceFileIds": ["file-a"],
            "templateSelection": {
                "mode": "user",
                "templateId": "operating-review",
                "version": 1,
            },
        }
    )

    assert runtime.extract_references("project_1", request) == []
    assert repository.calls == []


def test_content_planning_reuses_grounded_plan_through_geometry_free_adapter() -> None:
    captured: list[GenerateDeckRequest] = []

    def plan(request: GenerateDeckRequest) -> ContentPlan:
        captured.append(request)
        return _grounded_content_plan()

    runtime = _runtime(FakeS3Client(), content_plan_runner=plan)
    request = OoxmlReferenceTemplateGenerationRequest.model_validate(
        {
            "topic": "운영 리뷰",
            "slideCountRange": {"min": 2, "max": 2},
            "referencePolicy": "references-first",
            "referenceFileIds": ["file-a"],
            "templateSelection": {
                "mode": "user",
                "templateId": "operating-review",
                "version": 1,
            },
        }
    )

    result = runtime.plan_content(
        "project_1",
        request,
        [
            {
                "fileId": "file-a",
                "content": "검증된 운영 근거",
                "metadata": {"chunkIds": ["chunk-a"]},
            }
        ],
    )

    assert result.template_id == "operating-review"
    assert [slide.semantic_role for slide in result.slides] == ["cover", "closing"]
    assert captured[0].project_id == "project_1"
    assert captured[0].reference_context[0].content == "검증된 운영 근거"
    dumped = result.model_dump(by_alias=True, mode="json")
    assert not any(
        forbidden in str(dumped) for forbidden in ('"x"', '"y"', "zIndex", "geometry")
    )


def test_render_runtime_uses_pptx_renderer_and_existing_fidelity_gate() -> None:
    source = _source_pptx()
    catalog = CatalogRuntime(source)
    png = _png()

    def render(
        package_bytes: bytes,
        canvas: object,
    ) -> list[ImportedDesignAsset]:
        del package_bytes, canvas
        return [
            ImportedDesignAsset(
                assetId=f"slide_render_{order}",
                fileName=f"slide-{order:02d}.png",
                mimeType="image/png",
                contentBase64=base64.b64encode(png).decode("ascii"),
            )
            for order in (1, 2)
        ]

    client = FakeS3Client()
    runtime = _runtime(
        client,
        catalog=catalog,
        render_deck=render,
        render_environment=lambda: {
            "renderer": "libreoffice-pdf-pymupdf",
            "rendererVersion": "25.2.3.2",
            "fontFiles": [{"family": "Arial", "role": "body", "sha256": "f" * 64}],
        },
        fidelity_calibration=_calibration(),
    )
    current = runtime.store_current_package(
        "generation-1",
        "project_1",
        "operating-review",
        source,
    )

    validation = runtime.render_and_prepare_fidelity(
        "generation-1",
        "project_1",
        "operating-review",
        1,
        current.file_id,
        source,
        ["cover-01", "closing-02"],
    )
    report = evaluate_ooxml_reference_fidelity(
        template_id="operating-review",
        template_version=1,
        mode="generated-comparison",
        slides=validation.slides,
        package_warnings=[],
        environment=validation.environment,
        calibration=validation.calibration,
    )

    assert report["status"] == "passed"
    assert [asset.original_name for asset in validation.assets] == [
        "slide-001.png",
        "slide-002.png",
    ]
    assert (
        runtime.render_assets(
            "generation-1",
            "project_1",
            current.file_id,
            list(validation.assets),
        )
        == validation.assets
    )
    assert all("projects/" not in str(slide) for slide in validation.slides)
    assert all(
        {shape["shapeId"] for shape in slide["sourceLockedSnapshot"]["shapes"]}
        == {"2"}
        for slide in validation.slides
    )
    assert all(
        set(slide["sourceLockedSnapshot"]["relationships"])
        == {"layout", "master", "theme"}
        for slide in validation.slides
    )


def test_image_slot_reads_only_project_scoped_private_asset_metadata() -> None:
    content = _png()
    key = "projects/project_1/assets/file-image-source.png"
    client = FakeS3Client()
    client.objects[key] = {
        "content": content,
        "content_type": "image/png",
        "sha256": hashlib.sha256(content).hexdigest(),
    }
    reader = ProjectAssetReader(
        ProjectImageAsset(
            storage_key=key,
            mime_type="image/png",
            size=len(content),
            sha256=hashlib.sha256(content).hexdigest(),
        )
    )
    runtime = _runtime(client, project_asset_reader=reader)

    loaded, mime_type = runtime.read_image_asset("project_1", "file-image")

    assert loaded == content
    assert mime_type == "image/png"
    assert reader.calls == [("project_1", "file-image")]


def test_http_stage_uses_injected_private_runtime_for_topic_only_success() -> None:
    runtime = _runtime(FakeS3Client())
    app.state.ooxml_reference_generation_runtime = runtime
    try:
        response = TestClient(app).post(
            "/internal/ai/ooxml-reference-template-generation/stage",
            json={
                "jobId": "generation-1",
                "projectId": "project_1",
                "stage": "reference-extract-file",
                "templateId": "operating-review",
                "templateVersion": 1,
                "request": {
                    "topic": "운영 리뷰",
                    "slideCountRange": {"min": 2, "max": 2},
                    "referencePolicy": "topic-only",
                    "templateSelection": {
                        "mode": "user",
                        "templateId": "operating-review",
                        "version": 1,
                    },
                },
                "dependencies": [],
            },
        )
    finally:
        del app.state.ooxml_reference_generation_runtime

    assert response.status_code == 200
    assert response.json() == {
        "stage": "reference-extract-file",
        "templateId": "operating-review",
        "templateVersion": 1,
        "sourceSlideCount": 0,
        "slotCount": 0,
        "artifact": {"references": []},
        "issueCodes": [],
    }


def test_main_wiring_is_flag_and_exact_allowlist_gated() -> None:
    disabled = SimpleNamespace()
    configure_ooxml_reference_template_catalog(
        disabled,
        load_config(VALID_ENV),
        client=FakeS3Client(),
    )
    enabled = SimpleNamespace()
    enabled_client = FakeS3Client()
    _seed_calibration(enabled_client)
    configure_ooxml_reference_template_catalog(
        enabled,
        load_config(
            {
                **VALID_ENV,
                "AI_PPT_OOXML_REFERENCE_TEMPLATES_ENABLED": "true",
                "AI_PPT_OOXML_REFERENCE_TEMPLATE_ALLOWLIST": (
                    "operating-review@1"
                ),
            }
        ),
        client=enabled_client,
    )

    assert disabled.ooxml_reference_generation_runtime is None
    assert isinstance(
        enabled.ooxml_reference_generation_runtime,
        PrivateOoxmlReferenceGenerationRuntime,
    )


def test_enabled_main_wiring_fails_closed_without_private_calibration() -> None:
    state = SimpleNamespace()

    with pytest.raises(PrivateFidelityCalibrationError) as caught:
        configure_ooxml_reference_template_catalog(
            state,
            load_config(
                {
                    **VALID_ENV,
                    "AI_PPT_OOXML_REFERENCE_TEMPLATES_ENABLED": "true",
                    "AI_PPT_OOXML_REFERENCE_TEMPLATE_ALLOWLIST": (
                        "operating-review@1"
                    ),
                }
            ),
            client=FakeS3Client(),
        )

    assert (
        caught.value.code
        == "OOXML_REFERENCE_FIDELITY_CALIBRATION_UNAVAILABLE"
    )
    assert CALIBRATION_OBJECT_KEY not in str(caught.value)


def _runtime(
    client: FakeS3Client,
    *,
    reference_repository: object | None = None,
    embedding_client: object | None = None,
    content_plan_runner: object | None = None,
    catalog: CatalogRuntime | None = None,
    render_deck: object | None = None,
    render_environment: object | None = None,
    fidelity_calibration: dict[str, object] | None = None,
    project_asset_reader: object | None = None,
) -> PrivateOoxmlReferenceGenerationRuntime:
    return PrivateOoxmlReferenceGenerationRuntime(
        catalog=catalog or CatalogRuntime(),
        storage_client=client,
        bucket="private-bucket",
        database_url="postgres://unused",
        embedding_model="embedding-test",
        content_model="content-test",
        api_key=None,
        reference_repository=reference_repository,
        embedding_client=embedding_client,
        content_plan_runner=content_plan_runner,
        render_deck=render_deck,
        render_environment=render_environment,
        fidelity_calibration=fidelity_calibration,
        project_asset_reader=project_asset_reader,
    )


def _grounded_content_plan() -> ContentPlan:
    slides = [
        SlidePlan(
            order=1,
            slide_type="cover",
            title="운영 리뷰",
            message="핵심 의사결정",
            speaker_notes="표지",
            keywords=[],
            evidence=[],
        ),
        SlidePlan(
            order=2,
            slide_type="summary",
            title="다음 단계",
            message="실행 항목 확정",
            speaker_notes="마무리",
            keywords=[],
            evidence=[],
        ),
    ]
    return ContentPlan(
        outline=DeckOutline(
            title="운영 리뷰",
            slide_titles=[slide.title for slide in slides],
        ),
        slidePlans=slides,
        slideCount=2,
        timingPlan=PresentationTimingPlan(
            charsPerMinute=400,
            speakingTimeRatio=0.8,
            targetTotalChars=800,
            targetSpokenSeconds=120,
            targetSlideCount=2,
            targetSecondsPerSlide=60,
            targetSpeakerNotesCharsPerSlide=200,
        ),
        repairAttempted=False,
        repairReasonCodes=[],
    )


def _source_pptx() -> bytes:
    presentation = Presentation()
    presentation.slide_width = 12_192_000
    presentation.slide_height = 6_858_000
    for title in ("Source cover", "Source closing"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(600_000, 500_000, 8_000_000, 900_000)
        shape.text = title
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _png() -> bytes:
    output = BytesIO()
    Image.new("RGB", (160, 90), "white").save(output, format="PNG")
    return output.getvalue()


@pytest.mark.parametrize("layout_index", [0, 1])
def test_fidelity_geometry_resolves_inherited_placeholder_chain(
    layout_index: int,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    title = slide.shapes.title
    assert title is not None
    title.text = "Inherited placeholder"
    output = BytesIO()
    presentation.save(output)
    package = output.getvalue()
    shape_id = str(title.shape_id)

    mask = _slot_mask_png(
        package,
        "ppt/slides/slide1.xml",
        {shape_id},
        presentation.slide_width,
        presentation.slide_height,
        _png(),
    )
    with Image.open(BytesIO(mask)) as image:
        assert image.getbbox() is not None

    snapshot = _locked_snapshot(
        package,
        "ppt/slides/slide1.xml",
        {shape_id},
    )
    title_snapshot = next(
        shape for shape in snapshot["shapes"] if shape["shapeId"] == shape_id
    )
    assert title_snapshot["geometry"]["transform"] is not None


def _calibration() -> dict[str, object]:
    return {
        "schemaVersion": 1,
        "status": "calibrated",
        "lockedRegionSsimThreshold": 0.998,
        "geometryEdgeTolerancePx": 0,
        "rationale": "deterministic renderer identity baselines",
        "identityBaselines": [
            {
                "version": 1,
                "templateId": template_id,
                "renderer": "libreoffice-pdf-pymupdf",
                "rendererVersion": "25.2.3.2",
                "reportSha256": "e" * 64,
            }
            for template_id in EXPECTED_TEMPLATE_IDS
        ],
    }


def _seed_calibration(client: FakeS3Client) -> None:
    content = json.dumps(
        _calibration(),
        ensure_ascii=True,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    client.objects[CALIBRATION_OBJECT_KEY] = {
        "content": content,
        "content_type": CALIBRATION_CONTENT_TYPE,
        "sha256": hashlib.sha256(content).hexdigest(),
    }


def _manifest(source: bytes = SOURCE) -> OoxmlReferenceTemplateManifest:
    slides = []
    for order, role in ((1, "cover"), (2, "closing")):
        slides.append(
            {
                "sourceSlideId": f"{role}-{order:02d}",
                "sourceSlidePart": f"ppt/slides/slide{order}.xml",
                "sourceOrder": order,
                "semanticRole": role,
                "relationships": {
                    "layoutPart": "ppt/slideLayouts/slideLayout1.xml",
                    "masterPart": "ppt/slideMasters/slideMaster1.xml",
                    "themePart": "ppt/theme/theme1.xml",
                },
                "capacity": {
                    "textSlotCount": 1,
                    "imageSlotCount": 0,
                    "tableSlotCount": 0,
                    "chartSlotCount": 0,
                },
                "previewId": role,
                "lockedInventorySha256": "a" * 64,
                "slots": [
                    {
                        "slotId": f"{role}-title",
                        "contentType": "text",
                        "semanticRole": "title",
                        "required": True,
                        "locator": {
                            "slidePart": f"ppt/slides/slide{order}.xml",
                            "shapeId": "2",
                            "placeholderType": "title",
                            "relationshipId": None,
                        },
                        "mutationPolicy": ["text-content"],
                        "replacementPolicy": {"overflow": "fail"},
                        "capacity": {"maxChars": 120, "maxLines": 2},
                    }
                ],
            }
        )
    return OoxmlReferenceTemplateManifest.model_validate(
        {
            "templateId": "operating-review",
            "version": 1,
            "status": "active",
            "sourceFormat": "pptx",
            "sourceSha256": hashlib.sha256(source).hexdigest(),
            "slideCount": 2,
            "canvas": {
                "aspectRatio": "16:9",
                "widthEmu": 12_192_000,
                "heightEmu": 6_858_000,
            },
            "name": "Operating Review",
            "description": "운영 리뷰",
            "preview": {
                "coverPreviewId": "cover",
                "coverPreviewSha256": "b" * 64,
                "bodyPreviewId": "closing",
                "bodyPreviewSha256": "c" * 64,
            },
            "sourceSlides": slides,
            "provenance": {
                "authorizationStatus": "approved",
                "inventoryVersion": 1,
            },
        }
    )
