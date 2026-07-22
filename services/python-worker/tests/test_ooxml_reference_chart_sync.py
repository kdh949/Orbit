from __future__ import annotations

import copy
import hashlib
import posixpath
import zipfile
from io import BytesIO
from pathlib import PurePosixPath
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from app.ai.ooxml_reference_templates.chart_sync import (
    ChartDataSyncError,
    build_chart_data_locator,
    sync_chart_data,
)
from app.ai.pptx_ooxml_generation import (
    PackageFrameScale,
    apply_patch_operations_to_package,
)


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _parts(package_bytes: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return {
            item.filename: package.read(item.filename)
            for item in package.infolist()
            if not item.is_dir()
        }


def _write_parts(parts: dict[str, bytes]) -> bytes:
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as package:
        for name, content in sorted(parts.items()):
            package.writestr(name, content)
    return output.getvalue()


def _rels_part(part: str) -> str:
    path = PurePosixPath(part)
    return str(path.parent / "_rels" / f"{path.name}.rels")


def _resolve(source_part: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _fixture() -> tuple[bytes, dict[str, object]]:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    data = ChartData()
    data.categories = ["North", "South"]
    data.add_series("Revenue", (10.5, 20.25), number_format="0.00")
    data.add_series("Cost", (4.25, 8.5), number_format="0.00")
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.25),
        Inches(1.5),
        Inches(9),
        Inches(4.5),
        data,
    )
    output = BytesIO()
    presentation.save(output)
    package_bytes = output.getvalue()
    parts = _parts(package_bytes)
    slide_part = "ppt/slides/slide1.xml"
    slide_root = ET.fromstring(parts[slide_part])
    frame = slide_root.find(f".//{{{PRESENTATION_NS}}}graphicFrame")
    assert frame is not None
    shape = frame.find(f".//{{{PRESENTATION_NS}}}cNvPr")
    chart_ref = frame.find(f".//{{{CHART_NS}}}chart")
    assert shape is not None and chart_ref is not None
    relationship_id = chart_ref.attrib[f"{{{REL_NS}}}id"]
    source: dict[str, object] = {
        "elementId": "el_chart_1",
        "elementType": "chart",
        "ooxmlOrigin": "imported",
        "slidePart": slide_part,
        "shapeId": shape.attrib["id"],
        "relationshipId": relationship_id,
        "sourceType": "chart",
        "writable": True,
    }
    return package_bytes, source


def _style_signature(chart_xml: bytes) -> tuple[bytes, ...]:
    root = ET.fromstring(chart_xml)
    return tuple(
        ET.tostring(node)
        for path in ("barDir", "grouping", "legend", "axId", "spPr")
        for node in root.findall(f".//{{{CHART_NS}}}{path}")
    )


def _authoritative_source() -> tuple[bytes, dict[str, object]]:
    package_bytes, source = _fixture()
    locator = build_chart_data_locator(package_bytes, source=source)
    source["chartDataLocator"] = locator
    source["ooxmlEditCapabilities"] = {
        "richText": "none",
        "crop": "none",
        "tableCellText": False,
        "chartData": True,
        "frame": False,
        "delete": False,
        "imageSource": False,
    }
    return package_bytes, source


def test_chart_data_sync_updates_cache_and_workbook_but_preserves_frame_and_style() -> None:
    package_bytes, source = _authoritative_source()
    locator = source["chartDataLocator"]
    assert isinstance(locator, dict)
    before = _parts(package_bytes)

    result = sync_chart_data(
        package_bytes,
        source=source,
        props={
            "type": "bar",
            "data": [
                {"label": "Enterprise", "series": "Revenue", "value": 101.25},
                {"label": "SMB", "series": "Revenue", "value": 81.5},
                {"label": "Enterprise", "series": "Cost", "value": 50},
                {"label": "SMB", "series": "Cost", "value": 35.75},
            ],
        },
    )
    after = _parts(result.package_bytes)

    assert before[source["slidePart"]] == after[source["slidePart"]]
    assert _style_signature(before[locator["chartPart"]]) == _style_signature(
        after[locator["chartPart"]]
    )
    assert result.updated_source["chartDataLocator"]["workbookFingerprint"] == (
        hashlib.sha256(after[locator["workbookPart"]]).hexdigest()
    )
    assert result.updated_source["chartDataLocator"]["workbookFingerprint"] != (
        locator["workbookFingerprint"]
    )


@pytest.mark.parametrize(
    ("drift", "code"),
    [
        ("type", "CHART_TYPE_UNSUPPORTED"),
        ("count", "CHART_STRUCTURE_UNSUPPORTED"),
        ("formula", "CHART_FORMULA_DRIFT"),
        ("fingerprint", "CHART_WORKBOOK_FINGERPRINT_MISMATCH"),
    ],
)
def test_chart_sync_drift_is_fail_closed(drift: str, code: str) -> None:
    package_bytes, source = _authoritative_source()
    props: dict[str, object] = {
        "type": "bar",
        "data": [
            {"label": "North", "series": "Revenue", "value": 1},
            {"label": "South", "series": "Revenue", "value": 2},
            {"label": "North", "series": "Cost", "value": 3},
            {"label": "South", "series": "Cost", "value": 4},
        ],
    }
    if drift == "type":
        props["type"] = "line"
    elif drift == "count":
        props["data"] = list(props["data"])[0:2]
    elif drift == "formula":
        parts = _parts(package_bytes)
        locator = source["chartDataLocator"]
        assert isinstance(locator, dict)
        chart_part = str(locator["chartPart"])
        chart = ET.fromstring(parts[chart_part])
        formula = chart.find(f".//{{{CHART_NS}}}cat/{{{CHART_NS}}}strRef/{{{CHART_NS}}}f")
        assert formula is not None
        formula.text = "Sheet1!$A$3:$A$4"
        parts[chart_part] = ET.tostring(chart, encoding="utf-8", xml_declaration=True)
        package_bytes = _write_parts(parts)
    else:
        changed = copy.deepcopy(source)
        changed_locator = changed["chartDataLocator"]
        assert isinstance(changed_locator, dict)
        changed_locator["workbookFingerprint"] = "0" * 64
        source = changed

    with pytest.raises(ChartDataSyncError) as caught:
        sync_chart_data(package_bytes, source=source, props=props)

    assert caught.value.code == code
    assert caught.value.package_bytes == package_bytes


def test_pptx_ooxml_sync_routes_authoritative_chart_props_and_returns_new_fingerprint() -> None:
    package_bytes, source = _authoritative_source()
    blueprint = {
        "templateId": "template_chart",
        "sourceFileId": "file_chart",
        "slides": [
            {
                "slideId": "slide_1",
                "slideIndex": 1,
                "sourceSlideIndex": 1,
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "elementSources": [source],
            }
        ],
    }
    operation = {
        "type": "update_element_props",
        "slideId": "slide_1",
        "elementId": "el_chart_1",
        "props": {
            "type": "bar",
            "data": [
                {"label": "North", "series": "Revenue", "value": 11},
                {"label": "South", "series": "Revenue", "value": 22},
                {"label": "North", "series": "Cost", "value": 5},
                {"label": "South", "series": "Cost", "value": 9},
            ],
        },
    }

    synced, _, updated, applied, unsupported, _, _ = apply_patch_operations_to_package(
        package_bytes,
        blueprint,
        [operation],
        PackageFrameScale(1920, 1080, 12_192_000, 6_858_000),
    )

    assert synced != package_bytes
    assert [item.operation_type for item in applied] == ["update_element_props"]
    assert unsupported == []
    assert updated[0]["ooxmlEditCapabilities"]["chartData"] is True
    assert updated[0]["chartDataLocator"]["workbookFingerprint"] != (
        source["chartDataLocator"]["workbookFingerprint"]
    )


def test_pptx_ooxml_sync_keeps_original_package_when_chart_gate_rejects() -> None:
    package_bytes, source = _authoritative_source()
    locator = source["chartDataLocator"]
    assert isinstance(locator, dict)
    locator["workbookFingerprint"] = "0" * 64
    blueprint = {
        "templateId": "template_chart",
        "sourceFileId": "file_chart",
        "slides": [
            {
                "slideId": "slide_1",
                "slideIndex": 1,
                "sourceSlideIndex": 1,
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "elementSources": [source],
            }
        ],
    }
    operation = {
        "type": "update_element_props",
        "slideId": "slide_1",
        "elementId": "el_chart_1",
        "props": {
            "type": "bar",
            "data": [
                {"label": "North", "series": "Revenue", "value": 11},
                {"label": "South", "series": "Revenue", "value": 22},
                {"label": "North", "series": "Cost", "value": 5},
                {"label": "South", "series": "Cost", "value": 9},
            ],
        },
    }

    synced, _, updated, applied, unsupported, _, _ = apply_patch_operations_to_package(
        package_bytes,
        blueprint,
        [operation],
        PackageFrameScale(1920, 1080, 12_192_000, 6_858_000),
    )

    assert synced == package_bytes
    assert updated == []
    assert applied == []
    assert [item.reason_code for item in unsupported] == [
        "CHART_WORKBOOK_FINGERPRINT_MISMATCH"
    ]


def test_locator_is_not_issued_for_ambiguous_chart_relationship() -> None:
    package_bytes, source = _fixture()
    parts = _parts(package_bytes)
    slide_rels_part = _rels_part(str(source["slidePart"]))
    relationships = ET.fromstring(parts[slide_rels_part])
    original = next(
        item
        for item in relationships.findall(f"{{{PACKAGE_REL_NS}}}Relationship")
        if item.attrib.get("Id") == source["relationshipId"]
    )
    duplicate = copy.deepcopy(original)
    duplicate.attrib["Id"] = str(source["relationshipId"])
    relationships.append(duplicate)
    parts[slide_rels_part] = ET.tostring(
        relationships, encoding="utf-8", xml_declaration=True
    )

    with pytest.raises(ChartDataSyncError) as caught:
        build_chart_data_locator(_write_parts(parts), source=source)

    assert caught.value.code == "CHART_RELATIONSHIP_UNSAFE"
