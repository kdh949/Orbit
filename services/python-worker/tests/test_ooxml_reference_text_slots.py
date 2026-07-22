from __future__ import annotations

import copy
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from app.ai.ooxml_reference_templates.capacity import SlotCapacityError
from app.ai.ooxml_reference_templates.clone import clone_source_slides
from app.ai.ooxml_reference_templates.models import OoxmlTextTemplateSlot
from app.ai.ooxml_reference_templates.package import validate_cloned_package
from app.ai.ooxml_reference_templates.text_slots import replace_text_slot


DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _text_fixture(tmp_path: Path, *, font_family: str = "Aptos") -> tuple[bytes, str]:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
    first = textbox.text_frame.paragraphs[0]
    first.alignment = PP_ALIGN.CENTER
    first_properties = first._p.get_or_add_pPr()
    first_properties.set("marL", str(int(Inches(0.3))))
    first_properties.append(
        parse_xml(f'<a:buChar xmlns:a="{DML_NS}" char="▪"/>')
    )
    run = first.add_run()
    run.text = "Original one"
    run.font.name = font_family
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.italic = True
    second = textbox.text_frame.add_paragraph()
    second.alignment = PP_ALIGN.RIGHT
    second.level = 1
    second_run = second.add_run()
    second_run.text = "Original two"
    second_run.font.name = font_family
    second_run.font.size = Pt(18)
    second_run.font.underline = True

    output = BytesIO()
    presentation.save(output)
    cloned = clone_source_slides(
        output.getvalue(), source_slide_parts=["ppt/slides/slide1.xml"]
    )
    return cloned.package_bytes, str(textbox.shape_id)


def _slot(shape_id: str, **capacity: int) -> OoxmlTextTemplateSlot:
    return OoxmlTextTemplateSlot.model_validate(
        {
            "slotId": "fixture-v1-slide-01-body",
            "semanticRole": "body",
            "contentType": "text",
            "required": True,
            "locator": {
                "slidePart": "ppt/slides/slide1.xml",
                "shapeId": shape_id,
                "placeholderType": None,
                "relationshipId": None,
            },
            "capacity": {
                "maxChars": capacity.get("maxChars", 160),
                "maxLines": capacity.get("maxLines", 4),
                "maxParagraphs": capacity.get("maxParagraphs", 2),
                "maxBulletDepth": capacity.get("maxBulletDepth", 1),
            },
            "mutationPolicy": ["text-content"],
            "replacementPolicy": {"overflow": "fail"},
        }
    )


def _shape(package_bytes: bytes, shape_id: str) -> ET.Element:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        root = ET.fromstring(package.read("ppt/slides/slide1.xml"))
    for shape in root:
        for candidate in shape.iter():
            if candidate.tag.endswith("}cNvPr") and candidate.get("id") == shape_id:
                return shape
    raise AssertionError(f"shape {shape_id} not found")


def _paragraph_signature(shape: ET.Element) -> list[tuple[bytes, bytes]]:
    signature: list[tuple[bytes, bytes]] = []
    for paragraph in shape.findall(f".//{{{DML_NS}}}p"):
        paragraph_properties = paragraph.find(f"{{{DML_NS}}}pPr")
        run_properties = paragraph.find(f"{{{DML_NS}}}r/{{{DML_NS}}}rPr")
        signature.append(
            (
                ET.tostring(paragraph_properties) if paragraph_properties is not None else b"",
                ET.tostring(run_properties) if run_properties is not None else b"",
            )
        )
    return signature


def test_text_replacement_preserves_paragraph_and_run_formatting(tmp_path: Path) -> None:
    source, shape_id = _text_fixture(tmp_path)
    before = _shape(source, shape_id)

    result = replace_text_slot(
        source,
        slot=_slot(shape_id),
        text="새로운 첫 문장\n새로운 둘째 문장",
        available_fonts={"Aptos", "Noto Sans CJK KR"},
    )
    after = _shape(result.package_bytes, shape_id)

    assert _paragraph_signature(after) == _paragraph_signature(before)
    assert [node.text for node in after.findall(f".//{{{DML_NS}}}t")] == [
        "새로운 첫 문장",
        "새로운 둘째 문장",
    ]
    assert result.warning_codes == []
    assert validate_cloned_package(result.package_bytes) == []


def test_korean_line_break_uses_explicit_font_fallback_without_geometry_change(
    tmp_path: Path,
) -> None:
    source, shape_id = _text_fixture(tmp_path, font_family="Missing Brand Font")
    before_shape = _shape(source, shape_id)
    before_geometry = copy.deepcopy(before_shape.find(f".//{{{DML_NS}}}xfrm"))

    result = replace_text_slot(
        source,
        slot=_slot(shape_id),
        text="한글 첫 줄\n한글 둘째 줄",
        available_fonts={"Noto Sans CJK KR"},
        font_fallbacks={"Missing Brand Font": "Noto Sans CJK KR"},
    )
    after = _shape(result.package_bytes, shape_id)

    assert ET.tostring(after.find(f".//{{{DML_NS}}}xfrm")) == ET.tostring(before_geometry)
    assert len(after.findall(f".//{{{DML_NS}}}p")) == 2
    assert all(
        typeface.get("typeface") == "Noto Sans CJK KR"
        for typeface in after.findall(f".//{{{DML_NS}}}rPr/{{{DML_NS}}}latin")
    )
    assert result.font_substitutions == {
        "Missing Brand Font": "Noto Sans CJK KR"
    }
    assert validate_cloned_package(result.package_bytes) == []


def test_text_capacity_excess_fails_without_shrink_or_authored_fallback(
    tmp_path: Path,
) -> None:
    source, shape_id = _text_fixture(tmp_path)

    with pytest.raises(SlotCapacityError) as caught:
        replace_text_slot(
            source,
            slot=_slot(shape_id, maxChars=8, maxLines=1, maxParagraphs=1),
            text="한글 용량을 명백히 초과하는 두 줄\nfallback 금지",
            available_fonts={"Noto Sans CJK KR"},
        )

    assert caught.value.code == "OOXML_REFERENCE_CAPACITY_TEXT_EXCEEDED"
    assert caught.value.retryable is False
    assert caught.value.authored_fallback_created is False
    assert caught.value.package_bytes == source
