from __future__ import annotations

import hashlib
import json
from dataclasses import dataclass
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from app.ai.ooxml_reference_templates.models import OoxmlReferenceTemplateManifest
from app.ai.ooxml_reference_templates.registry import (
    LocalDirectoryObjectStorage,
    RegistryError,
    StoredObjectMetadata,
    canonical_manifest_bytes,
    ingest_reference_template,
    load_active_reference_template,
    load_repository_catalog,
)
from scripts.ingest_ooxml_reference_templates import build_dry_run_report


def _sha256(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _source_pptx(tmp_path: Path) -> Path:
    tmp_path.mkdir(parents=True, exist_ok=True)
    source_path = tmp_path / "reference.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[5])
    presentation.slides.add_slide(presentation.slide_layouts[5])
    presentation.save(source_path)
    return source_path


def _preview_png(tmp_path: Path, name: str, color: tuple[int, int, int]) -> Path:
    preview_path = tmp_path / name
    Image.new("RGB", (16, 9), color).save(preview_path)
    return preview_path


def _manifest(
    source_path: Path,
    cover_preview: Path,
    body_preview: Path,
    **overrides: object,
) -> OoxmlReferenceTemplateManifest:
    value: dict[str, object] = {
        "templateId": "operating-review",
        "version": 1,
        "status": "active",
        "sourceFormat": "pptx",
        "sourceSha256": _sha256(source_path.read_bytes()),
        "slideCount": 2,
        "canvas": {
            "aspectRatio": "16:9",
            "widthEmu": 12_192_000,
            "heightEmu": 6_858_000,
        },
        "name": "Operating Review",
        "description": "경영 보고와 KPI 중심",
        "preview": {
            "coverPreviewId": "cover",
            "coverPreviewSha256": _sha256(cover_preview.read_bytes()),
            "bodyPreviewId": "body",
            "bodyPreviewSha256": _sha256(body_preview.read_bytes()),
        },
        "sourceSlides": [
            {
                "sourceSlideId": "cover-01",
                "sourceSlidePart": "ppt/slides/slide1.xml",
                "sourceOrder": 1,
                "semanticRole": "cover",
                "relationships": {
                    "layoutPart": "ppt/slideLayouts/slideLayout1.xml",
                    "masterPart": "ppt/slideMasters/slideMaster1.xml",
                    "themePart": "ppt/theme/theme1.xml",
                },
                "capacity": {
                    "textSlotCount": 1,
                    "imageSlotCount": 0,
                    "tableSlotCount": 0,
                    "chartSlotCount": 0,
                },
                "previewId": "cover",
                "lockedInventorySha256": "a" * 64,
                "slots": [
                    {
                        "slotId": "operating-review-v1-slide-01-title",
                        "semanticRole": "title",
                        "contentType": "text",
                        "required": True,
                        "locator": {
                            "slidePart": "ppt/slides/slide1.xml",
                            "shapeId": "2",
                            "placeholderType": "title",
                            "relationshipId": None,
                        },
                        "capacity": {"maxChars": 80, "maxLines": 2},
                        "mutationPolicy": ["text-content"],
                        "replacementPolicy": {"overflow": "fail"},
                    }
                ],
            },
            {
                "sourceSlideId": "closing-02",
                "sourceSlidePart": "ppt/slides/slide2.xml",
                "sourceOrder": 2,
                "semanticRole": "closing",
                "relationships": {
                    "layoutPart": "ppt/slideLayouts/slideLayout1.xml",
                    "masterPart": "ppt/slideMasters/slideMaster1.xml",
                    "themePart": "ppt/theme/theme1.xml",
                },
                "capacity": {
                    "textSlotCount": 0,
                    "imageSlotCount": 0,
                    "tableSlotCount": 0,
                    "chartSlotCount": 0,
                },
                "previewId": "body",
                "lockedInventorySha256": "b" * 64,
                "slots": [],
            },
        ],
        "provenance": {"authorizationStatus": "approved", "inventoryVersion": 1},
    }
    value.update(overrides)
    return OoxmlReferenceTemplateManifest.model_validate(value)


@dataclass
class _StoredObject:
    content: bytes
    content_type: str


class MemoryObjectStorage:
    def __init__(self) -> None:
        self.objects: dict[str, _StoredObject] = {}
        self.created_count = 0

    def put_if_absent(self, key: str, content: bytes, content_type: str) -> bool:
        if key in self.objects:
            return False
        self.objects[key] = _StoredObject(content=content, content_type=content_type)
        self.created_count += 1
        return True

    def head_object(self, key: str) -> StoredObjectMetadata | None:
        stored = self.objects.get(key)
        if stored is None:
            return None
        return StoredObjectMetadata(
            sha256=_sha256(stored.content),
            size=len(stored.content),
            content_type=stored.content_type,
        )

    def read_object(self, key: str) -> bytes:
        return self.objects[key].content


def _active_fixture(
    tmp_path: Path,
) -> tuple[Path, Path, Path, OoxmlReferenceTemplateManifest]:
    source = _source_pptx(tmp_path)
    cover = _preview_png(tmp_path, "cover.png", (12, 34, 56))
    body = _preview_png(tmp_path, "body.png", (78, 90, 123))
    return source, cover, body, _manifest(source, cover, body)


def test_ingestion_puts_immutable_source_previews_and_canonical_manifest(
    tmp_path: Path,
) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = MemoryObjectStorage()

    result = ingest_reference_template(storage, manifest, source, cover, body)

    prefix = "system/ooxml-reference-templates/operating-review/v1"
    assert set(storage.objects) == {
        f"{prefix}/source.pptx",
        f"{prefix}/previews/cover.png",
        f"{prefix}/previews/body.png",
        f"{prefix}/manifest.json",
    }
    assert storage.read_object(f"{prefix}/source.pptx") == source.read_bytes()
    assert storage.read_object(f"{prefix}/manifest.json") == canonical_manifest_bytes(
        manifest
    )
    assert result.created_object_count == 4
    assert result.manifest_sha256 == _sha256(canonical_manifest_bytes(manifest))


def test_same_template_version_and_checksums_are_idempotent(tmp_path: Path) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = MemoryObjectStorage()
    first = ingest_reference_template(storage, manifest, source, cover, body)

    second = ingest_reference_template(storage, manifest, source, cover, body)

    assert first.created_object_count == 4
    assert second.created_object_count == 0
    assert storage.created_count == 4


def test_same_template_version_with_different_checksum_never_overwrites(
    tmp_path: Path,
) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = MemoryObjectStorage()
    ingest_reference_template(storage, manifest, source, cover, body)
    changed_source = _source_pptx(tmp_path / "changed")
    changed_presentation = Presentation(changed_source)
    changed_presentation.core_properties.title = "different immutable source"
    changed_presentation.save(changed_source)
    changed_manifest = manifest.model_copy(
        update={"source_sha256": _sha256(changed_source.read_bytes())}
    )

    with pytest.raises(RegistryError, match="IMMUTABLE_OBJECT_CONFLICT"):
        ingest_reference_template(
            storage,
            changed_manifest,
            changed_source,
            cover,
            body,
        )

    prefix = "system/ooxml-reference-templates/operating-review/v1"
    assert storage.read_object(f"{prefix}/source.pptx") == source.read_bytes()


@pytest.mark.parametrize(
    "missing_suffix",
    ["source.pptx", "previews/cover.png", "previews/body.png", "manifest.json"],
)
def test_active_load_requires_every_verified_object(
    tmp_path: Path,
    missing_suffix: str,
) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = MemoryObjectStorage()
    ingest_reference_template(storage, manifest, source, cover, body)
    prefix = "system/ooxml-reference-templates/operating-review/v1"
    del storage.objects[f"{prefix}/{missing_suffix}"]

    with pytest.raises(RegistryError, match="CATALOG_OBJECT_MISSING"):
        load_active_reference_template(storage, manifest)


@pytest.mark.parametrize(
    "tampered_suffix",
    ["source.pptx", "previews/cover.png", "previews/body.png", "manifest.json"],
)
def test_active_load_rejects_checksum_drift(
    tmp_path: Path,
    tampered_suffix: str,
) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = MemoryObjectStorage()
    ingest_reference_template(storage, manifest, source, cover, body)
    prefix = "system/ooxml-reference-templates/operating-review/v1"
    stored = storage.objects[f"{prefix}/{tampered_suffix}"]
    storage.objects[f"{prefix}/{tampered_suffix}"] = _StoredObject(
        content=stored.content + b"tampered",
        content_type=stored.content_type,
    )

    with pytest.raises(RegistryError, match="CATALOG_CHECKSUM_MISMATCH"):
        load_active_reference_template(storage, manifest)


def test_active_load_rechecks_approved_provenance_gate(tmp_path: Path) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = MemoryObjectStorage()
    ingest_reference_template(storage, manifest, source, cover, body)
    unapproved = manifest.model_copy(
        update={
            "provenance": manifest.provenance.model_copy(
                update={"authorization_status": "pending"}
            )
        }
    )

    with pytest.raises(RegistryError, match="PROVENANCE_NOT_APPROVED"):
        load_active_reference_template(storage, unapproved)


def test_local_directory_fake_has_read_after_write_immutable_semantics(
    tmp_path: Path,
) -> None:
    source, cover, body, manifest = _active_fixture(tmp_path)
    storage = LocalDirectoryObjectStorage(tmp_path / "private-test-storage")

    ingest_reference_template(storage, manifest, source, cover, body)
    loaded = load_active_reference_template(storage, manifest)

    assert loaded.template_id == "operating-review"
    assert loaded.version == 1


def test_repository_catalog_is_strict_bounded_and_contains_no_private_locator() -> None:
    catalog_path = (
        Path(__file__).parents[1]
        / "app/ai/design_library/ooxml-reference-templates/catalog.json"
    )

    catalog = load_repository_catalog(catalog_path)
    serialized = catalog_path.read_text(encoding="utf-8")

    assert len(catalog.templates) == 7
    assert sum(template.slide_count for template in catalog.templates) == 139
    assert sum(
        template.annotation_review.slot_count for template in catalog.templates
    ) == 253
    assert all(template.status == "disabled" for template in catalog.templates)
    assert all(
        template.provenance.authorization_status == "approved"
        for template in catalog.templates
    )
    assert all(
        template.annotation_review.status == "approved"
        and template.annotation_review.reviewed_on.isoformat() == "2026-07-23"
        and template.annotation_review.slide_count == template.slide_count
        and template.annotation_review.content_types == ["text"]
        for template in catalog.templates
    )
    assert all(template.activation_blockers for template in catalog.templates)
    assert all(
        "SOURCE_AUTHORIZATION_PENDING" not in template.activation_blockers
        and "SOURCE_SLIDE_ANNOTATION_MISSING" not in template.activation_blockers
        and "COVER_BODY_PREVIEW_BASELINE_MISSING"
        not in template.activation_blockers
        and "POWERPOINT_QA_PENDING" in template.activation_blockers
        and "FONT_AVAILABILITY_VALIDATION_PENDING" in template.activation_blockers
        and "PRIVATE_MANAGED_STORAGE_ADAPTER_UNCONFIGURED"
        in template.activation_blockers
        for template in catalog.templates
    )
    assert all(
        template.preview.cover_preview_sha256 is not None
        and template.preview.body_preview_sha256 is not None
        for template in catalog.templates
    )
    assert not any(
        forbidden in serialized
        for forbidden in (
            "/Users/",
            "/private/",
            "plugin/cache",
            "file://",
            "signedUrl",
            "storageKey",
            "credential",
            "password",
            "secret",
            "base64",
        )
    )
    assert json.loads(serialized)["schemaVersion"] == 2


@pytest.mark.parametrize(
    "mutation",
    [
        lambda value: value["templates"][0]["activationBlockers"].append(
            "SOURCE_AUTHORIZATION_PENDING"
        ),
        lambda value: value["templates"][0]["annotationReview"].update(
            {"slideCount": 1}
        ),
        lambda value: value["templates"][0]["activationBlockers"].remove(
            "POWERPOINT_QA_PENDING"
        ),
        lambda value: value["templates"][0]["preview"].update(
            {"coverPreviewSha256": None}
        ),
    ],
)
def test_repository_catalog_rejects_approval_and_activation_blocker_drift(
    tmp_path: Path,
    mutation: object,
) -> None:
    source = (
        Path(__file__).parents[1]
        / "app/ai/design_library/ooxml-reference-templates/catalog.json"
    )
    value = json.loads(source.read_text(encoding="utf-8"))
    assert callable(mutation)
    mutation(value)
    path = tmp_path / "catalog.json"
    path.write_text(json.dumps(value), encoding="utf-8")

    with pytest.raises(RegistryError, match="REPOSITORY_CATALOG_INVALID"):
        load_repository_catalog(path)


def test_ingestion_dry_run_reports_approved_inputs_without_enabling_catalog() -> None:
    catalog_path = (
        Path(__file__).parents[1]
        / "app/ai/design_library/ooxml-reference-templates/catalog.json"
    )
    catalog = load_repository_catalog(catalog_path)

    report = build_dry_run_report(
        catalog,
        {"sourceCount": 7, "slideCount": 139},
    )

    assert report["catalogSchemaVersion"] == 2
    assert report["approvalSummary"] == {
        "sourceAuthorizationApproved": True,
        "annotationReviewApproved": True,
        "reviewedSlideCount": 139,
        "reviewedTextSlotCount": 253,
    }
    assert report["readyForPrivateIngestion"] is False
    assert report["privateManagedStorageConfigured"] is False
    assert "SOURCE_AUTHORIZATION_PENDING" not in report["blockerCodes"]
    assert "SOURCE_SLIDE_ANNOTATION_MISSING" not in report["blockerCodes"]
    assert "POWERPOINT_QA_PENDING" in report["blockerCodes"]
    assert "FONT_AVAILABILITY_VALIDATION_PENDING" in report["blockerCodes"]
    assert all(
        template["status"] == "disabled" for template in report["templates"]
    )
