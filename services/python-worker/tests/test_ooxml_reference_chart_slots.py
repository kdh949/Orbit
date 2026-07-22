from __future__ import annotations

import copy
import hashlib
import posixpath
import zipfile
from io import BytesIO
from pathlib import PurePosixPath
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from app.ai.ooxml_reference_templates.capacity import SlotCapacityError
from app.ai.ooxml_reference_templates.chart_slots import (
    ChartSeriesData,
    replace_chart_slot,
)
from app.ai.ooxml_reference_templates.models import OoxmlChartTemplateSlot
from app.ai.ooxml_reference_templates.package import validate_cloned_package


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SHEET_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"


def _parts(package_bytes: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        return {
            item.filename: package.read(item.filename)
            for item in package.infolist()
            if not item.is_dir()
        }


def _write_parts(parts: dict[str, bytes]) -> bytes:
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as package:
        for name, content in sorted(parts.items()):
            package.writestr(name, content)
    return output.getvalue()


def _rels_part(part: str) -> str:
    path = PurePosixPath(part)
    return str(path.parent / "_rels" / f"{path.name}.rels")


def _resolve(source_part: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _fixture() -> tuple[bytes, OoxmlChartTemplateSlot, str, str]:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    data = ChartData()
    data.categories = ["North", "South"]
    data.add_series("Revenue", (10.5, 20.25), number_format="0.00")
    data.add_series("Cost", (4.25, 8.5), number_format="0.00")
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.25),
        Inches(1.5),
        Inches(9),
        Inches(4.5),
        data,
    )
    chart_frame.chart.has_legend = True
    chart_frame.chart.legend.include_in_layout = False
    output = BytesIO()
    presentation.save(output)
    package_bytes = output.getvalue()
    parts = _parts(package_bytes)

    slide_part = "ppt/slides/slide1.xml"
    slide_xml = ET.fromstring(parts[slide_part])
    frame = slide_xml.find(f".//{{{PRESENTATION_NS}}}graphicFrame")
    assert frame is not None
    shape = frame.find(f".//{{{PRESENTATION_NS}}}cNvPr")
    chart_ref = frame.find(f".//{{{CHART_NS}}}chart")
    assert shape is not None and chart_ref is not None
    relationship_id = chart_ref.attrib[f"{{{REL_NS}}}id"]
    slide_rels = ET.fromstring(parts[_rels_part(slide_part)])
    chart_relationship = next(
        item
        for item in slide_rels.findall(f"{{{PACKAGE_REL_NS}}}Relationship")
        if item.attrib.get("Id") == relationship_id
    )
    chart_part = _resolve(slide_part, chart_relationship.attrib["Target"])
    chart_rels = ET.fromstring(parts[_rels_part(chart_part)])
    workbook_relationship = next(
        item
        for item in chart_rels.findall(f"{{{PACKAGE_REL_NS}}}Relationship")
        if item.attrib.get("Type", "").endswith("/package")
    )
    workbook_part = _resolve(chart_part, workbook_relationship.attrib["Target"])
    slot = OoxmlChartTemplateSlot.model_validate(
        {
            "slotId": "chart-1",
            "semanticRole": "chart",
            "contentType": "chart",
            "required": True,
            "locator": {
                "slidePart": slide_part,
                "shapeId": shape.attrib["id"],
                "placeholderType": None,
                "relationshipId": relationship_id,
            },
            "capacity": {
                "chartType": "column",
                "maxCategories": 4,
                "maxSeries": 2,
                "workbookUpdatePolicy": "atomic",
                "workbookFingerprint": hashlib.sha256(
                    parts[workbook_part]
                ).hexdigest(),
            },
            "mutationPolicy": ["chart-data"],
            "replacementPolicy": {"overflow": "fail"},
        }
    )
    return package_bytes, slot, chart_part, workbook_part


def _cache_values(chart_xml: bytes, selector: str) -> list[str]:
    root = ET.fromstring(chart_xml)
    result: list[str] = []
    for cache in root.findall(selector):
        values = [
            str(value.text or "")
            for value in cache.findall(
                f"{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
            )
        ]
        result.extend(values)
    return result


def _workbook_values(workbook_bytes: bytes) -> dict[str, str]:
    parts = _parts(workbook_bytes)
    sheet = ET.fromstring(parts["xl/worksheets/sheet1.xml"])
    values: dict[str, str] = {}
    for cell in sheet.findall(f".//{{{SHEET_NS}}}c"):
        reference = cell.attrib["r"]
        inline = cell.find(f"{{{SHEET_NS}}}is/{{{SHEET_NS}}}t")
        value = cell.find(f"{{{SHEET_NS}}}v")
        if inline is not None:
            values[reference] = str(inline.text or "")
        elif value is not None:
            values[reference] = str(value.text or "")
    return values


def _workbook_cell_styles(workbook_bytes: bytes) -> dict[str, str | None]:
    parts = _parts(workbook_bytes)
    sheet = ET.fromstring(parts["xl/worksheets/sheet1.xml"])
    return {
        cell.attrib["r"]: cell.attrib.get("s")
        for cell in sheet.findall(f".//{{{SHEET_NS}}}c")
    }


def _chart_style_signature(chart_xml: bytes) -> tuple[bytes, ...]:
    chart = ET.fromstring(chart_xml)
    paths = [
        f".//{{{CHART_NS}}}barDir",
        f".//{{{CHART_NS}}}grouping",
        f".//{{{CHART_NS}}}varyColors",
        f".//{{{CHART_NS}}}legend",
        f".//{{{CHART_NS}}}axId",
        f".//{{{CHART_NS}}}spPr",
    ]
    return tuple(
        ET.tostring(node)
        for path in paths
        for node in chart.findall(path)
    )


def test_replaces_chart_cache_and_embedded_workbook_atomically() -> None:
    package_bytes, slot, chart_part, workbook_part = _fixture()
    before = _parts(package_bytes)

    result = replace_chart_slot(
        package_bytes,
        slot=slot,
        categories=["Enterprise", "SMB", "Public"],
        series=[
            ChartSeriesData(name="Revenue", values=(101.25, 81.5, 42.0)),
            ChartSeriesData(name="Cost", values=(50.0, 35.75, 18.25)),
        ],
    )
    after = _parts(result.package_bytes)

    assert result.warning_codes == []
    assert result.chart_part == chart_part
    assert result.workbook_part == workbook_part
    assert validate_cloned_package(result.package_bytes) == []
    assert before[slot.locator.slide_part] == after[slot.locator.slide_part]
    assert before[_rels_part(slot.locator.slide_part)] == after[
        _rels_part(slot.locator.slide_part)
    ]
    assert before[_rels_part(chart_part)] == after[_rels_part(chart_part)]
    assert _chart_style_signature(before[chart_part]) == _chart_style_signature(
        after[chart_part]
    )
    assert _cache_values(
        after[chart_part],
        f".//{{{CHART_NS}}}cat/{{{CHART_NS}}}strRef/{{{CHART_NS}}}strCache",
    ) == ["Enterprise", "SMB", "Public"] * 2
    assert _cache_values(
        after[chart_part],
        f".//{{{CHART_NS}}}val/{{{CHART_NS}}}numRef/{{{CHART_NS}}}numCache",
    ) == ["101.25", "81.5", "42", "50", "35.75", "18.25"]
    before_formats = [
        node.text
        for node in ET.fromstring(before[chart_part]).findall(
            f".//{{{CHART_NS}}}numCache/{{{CHART_NS}}}formatCode"
        )
    ]
    after_formats = [
        node.text
        for node in ET.fromstring(after[chart_part]).findall(
            f".//{{{CHART_NS}}}numCache/{{{CHART_NS}}}formatCode"
        )
    ]
    assert after_formats == before_formats
    assert _workbook_values(after[workbook_part]) == {
        "B1": "Revenue",
        "C1": "Cost",
        "A2": "Enterprise",
        "A3": "SMB",
        "A4": "Public",
        "B2": "101.25",
        "B3": "81.5",
        "B4": "42",
        "C2": "50",
        "C3": "35.75",
        "C4": "18.25",
    }
    before_workbook = _parts(before[workbook_part])
    after_workbook = _parts(after[workbook_part])
    assert before_workbook["xl/styles.xml"] == after_workbook["xl/styles.xml"]
    before_styles = _workbook_cell_styles(before[workbook_part])
    after_styles = _workbook_cell_styles(after[workbook_part])
    assert after_styles["B4"] == before_styles["B3"]
    assert after_styles["C4"] == before_styles["C3"]
    formulas = [
        node.text
        for node in ET.fromstring(after[chart_part]).findall(
            f".//{{{CHART_NS}}}f"
        )
    ]
    assert "Sheet1!$A$2:$A$4" in formulas
    assert "Sheet1!$B$2:$B$4" in formulas
    assert "Sheet1!$C$2:$C$4" in formulas


@pytest.mark.parametrize(
    ("categories", "series", "code"),
    [
        (
            ["A", "B", "C", "D", "E"],
            [
                ChartSeriesData(name="Revenue", values=(1, 2, 3, 4, 5)),
                ChartSeriesData(name="Cost", values=(1, 2, 3, 4, 5)),
            ],
            "OOXML_REFERENCE_CAPACITY_CHART_CATEGORIES",
        ),
        (
            ["A", "B"],
            [
                ChartSeriesData(name="A", values=(1, 2)),
                ChartSeriesData(name="B", values=(1, 2)),
                ChartSeriesData(name="C", values=(1, 2)),
            ],
            "OOXML_REFERENCE_CAPACITY_CHART_SERIES",
        ),
    ],
)
def test_capacity_mismatch_fails_closed(
    categories: list[str], series: list[ChartSeriesData], code: str
) -> None:
    package_bytes, slot, _, _ = _fixture()

    with pytest.raises(SlotCapacityError) as caught:
        replace_chart_slot(
            package_bytes,
            slot=slot,
            categories=categories,
            series=series,
        )

    assert caught.value.code == code
    assert caught.value.package_bytes == package_bytes
    assert not caught.value.authored_fallback_created


def test_workbook_fingerprint_drift_fails_closed() -> None:
    package_bytes, slot, _, _ = _fixture()
    drifted = slot.model_copy(deep=True)
    drifted.capacity.workbook_fingerprint = "0" * 64

    with pytest.raises(SlotCapacityError) as caught:
        replace_chart_slot(
            package_bytes,
            slot=drifted,
            categories=["A", "B"],
            series=[
                ChartSeriesData(name="Revenue", values=(1, 2)),
                ChartSeriesData(name="Cost", values=(3, 4)),
            ],
        )

    assert caught.value.code == "OOXML_REFERENCE_CHART_WORKBOOK_FINGERPRINT_MISMATCH"
    assert caught.value.package_bytes == package_bytes


def test_external_workbook_is_preserve_only() -> None:
    package_bytes, slot, chart_part, _ = _fixture()
    parts = _parts(package_bytes)
    rels_part = _rels_part(chart_part)
    relationships = ET.fromstring(parts[rels_part])
    package_relationship = next(
        item
        for item in relationships.findall(f"{{{PACKAGE_REL_NS}}}Relationship")
        if item.attrib.get("Type", "").endswith("/package")
    )
    package_relationship.attrib["TargetMode"] = "External"
    package_relationship.attrib["Target"] = "https://example.invalid/data.xlsx"
    parts[rels_part] = ET.tostring(
        relationships, encoding="utf-8", xml_declaration=True
    )
    modified = _write_parts(parts)

    with pytest.raises(SlotCapacityError) as caught:
        replace_chart_slot(
            modified,
            slot=slot,
            categories=["A", "B"],
            series=[
                ChartSeriesData(name="Revenue", values=(1, 2)),
                ChartSeriesData(name="Cost", values=(3, 4)),
            ],
        )

    assert caught.value.code == "OOXML_REFERENCE_CHART_WORKBOOK_UNSUPPORTED"
    assert caught.value.package_bytes == modified


def test_unsupported_chart_and_locator_drift_fail_closed() -> None:
    package_bytes, slot, chart_part, _ = _fixture()
    parts = _parts(package_bytes)
    chart = ET.fromstring(parts[chart_part])
    chart_type = chart.find(f".//{{{CHART_NS}}}barChart")
    assert chart_type is not None
    chart_type.tag = f"{{{CHART_NS}}}areaChart"
    parts[chart_part] = ET.tostring(chart, encoding="utf-8", xml_declaration=True)
    unsupported = _write_parts(parts)

    with pytest.raises(SlotCapacityError) as chart_error:
        replace_chart_slot(
            unsupported,
            slot=slot,
            categories=["A", "B"],
            series=[
                ChartSeriesData(name="Revenue", values=(1, 2)),
                ChartSeriesData(name="Cost", values=(3, 4)),
            ],
        )
    assert chart_error.value.code == "OOXML_REFERENCE_CHART_TYPE_UNSUPPORTED"
    assert chart_error.value.package_bytes == unsupported

    bad_locator = copy.deepcopy(slot)
    bad_locator.locator.relationship_id = "rIdMissing"
    with pytest.raises(SlotCapacityError) as locator_error:
        replace_chart_slot(
            package_bytes,
            slot=bad_locator,
            categories=["A", "B"],
            series=[
                ChartSeriesData(name="Revenue", values=(1, 2)),
                ChartSeriesData(name="Cost", values=(3, 4)),
            ],
        )
    assert locator_error.value.code == "OOXML_REFERENCE_CHART_LOCATOR_INVALID"
    assert locator_error.value.package_bytes == package_bytes


def test_source_series_capacity_and_value_shape_are_bounded() -> None:
    package_bytes, slot, _, _ = _fixture()

    with pytest.raises(SlotCapacityError) as series_error:
        replace_chart_slot(
            package_bytes,
            slot=slot,
            categories=["A", "B"],
            series=[ChartSeriesData(name="Revenue", values=(1, 2))],
        )
    assert series_error.value.code == "OOXML_REFERENCE_CHART_STRUCTURE_UNSUPPORTED"

    with pytest.raises(SlotCapacityError) as value_error:
        replace_chart_slot(
            package_bytes,
            slot=slot,
            categories=["A", "B"],
            series=[
                ChartSeriesData(name="Revenue", values=(1,)),
                ChartSeriesData(name="Cost", values=(3, 4)),
            ],
        )
    assert value_error.value.code == "OOXML_REFERENCE_CHART_DATA_INVALID"
