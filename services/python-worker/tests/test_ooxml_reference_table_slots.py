from __future__ import annotations

import copy
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.ai.ooxml_reference_templates.capacity import SlotCapacityError
from app.ai.ooxml_reference_templates.models import OoxmlTableTemplateSlot
from app.ai.ooxml_reference_templates.package import validate_cloned_package
from app.ai.ooxml_reference_templates.table_slots import replace_table_slot
from app.ai.pptx_ooxml_generation import generate_pptx_ooxml, sync_pptx_ooxml
from app.ai.pptx_ooxml_vector_importer import table_cell_fingerprint


DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
SLIDE_PART = "ppt/slides/slide1.xml"


def _table_fixture() -> tuple[bytes, str]:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    frame = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(7),
        Inches(2),
    )
    for row_index, row in enumerate(frame.table.rows):
        for column_index, cell in enumerate(row.cells):
            cell.text = f"Source {row_index},{column_index}"

    output = BytesIO()
    presentation.save(output)
    package_bytes = _rewrite_slide(
        output.getvalue(),
        lambda root: _add_cell_style(first_descendant(root, "tc")),
    )
    return package_bytes, str(frame.shape_id)


def _slot(
    package_bytes: bytes,
    shape_id: str,
    *,
    row_count: int = 2,
    column_count: int = 2,
    editable_indexes: tuple[tuple[int, int], ...] = ((0, 0), (0, 1), (1, 0), (1, 1)),
) -> OoxmlTableTemplateSlot:
    table = _table(package_bytes, shape_id)
    rows = direct_children(table, "tr")
    cells = [direct_children(row, "tc") for row in rows]
    return OoxmlTableTemplateSlot.model_validate(
        {
            "slotId": "fixture-v1-slide-01-table",
            "semanticRole": "table",
            "contentType": "table",
            "required": True,
            "locator": {
                "slidePart": SLIDE_PART,
                "shapeId": shape_id,
                "placeholderType": None,
                "relationshipId": None,
            },
            "capacity": {
                "rowCount": row_count,
                "columnCount": column_count,
                "mergedCellPolicy": "preserve",
                "editableCells": [
                    {
                        "rowIndex": row_index,
                        "columnIndex": column_index,
                        "fingerprint": table_cell_fingerprint(
                            cells[row_index][column_index]
                        ),
                    }
                    for row_index, column_index in editable_indexes
                ],
            },
            "mutationPolicy": ["table-cell-text"],
            "replacementPolicy": {"overflow": "fail"},
        }
    )


def test_table_replacement_changes_only_text_and_preserves_table_structure() -> None:
    source, shape_id = _table_fixture()
    before_frame = _frame(source, shape_id)
    before_signature = _structure_signature(before_frame)

    result = replace_table_slot(
        source,
        slot=_slot(source, shape_id),
        rows=[["분기", "매출"], ["1분기", "120"]],
    )

    after_frame = _frame(result.package_bytes, shape_id)
    after = first_descendant(after_frame, "tbl")
    assert _table_text(after) == [["분기", "매출"], ["1분기", "120"]]
    assert _structure_signature(after_frame) == before_signature
    assert _changed_parts(source, result.package_bytes) == [SLIDE_PART]
    assert result.warning_codes == []
    assert validate_cloned_package(result.package_bytes) == []


@pytest.mark.parametrize(
    ("row_count", "column_count", "rows", "expected_code"),
    [
        (3, 2, [["A", "B"], ["C", "D"]], "OOXML_REFERENCE_CAPACITY_TABLE_SHAPE"),
        (2, 3, [["A", "B"], ["C", "D"]], "OOXML_REFERENCE_CAPACITY_TABLE_SHAPE"),
        (2, 2, [["A", "B"], ["C"]], "OOXML_REFERENCE_CAPACITY_TABLE_SHAPE"),
    ],
)
def test_table_replacement_rejects_capacity_or_rectangular_input_mismatch(
    row_count: int,
    column_count: int,
    rows: list[list[str]],
    expected_code: str,
) -> None:
    source, shape_id = _table_fixture()

    with pytest.raises(SlotCapacityError) as caught:
        replace_table_slot(
            source,
            slot=_slot(
                source,
                shape_id,
                row_count=row_count,
                column_count=column_count,
            ),
            rows=rows,
        )

    assert caught.value.code == expected_code
    assert caught.value.package_bytes == source
    assert caught.value.authored_fallback_created is False


def test_table_replacement_rejects_locator_fingerprint_drift() -> None:
    source, shape_id = _table_fixture()
    slot = _slot(source, shape_id)
    slot.capacity.editable_cells[0].fingerprint = "0" * 64

    with pytest.raises(SlotCapacityError) as caught:
        replace_table_slot(
            source,
            slot=slot,
            rows=[["A", "B"], ["C", "D"]],
        )

    assert caught.value.code == "OOXML_REFERENCE_TABLE_LOCATOR_FINGERPRINT_MISMATCH"
    assert caught.value.package_bytes == source


@pytest.mark.parametrize("mutation", ["jagged", "merged", "indirect"])
def test_table_replacement_rejects_unsupported_source_structure(mutation: str) -> None:
    source, shape_id = _table_fixture()
    slot = _slot(source, shape_id)

    def mutate(root: ET.Element) -> None:
        table = first_descendant(root, "tbl")
        if mutation == "jagged":
            direct_children(table, "tr")[-1].remove(
                direct_children(direct_children(table, "tr")[-1], "tc")[-1]
            )
        elif mutation == "merged":
            direct_children(direct_children(table, "tr")[0], "tc")[0].set(
                "gridSpan", "2"
            )
        else:
            graphic_data = first_descendant(root, "graphicData")
            graphic_data.set("uri", "urn:unsupported:wrapper")

    drifted = _rewrite_slide(source, mutate)
    with pytest.raises(SlotCapacityError) as caught:
        replace_table_slot(
            drifted,
            slot=slot,
            rows=[["A", "B"], ["C", "D"]],
        )

    assert caught.value.code == "OOXML_REFERENCE_TABLE_STRUCTURE_UNSUPPORTED"
    assert caught.value.package_bytes == drifted


def test_table_replacement_rejects_changes_outside_editable_cells() -> None:
    source, shape_id = _table_fixture()

    with pytest.raises(SlotCapacityError) as caught:
        replace_table_slot(
            source,
            slot=_slot(source, shape_id, editable_indexes=((0, 0),)),
            rows=[["Allowed", "Not allowed"], ["Source 1,0", "Source 1,1"]],
        )

    assert caught.value.code == "OOXML_REFERENCE_TABLE_CELL_NOT_EDITABLE"
    assert caught.value.package_bytes == source


def test_generated_table_supports_warning_free_single_cell_targeted_sync(
    tmp_path: Path,
) -> None:
    source, shape_id = _table_fixture()
    replaced = replace_table_slot(
        source,
        slot=_slot(source, shape_id),
        rows=[["분기", "매출"], ["1분기", "120"]],
    )
    generated_path = tmp_path / "generated-table.pptx"
    generated_path.write_bytes(replaced.package_bytes)
    generated = generate_pptx_ooxml(generated_path, "file_table", render=False)
    element = next(
        item
        for item in generated.blueprint["slides"][0]["elements"]
        if item.get("type") == "table"
    )
    source_record = next(
        item
        for item in generated.template_blueprint["slides"][0]["elementSources"]
        if item.get("elementId") == element["elementId"]
    )
    assert source_record["ooxmlEditCapabilities"]["tableCellText"] is True
    props = copy.deepcopy(element["props"])
    props["rows"][1][1]["text"] = "135"

    synced = sync_pptx_ooxml(
        generated_path,
        template_blueprint=generated.template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": generated.template_blueprint["slides"][0]["slideId"],
                "elementId": element["elementId"],
                "props": props,
            }
        ],
        deck_canvas=generated.canvas,
        synced_deck_version=2,
        render=False,
    )

    assert synced.unsupported_operations == []
    assert synced.warnings == []
    assert [item.operation_type for item in synced.applied_operations] == [
        "update_element_props"
    ]


def _add_cell_style(cell: ET.Element) -> None:
    cell_properties = next(
        (child for child in cell if child.tag.endswith("}tcPr")), None
    )
    assert cell_properties is not None
    cell_properties.set("anchor", "ctr")
    line = ET.SubElement(cell_properties, f"{{{DML_NS}}}lnL", {"w": "12700"})
    solid = ET.SubElement(line, f"{{{DML_NS}}}solidFill")
    ET.SubElement(solid, f"{{{DML_NS}}}srgbClr", {"val": "112233"})


def _rewrite_slide(package_bytes: bytes, mutation: object) -> bytes:
    source = BytesIO(package_bytes)
    output = BytesIO()
    with (
        zipfile.ZipFile(source, "r") as archive,
        zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as target,
    ):
        for info in archive.infolist():
            payload = archive.read(info.filename)
            if info.filename == SLIDE_PART:
                root = ET.fromstring(payload)
                assert callable(mutation)
                mutation(root)
                payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info, payload)
    return output.getvalue()


def _table(package_bytes: bytes, shape_id: str) -> ET.Element:
    return first_descendant(_frame(package_bytes, shape_id), "tbl")


def _frame(package_bytes: bytes, shape_id: str) -> ET.Element:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        root = ET.fromstring(package.read(SLIDE_PART))
    for frame in root.iter():
        if not frame.tag.endswith("}graphicFrame"):
            continue
        non_visual = next(
            (item for item in frame.iter() if item.tag.endswith("}cNvPr")), None
        )
        if non_visual is not None and non_visual.get("id") == shape_id:
            return frame
    raise AssertionError(f"table shape {shape_id} not found")


def _table_text(table: ET.Element) -> list[list[str]]:
    return [
        [
            "".join(node.text or "" for node in cell.iter() if node.tag.endswith("}t"))
            for cell in direct_children(row, "tc")
        ]
        for row in direct_children(table, "tr")
    ]


def _structure_signature(root: ET.Element) -> bytes:
    payload = copy.deepcopy(root)
    for node in payload.iter():
        if node.tag.endswith("}t"):
            node.text = ""
            node.attrib.pop("{http://www.w3.org/XML/1998/namespace}space", None)
    return canonical(payload)


def _changed_parts(before: bytes, after: bytes) -> list[str]:
    with (
        zipfile.ZipFile(BytesIO(before), "r") as left,
        zipfile.ZipFile(BytesIO(after), "r") as right,
    ):
        return sorted(
            name for name in left.namelist() if left.read(name) != right.read(name)
        )


def first_descendant(root: ET.Element, local_name: str) -> ET.Element:
    return next(item for item in root.iter() if item.tag.endswith(f"}}{local_name}"))


def direct_children(root: ET.Element, local_name: str) -> list[ET.Element]:
    return [item for item in root if item.tag.endswith(f"}}{local_name}")]


def canonical(root: ET.Element) -> bytes:
    value = ET.canonicalize(ET.tostring(root, encoding="unicode"), with_comments=False)
    return value.encode("utf-8")
