from __future__ import annotations

import base64
from io import BytesIO

from pptx import Presentation

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.deck_generation.design_pack_registry import load_design_pack_catalog
from app.ai.deck_generation.design_pack_selector import (
    DESIGN_PACK_DIRECTORY,
    select_system_design_pack,
)
from app.ai.deck_generation.models import GenerateDeckRequest
from app.ai.deck_generation.pipeline import analyze_input
from app.ai.design_pack_layouts.editorial_insight import (
    compile_editorial_insight_layout,
)
from app.ai.design_program import DeckDesignProgram


def test_editorial_profile_and_market_insight_select_pack() -> None:
    for raw_input in (editorial_input(profile=True), editorial_input(profile=False)):
        selection = select_system_design_pack(raw_input, golden_slides())
        registry = load_design_pack_catalog(DESIGN_PACK_DIRECTORY)
        layouts = {layout.layout_id: layout for layout in registry.layouts}
        silhouettes = [
            layouts[layout_id].silhouette_id for layout_id in selection.layout_ids
        ]

        assert selection.pack_id == "editorial-insight"
        assert selection.fallback_used is False
        assert len(set(silhouettes)) >= 4
        assert all(
            left != right for left, right in zip(silhouettes, silhouettes[1:])
        )


def test_statement_has_one_primary_claim_and_no_invented_metric() -> None:
    program = design_program()
    compiled = compile_editorial_insight_layout(
        "editorial-statement-01",
        program.slides[0],
        slide("solution", 1, "시장이 바뀌는 단 하나의 이유"),
        program,
    )

    assert sum(
        element["elementId"] == compiled.primary_focal_element_id
        for element in compiled.elements
    ) == 1
    assert all(element.get("role") != "metric" for element in compiled.elements)
    response = export_deck_pptx(
        DeckPptxExportRequest(deck=export_deck([compiled.elements]))
    )
    assert len(Presentation(BytesIO(base64.b64decode(response.content_base64))).slides) == 1


def test_evidence_without_typed_metric_stays_non_metric() -> None:
    program = design_program()
    compiled = compile_editorial_insight_layout(
        "editorial-evidence-01",
        program.slides[0],
        slide("data", 3, "시장 근거"),
        program,
    )

    assert all(element.get("role") != "metric" for element in compiled.elements)


def editorial_input(*, profile: bool):
    request = GenerateDeckRequest.model_validate(
        {
            "projectId": "project_demo_1",
            "topic": "2027 시장 동향과 고객 행동 인사이트",
            "metadata": {"purpose": "inform"},
            "design": {
                "profile": "editorial" if profile else None,
                "mediaPolicy": "minimal",
            },
        }
    )
    return analyze_input(request).model_copy(
        update={"presentation_profile": "general-inform"}
    )


def golden_slides() -> list[dict]:
    return [
        slide("cover", 0, "시장 동향"),
        slide("solution", 1, "핵심 논지"),
        slide("data", 3, "관찰 근거"),
        slide("comparison", 3, "트렌드 변화"),
        slide("quote", 1, "현장의 목소리"),
        slide("solution", 2, "시사점"),
        slide("data", 4, "추가 근거"),
        slide("summary", 2, "결론"),
    ]


def slide(slide_type: str, item_count: int, title: str) -> dict:
    return {
        "title": title,
        "message": "하나의 중심 주장과 검증 가능한 근거를 연결합니다",
        "slideType": slide_type,
        "contentItems": [
            {"contentItemId": f"item-{index}", "text": f"시장 근거 {index}"}
            for index in range(1, item_count + 1)
        ],
        "typedMetrics": [],
        "mediaIntent": {"kind": "none"},
    }


def design_program() -> DeckDesignProgram:
    return DeckDesignProgram.model_validate(
        {
            "visualConcept": "Editorial insight",
            "paletteRoles": {
                "dominant": "#FAFAF9",
                "surface": "#F5F5F4",
                "text": "#1C1917",
                "focal": "#C2410C",
                "secondary": "#334155",
            },
            "typography": {
                "headingFont": "Pretendard",
                "bodyFont": "Pretendard",
                "typeScale": {"cover": 72, "title": 54, "body": 28},
            },
            "backgroundSequence": ["light"],
            "imageStyle": "Evidence-led editorial crop",
            "surfaceStyle": "Asymmetric editorial fields",
            "slides": [
                {
                    "order": 1,
                    "compositionId": "statement-poster",
                    "variant": "light",
                    "backgroundMode": "light",
                    "focalType": "message",
                    "assetRole": "none",
                    "requiredAsset": False,
                }
            ],
        }
    )


def export_deck(slide_elements: list[list[dict]]) -> dict:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {
            "backgroundColor": "#FAFAF9",
            "textColor": "#1C1917",
            "fontFamily": "Pretendard",
        },
        "slides": [
            {
                "order": index,
                "style": {"backgroundColor": "#FAFAF9"},
                "speakerNotes": "시장 근거와 시사점을 설명합니다.",
                "elements": elements,
            }
            for index, elements in enumerate(slide_elements, start=1)
        ],
    }
