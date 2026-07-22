from __future__ import annotations

import base64
from io import BytesIO

import pytest
from pptx import Presentation

from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx
from app.ai.deck_generation.design_pack_registry import load_design_pack_catalog
from app.ai.deck_generation.design_pack_selector import select_system_design_pack
from app.ai.deck_generation.models import GenerateDeckRequest
from app.ai.deck_generation.pipeline import analyze_input
from app.ai.design_pack_layouts.executive_review import (
    compile_executive_review_layout,
)
from app.ai.design_program import DeckDesignProgram


def test_executive_report_selects_review_pack_with_four_silhouettes() -> None:
    selection = select_system_design_pack(report_input(), golden_slides())
    registry = load_design_pack_catalog(selection_catalog_directory())
    layouts = {layout.layout_id: layout for layout in registry.layouts}
    silhouettes = [layouts[layout_id].silhouette_id for layout_id in selection.layout_ids]

    assert selection.pack_id == "executive-review"
    assert selection.fallback_used is False
    assert len(set(silhouettes)) >= 4
    assert all(left != right for left, right in zip(silhouettes, silhouettes[1:]))


def test_executive_table_and_grounded_chart_are_native_and_exportable() -> None:
    program = design_program()
    direction = program.slides[0]
    table = compile_executive_review_layout(
        "executive-table-01",
        direction,
        golden_slides()[2],
        program,
    )
    chart = compile_executive_review_layout(
        "executive-chart-01",
        direction,
        golden_slides()[3],
        program,
    )

    assert any(element["type"] == "table" for element in table.elements)
    assert any(element["type"] == "chart" for element in chart.elements)
    response = export_deck_pptx(
        DeckPptxExportRequest(deck=export_deck([table.elements, chart.elements]))
    )
    presentation = Presentation(BytesIO(base64.b64decode(response.content_base64)))
    assert len(presentation.slides) == 2


def test_executive_chart_rejects_ungrounded_or_oversized_data() -> None:
    program = design_program()
    direction = program.slides[0]
    for metrics in ([], typed_metrics(5)):
        candidate = slide("chart", 2, metrics=metrics)
        with pytest.raises(ValueError, match="grounded typed metrics"):
            compile_executive_review_layout(
                "executive-chart-01",
                direction,
                candidate,
                program,
            )


def report_input():
    request = GenerateDeckRequest.model_validate(
        {
            "projectId": "project_demo_1",
            "topic": "분기 운영 리뷰",
            "metadata": {"audience": "executive", "purpose": "report"},
            "design": {"mediaPolicy": "minimal"},
        }
    )
    return analyze_input(request).model_copy(
        update={"presentation_profile": "executive-report"}
    )


def selection_catalog_directory():
    from app.ai.deck_generation.design_pack_selector import DESIGN_PACK_DIRECTORY

    return DESIGN_PACK_DIRECTORY


def golden_slides() -> list[dict]:
    return [
        slide("cover", 0),
        slide("data", 3, metrics=typed_metrics(3)),
        slide("data", 4),
        slide("chart", 3, metrics=typed_metrics(3)),
        slide("comparison", 3),
        slide("solution", 3),
        slide("quote", 1),
        slide("data", 2, metrics=typed_metrics(2)),
        slide("summary", 1),
    ]


def slide(
    slide_type: str,
    item_count: int,
    *,
    metrics: list[dict[str, str]] | None = None,
) -> dict:
    return {
        "title": f"{slide_type} title",
        "message": "경영 판단에 필요한 핵심 결론",
        "slideType": slide_type,
        "contentItems": [
            {"contentItemId": f"item-{index}", "text": f"운영 근거 {index}"}
            for index in range(1, item_count + 1)
        ],
        "typedMetrics": metrics or [],
        "mediaIntent": {"kind": "none"},
    }


def typed_metrics(count: int) -> list[dict[str, str]]:
    return [
        {
            "value": str(10 + index),
            "unit": "%",
            "label": f"KPI {index}",
            "sourceRef": f"source:{index}",
        }
        for index in range(1, count + 1)
    ]


def design_program() -> DeckDesignProgram:
    return DeckDesignProgram.model_validate(
        {
            "visualConcept": "Executive review",
            "paletteRoles": {
                "dominant": "#FFFFFF",
                "surface": "#F1F5F9",
                "text": "#0F172A",
                "focal": "#2563EB",
                "secondary": "#0F766E",
            },
            "typography": {
                "headingFont": "Pretendard",
                "bodyFont": "Pretendard",
                "typeScale": {"cover": 72, "title": 56, "body": 32},
            },
            "backgroundSequence": ["light"],
            "imageStyle": "Evidence-first",
            "surfaceStyle": "Flat executive surfaces",
            "slides": [
                {
                    "order": 1,
                    "compositionId": "editorial-split",
                    "variant": "light",
                    "backgroundMode": "light",
                    "focalType": "evidence",
                    "assetRole": "none",
                    "requiredAsset": False,
                }
            ],
        }
    )


def export_deck(slide_elements: list[list[dict]]) -> dict:
    return {
        "canvas": {"width": 1920, "height": 1080},
        "theme": {
            "backgroundColor": "#FFFFFF",
            "textColor": "#0F172A",
            "fontFamily": "Pretendard",
        },
        "slides": [
            {
                "order": index,
                "style": {"backgroundColor": "#FFFFFF"},
                "speakerNotes": "경영 보고 내용을 설명합니다.",
                "elements": elements,
            }
            for index, elements in enumerate(slide_elements, start=1)
        ],
    }
