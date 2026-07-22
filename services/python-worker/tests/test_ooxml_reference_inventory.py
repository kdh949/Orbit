from __future__ import annotations

import hashlib
import io
import json
import subprocess
import sys
import zipfile
from collections.abc import Callable
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.ai.ooxml_reference_templates.inventory import (
    EXPECTED_REFERENCE_SLIDE_COUNTS,
    InventoryLimits,
    InventoryValidationError,
    ReferenceSource,
    build_reference_inventory,
)


def _source_pptx(tmp_path: Path, name: str = "source.pptx") -> Path:
    image_path = tmp_path / "private-image.png"
    Image.new("RGB", (8, 8), (12, 34, 56)).save(image_path)

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "TOP SECRET SOURCE TEXT"
    first.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(1))
    first.shapes.add_picture(str(image_path), Inches(5), Inches(1), Inches(1), Inches(1))

    presentation.slides.add_slide(presentation.slide_layouts[5])

    source_path = tmp_path / name
    presentation.save(source_path)
    with zipfile.ZipFile(source_path, "r") as package:
        content_types = package.read("[Content_Types].xml").replace(
            b"</Types>",
            b'<Override PartName="/ppt/charts/chart1.xml" '
            b'ContentType="application/vnd.openxmlformats-officedocument.'
            b'drawingml.chart+xml"/></Types>',
        )
    _rewrite_package(
        source_path,
        source_path,
        replacements={"[Content_Types].xml": content_types},
        additions={
            "ppt/charts/chart1.xml": (
                b'<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/'
                b'drawingml/2006/chart"/>'
            )
        },
    )
    return source_path


def _rewrite_package(
    source: Path,
    target: Path,
    replacements: dict[str, bytes] | None = None,
    additions: dict[str, bytes] | None = None,
) -> Path:
    replacements = replacements or {}
    additions = additions or {}
    with zipfile.ZipFile(source, "r") as original:
        entries = {
            item.filename: original.read(item.filename)
            for item in original.infolist()
        }
    entries.update(replacements)
    entries.update(additions)
    with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as rewritten:
        for filename, content in entries.items():
            rewritten.writestr(filename, content)
    return target


def _replace_xml(source: Path, target: Path, part: str, transform: Callable[[bytes], bytes]) -> Path:
    with zipfile.ZipFile(source, "r") as package:
        replacement = transform(package.read(part))
    return _rewrite_package(source, target, replacements={part: replacement})


def test_inventory_records_bounded_package_metadata_without_source_content(
    tmp_path: Path,
) -> None:
    source_path = _source_pptx(tmp_path)
    report = build_reference_inventory(
        [ReferenceSource("fixture", source_path)],
        expected_slide_counts={"fixture": 2},
    )

    assert report["schemaVersion"] == 1
    assert report["sourceCount"] == 1
    assert report["slideCount"] == 2
    source = report["sources"][0]
    assert source["templateId"] == "fixture"
    assert source["sha256"] == hashlib.sha256(source_path.read_bytes()).hexdigest()
    assert source["slideCount"] == 2
    assert source["masterCount"] >= 1
    assert source["layoutCount"] >= 1
    assert source["themeCount"] >= 1
    assert source["fontCount"] >= 1
    assert source["mediaCount"] == 1
    assert source["chartCount"] == 1
    assert source["tableCount"] == 1
    assert source["smartArtCount"] == 0
    assert source["animationCount"] == 0
    serialized = json.dumps(report, ensure_ascii=False)
    assert str(tmp_path) not in serialized
    assert "TOP SECRET SOURCE TEXT" not in serialized
    assert "private-image.png" not in serialized
    assert "<p:sld" not in serialized


def test_inventory_expected_baseline_is_exactly_seven_templates_and_139_slides() -> None:
    assert set(EXPECTED_REFERENCE_SLIDE_COUNTS) == {
        "simple-light",
        "simple-dark",
        "operating-review",
        "business-review",
        "project-kickoff",
        "team-alignment",
        "market-trends-report",
    }
    assert sum(EXPECTED_REFERENCE_SLIDE_COUNTS.values()) == 139


def test_inventory_counts_smart_art_and_animation_parts(tmp_path: Path) -> None:
    source_path = _source_pptx(tmp_path)
    with zipfile.ZipFile(source_path, "r") as package:
        content_types = package.read("[Content_Types].xml").replace(
            b"</Types>",
            b'<Override PartName="/ppt/diagrams/data1.xml" '
            b'ContentType="application/vnd.openxmlformats-officedocument.'
            b'drawingml.diagramData+xml"/></Types>',
        )
        slide = package.read("ppt/slides/slide1.xml").replace(
            b"</p:sld>", b"<p:timing/></p:sld>"
        )
    enriched_path = _rewrite_package(
        source_path,
        tmp_path / "smart-art-animation.pptx",
        replacements={
            "[Content_Types].xml": content_types,
            "ppt/slides/slide1.xml": slide,
        },
        additions={
            "ppt/diagrams/data1.xml": (
                b'<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/'
                b'drawingml/2006/diagram"/>'
            )
        },
    )

    report = build_reference_inventory(
        [ReferenceSource("fixture", enriched_path)],
        expected_slide_counts={"fixture": 2},
    )

    assert report["sources"][0]["smartArtCount"] == 1
    assert report["sources"][0]["animationCount"] == 1


def test_inventory_rejects_source_or_slide_count_drift(tmp_path: Path) -> None:
    source_path = _source_pptx(tmp_path)

    with pytest.raises(InventoryValidationError, match="inventory_drift"):
        build_reference_inventory(
            [ReferenceSource("fixture", source_path)],
            expected_slide_counts={"fixture": 3},
        )


def test_inventory_rejects_duplicate_template_ids(tmp_path: Path) -> None:
    source_path = _source_pptx(tmp_path)

    with pytest.raises(InventoryValidationError, match="inventory_drift"):
        build_reference_inventory(
            [
                ReferenceSource("fixture", source_path),
                ReferenceSource("fixture", source_path),
            ],
            expected_slide_counts={"fixture": 2},
        )


@pytest.mark.parametrize(
    ("unsafe_name", "error_code"),
    [
        ("../escape.xml", "zip_path_traversal"),
        ("/absolute.xml", "zip_path_traversal"),
        ("C:\\escape.xml", "zip_path_traversal"),
    ],
)
def test_security_preflight_rejects_zip_path_traversal(
    tmp_path: Path,
    unsafe_name: str,
    error_code: str,
) -> None:
    source_path = _source_pptx(tmp_path)
    unsafe_path = _rewrite_package(
        source_path,
        tmp_path / "unsafe.pptx",
        additions={unsafe_name: b"unsafe"},
    )

    with pytest.raises(InventoryValidationError, match=error_code):
        build_reference_inventory(
            [ReferenceSource("fixture", unsafe_path)],
            expected_slide_counts={"fixture": 2},
        )


def test_security_preflight_rejects_duplicate_zip_parts(tmp_path: Path) -> None:
    source_path = _source_pptx(tmp_path)
    duplicate_path = tmp_path / "duplicate.pptx"
    duplicate_path.write_bytes(source_path.read_bytes())
    with zipfile.ZipFile(duplicate_path, "a") as package:
        package.writestr("[Content_Types].xml", b"duplicate")

    with pytest.raises(InventoryValidationError, match="duplicate_zip_part"):
        build_reference_inventory(
            [ReferenceSource("fixture", duplicate_path)],
            expected_slide_counts={"fixture": 2},
        )


@pytest.mark.parametrize(
    ("limits", "error_code"),
    [
        (InventoryLimits(max_archive_bytes=1), "archive_size_limit"),
        (InventoryLimits(max_parts=1), "part_count_limit"),
        (InventoryLimits(max_total_uncompressed_bytes=1), "uncompressed_size_limit"),
        (InventoryLimits(max_part_uncompressed_bytes=1), "part_size_limit"),
        (InventoryLimits(max_compression_ratio=1.0), "compression_ratio_limit"),
    ],
)
def test_security_preflight_rejects_zip_bomb_resource_limits(
    tmp_path: Path,
    limits: InventoryLimits,
    error_code: str,
) -> None:
    source_path = _source_pptx(tmp_path)

    with pytest.raises(InventoryValidationError, match=error_code):
        build_reference_inventory(
            [ReferenceSource("fixture", source_path)],
            expected_slide_counts={"fixture": 2},
            limits=limits,
        )


@pytest.mark.parametrize(
    "part_name",
    [
        "ppt/vbaProject.bin",
        "ppt/activeX/activeX1.bin",
        "ppt/embeddings/oleObject1.bin",
        "ppt/embeddings/embeddedPackage1.xlsx",
    ],
)
def test_security_preflight_rejects_active_content_and_embedded_objects(
    tmp_path: Path,
    part_name: str,
) -> None:
    source_path = _source_pptx(tmp_path)
    unsafe_path = _rewrite_package(
        source_path,
        tmp_path / "active-content.pptx",
        additions={part_name: b"unsafe"},
    )

    with pytest.raises(InventoryValidationError, match="prohibited_active_content"):
        build_reference_inventory(
            [ReferenceSource("fixture", unsafe_path)],
            expected_slide_counts={"fixture": 2},
        )


def test_security_preflight_allows_only_internal_chart_workbook_after_nested_scan(
    tmp_path: Path,
) -> None:
    source_path = _source_pptx(tmp_path)
    workbook_buffer = io.BytesIO()
    with zipfile.ZipFile(workbook_buffer, "w", zipfile.ZIP_DEFLATED) as workbook:
        workbook.writestr(
            "[Content_Types].xml",
            b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            b'<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            b'<Default Extension="xml" ContentType="application/xml"/>'
            b'<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
            b"</Types>",
        )
        workbook.writestr(
            "_rels/.rels",
            b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            b'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>'
            b"</Relationships>",
        )
        workbook.writestr(
            "xl/workbook.xml",
            b'<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"/>',
        )

    with zipfile.ZipFile(source_path, "r") as package:
        content_types = package.read("[Content_Types].xml").replace(
            b"</Types>",
            b'<Default Extension="xlsx" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"/>'
            b"</Types>",
        )
    safe_path = _rewrite_package(
        source_path,
        tmp_path / "chart-workbook.pptx",
        replacements={"[Content_Types].xml": content_types},
        additions={
            "ppt/charts/_rels/chart1.xml.rels": (
                b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                b'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/package" '
                b'Target="../embeddings/workbook.xlsx"/>'
                b"</Relationships>"
            ),
            "ppt/embeddings/workbook.xlsx": workbook_buffer.getvalue(),
        },
    )

    report = build_reference_inventory(
        [ReferenceSource("fixture", safe_path)],
        expected_slide_counts={"fixture": 2},
    )

    assert report["sources"][0]["embeddedWorkbookCount"] == 1


def test_security_preflight_rejects_external_relationships(tmp_path: Path) -> None:
    source_path = _source_pptx(tmp_path)

    def add_external_relationship(content: bytes) -> bytes:
        return content.replace(
            b"</Relationships>",
            b'<Relationship Id="rIdExternal" Type="http://example.test/type" '
            b'Target="https://example.test/private" TargetMode="External"/>'
            b"</Relationships>",
        )

    unsafe_path = _replace_xml(
        source_path,
        tmp_path / "external.pptx",
        "_rels/.rels",
        add_external_relationship,
    )

    with pytest.raises(InventoryValidationError, match="external_relationship"):
        build_reference_inventory(
            [ReferenceSource("fixture", unsafe_path)],
            expected_slide_counts={"fixture": 2},
        )


@pytest.mark.parametrize(
    ("part", "replacement", "error_code"),
    [
        ("[Content_Types].xml", b"<Types>", "malformed_content_types"),
        ("_rels/.rels", b"<Relationships>", "malformed_relationships"),
        ("ppt/presentation.xml", b"<presentation>", "malformed_presentation"),
    ],
)
def test_security_preflight_rejects_malformed_package_control_xml(
    tmp_path: Path,
    part: str,
    replacement: bytes,
    error_code: str,
) -> None:
    source_path = _source_pptx(tmp_path)
    unsafe_path = _rewrite_package(
        source_path,
        tmp_path / "malformed.pptx",
        replacements={part: replacement},
    )

    with pytest.raises(InventoryValidationError, match=error_code):
        build_reference_inventory(
            [ReferenceSource("fixture", unsafe_path)],
            expected_slide_counts={"fixture": 2},
        )


def test_security_preflight_rejects_broken_presentation_slide_mapping(
    tmp_path: Path,
) -> None:
    source_path = _source_pptx(tmp_path)

    def break_slide_target(content: bytes) -> bytes:
        return content.replace(b"slides/slide1.xml", b"slides/missing.xml")

    unsafe_path = _replace_xml(
        source_path,
        tmp_path / "broken-mapping.pptx",
        "ppt/_rels/presentation.xml.rels",
        break_slide_target,
    )

    with pytest.raises(InventoryValidationError, match="malformed_relationships"):
        build_reference_inventory(
            [ReferenceSource("fixture", unsafe_path)],
            expected_slide_counts={"fixture": 2},
        )


def test_security_preflight_rejects_encrypted_or_non_zip_package(
    tmp_path: Path,
) -> None:
    encrypted_path = tmp_path / "encrypted.pptx"
    encrypted_path.write_bytes(bytes.fromhex("d0cf11e0a1b11ae1") + b"encrypted")

    with pytest.raises(InventoryValidationError, match="encrypted_package"):
        build_reference_inventory(
            [ReferenceSource("fixture", encrypted_path)],
            expected_slide_counts={"fixture": 2},
        )


def test_security_preflight_rejects_password_protected_presentation(
    tmp_path: Path,
) -> None:
    source_path = _source_pptx(tmp_path)

    def add_protection(content: bytes) -> bytes:
        return content.replace(
            b"</p:presentation>",
            b'<p:modifyVerifier cryptProviderType="rsaAES"/>'
            b"</p:presentation>",
        )

    unsafe_path = _replace_xml(
        source_path,
        tmp_path / "protected.pptx",
        "ppt/presentation.xml",
        add_protection,
    )

    with pytest.raises(InventoryValidationError, match="protected_presentation"):
        build_reference_inventory(
            [ReferenceSource("fixture", unsafe_path)],
            expected_slide_counts={"fixture": 2},
        )


def test_inventory_cli_returns_non_zero_on_drift_without_writing_report(
    tmp_path: Path,
) -> None:
    source_path = _source_pptx(tmp_path)
    output_path = tmp_path / "report.json"
    script = Path(__file__).parents[1] / "scripts/build_ooxml_reference_inventory.py"

    result = subprocess.run(
        [
            sys.executable,
            str(script),
            "--source",
            f"simple-light={source_path}",
            "--output",
            str(output_path),
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode != 0
    assert "inventory_drift" in result.stderr
    assert not output_path.exists()
