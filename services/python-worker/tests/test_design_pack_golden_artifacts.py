import base64
from io import BytesIO
from pathlib import Path

from pptx import Presentation

from app.ai.deck_generation.design_pack_artifacts import build_golden_deck
from app.ai.deck_generation.design_pack_evaluation import load_golden_briefs
from app.ai.deck_generation.quality import (
    detect_text_overlap_candidates,
    validate_design,
    validate_layout,
)
from app.ai.deck_pptx_export import DeckPptxExportRequest, export_deck_pptx


FIXTURES = Path(__file__).parent / "fixtures/design-pack-golden"


def test_four_family_golden_decks_export_without_publication_issues() -> None:
    for brief in load_golden_briefs(FIXTURES):
        deck = build_golden_deck(brief)
        export = export_deck_pptx(DeckPptxExportRequest(deck=deck))
        presentation = Presentation(BytesIO(base64.b64decode(export.content_base64)))

        assert len(presentation.slides) == len(brief.slides)
        assert export.warnings == []
        assert validate_layout(deck) == []
        assert validate_design(deck) == []
        assert detect_text_overlap_candidates(deck) == []
