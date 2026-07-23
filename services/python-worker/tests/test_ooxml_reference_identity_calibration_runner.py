from __future__ import annotations

import hashlib
import json
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation

from scripts import run_ooxml_reference_identity_calibration as runner


TEMPLATE_IDS = [
    "simple-light",
    "simple-dark",
    "operating-review",
    "business-review",
    "project-kickoff",
    "team-alignment",
    "market-trends-report",
]


def _package() -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[5])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _png(color: str) -> bytes:
    output = BytesIO()
    Image.new("RGB", (160, 90), color).save(output, format="PNG")
    return output.getvalue()


def _write_plan(tmp_path: Path) -> Path:
    package = _package()
    templates: list[dict[str, object]] = []
    for template_id in TEMPLATE_IDS:
        root = tmp_path / "private-input" / template_id
        root.mkdir(parents=True)
        source = root / "source.pptx"
        source.write_bytes(package)
        manifest = root / "manifest.json"
        manifest.write_text(
            json.dumps({"templateId": template_id, "version": 1}),
            encoding="utf-8",
        )
        templates.append(
            {
                "templateId": template_id,
                "templateVersion": 1,
                "sourcePath": str(source),
                "manifestPath": str(manifest),
            }
        )
    plan = tmp_path / "private-plan.json"
    plan.write_text(
        json.dumps({"schemaVersion": 1, "templates": templates}),
        encoding="utf-8",
    )
    return plan


def _fake_manifest(source: Path, value: dict[str, object]) -> SimpleNamespace:
    return SimpleNamespace(
        template_id=value["templateId"],
        version=value["version"],
        source_sha256=hashlib.sha256(source.read_bytes()).hexdigest(),
        provenance=SimpleNamespace(authorization_status="approved"),
        canvas=SimpleNamespace(aspect_ratio="16:9"),
        source_slides=[
            SimpleNamespace(
                source_slide_id=f"{value['templateId']}-slide-001",
                source_slide_part="ppt/slides/slide1.xml",
            )
        ],
    )


def _fake_render(
    _package_bytes: bytes,
    _aspect_ratio: str,
) -> runner.LibreOfficeRenderResult:
    return runner.LibreOfficeRenderResult(
        status="passed",
        version="26.8.0",
        pngs=(_png("#112233"),),
        issue_codes=(),
    )


def _fake_font_manifest(
    _source_package: bytes,
    _identity_package: bytes,
    template_id: str,
    template_version: int,
    renderer_version: str | None,
) -> dict[str, object]:
    return {
        "schemaVersion": 1,
        "templateId": template_id,
        "templateVersion": template_version,
        "renderer": "libreoffice-pdf-pymupdf",
        "rendererVersion": renderer_version,
        "status": "pending",
        "fontCount": 1,
        "issueCodes": ["FONT_AVAILABILITY_VALIDATION_PENDING"],
        "fonts": [
            {
                "requestedFamily": "Requested Font",
                "resolvedFamily": "Fallback Font",
                "roles": ["identity-clone", "source"],
                "status": "substituted",
                "sha256": "a" * 64,
            }
        ],
    }


def test_runner_builds_evidence_derived_candidate_without_approving_runtime(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)
    monkeypatch.setattr(runner, "_identity_font_manifest", _fake_font_manifest)

    summary = runner.run_identity_calibration(
        plan,
        tmp_path / "private-output",
        render_deck=_fake_render,
    )

    assert summary["status"] == "generated"
    assert summary["approvalStatus"] == "pending"
    assert summary["runtimeEligible"] is False
    assert summary["templateCount"] == 7
    assert summary["slideCount"] == 7
    candidate_path = tmp_path / "private-output/calibration-candidate.json"
    candidate = json.loads(candidate_path.read_text(encoding="utf-8"))
    assert candidate["status"] == "pending-approval"
    assert candidate["runtimeEligible"] is False
    assert candidate["applied"] is False
    assert candidate["measuredMinimumLockedRegionSsim"] == 1.0
    assert candidate["proposedLockedRegionSsimThreshold"] is None
    assert candidate["geometryEdgeTolerancePx"] == 0
    assert len(candidate["identityBaselines"]) == 7
    assert all(
        baseline["minimumLockedRegionSsim"] == 1.0
        and baseline["structuralStatus"] == "passed"
        and baseline["fontStatus"] == "pending"
        for baseline in candidate["identityBaselines"]
    )
    assert "FONT_AVAILABILITY_VALIDATION_PENDING" in candidate["issueCodes"]
    assert "HUMAN_FIDELITY_REVIEW_PENDING" in candidate["issueCodes"]
    serialized = "\n".join(
        path.read_text(encoding="utf-8")
        for path in (tmp_path / "private-output").rglob("*.json")
    )
    assert str(tmp_path) not in serialized
    assert "source.pptx" not in serialized


def test_runner_refuses_to_propose_threshold_on_structural_drift(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    plan = _write_plan(tmp_path)
    monkeypatch.setattr(runner, "validate_source_slide_annotations", _fake_manifest)
    monkeypatch.setattr(runner, "_identity_font_manifest", _fake_font_manifest)
    monkeypatch.setattr(runner, "_locked_snapshot_drift", lambda *_args: (1, 0, 0))

    summary = runner.run_identity_calibration(
        plan,
        tmp_path / "private-output",
        render_deck=_fake_render,
    )

    assert summary["status"] == "failed"
    assert summary["runtimeEligible"] is False
    assert not (tmp_path / "private-output/calibration-candidate.json").exists()
    assert "OOXML_REFERENCE_FIDELITY_LOCKED_GEOMETRY_DRIFT" in summary["issueCodes"]


def test_runner_rejects_output_inside_repository(tmp_path: Path) -> None:
    plan = _write_plan(tmp_path)

    with pytest.raises(
        runner.IdentityCalibrationRunnerError,
        match="IDENTITY_CALIBRATION_OUTPUT_INVALID",
    ):
        runner.run_identity_calibration(
            plan,
            runner.REPOSITORY_ROOT / "private-output",
            render_deck=_fake_render,
        )
