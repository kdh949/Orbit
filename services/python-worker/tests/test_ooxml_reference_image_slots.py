from __future__ import annotations

import base64
import zipfile
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.ai.ooxml_reference_templates.capacity import SlotCapacityError
from app.ai.ooxml_reference_templates.clone import clone_source_slides
from app.ai.ooxml_reference_templates.image_slots import replace_image_slot
from app.ai.ooxml_reference_templates.models import OoxmlImageTemplateSlot
from app.ai.ooxml_reference_templates.package import validate_cloned_package
from app.ai.pptx_ooxml_generation import (
    PptxOoxmlSyncResult,
    generate_pptx_ooxml,
    sync_pptx_ooxml,
)


DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


def _png(size: tuple[int, int], color: tuple[int, int, int, int]) -> bytes:
    output = BytesIO()
    Image.new("RGBA", size, color).save(output, format="PNG")
    return output.getvalue()


def _jpeg(size: tuple[int, int], color: tuple[int, int, int]) -> bytes:
    output = BytesIO()
    Image.new("RGB", size, color).save(output, format="JPEG")
    return output.getvalue()


def _image_fixture(tmp_path: Path) -> tuple[bytes, str, str]:
    image_path = tmp_path / "original.png"
    image_path.write_bytes(_png((400, 300), (30, 90, 150, 180)))
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path), Inches(2), Inches(1), Inches(6), Inches(4)
    )
    picture.crop_left = 0.1
    picture.crop_top = 0.05
    picture.crop_right = 0.15
    picture.crop_bottom = 0.08
    picture.rotation = 17
    relationship_id = picture._pic.blipFill.blip.rEmbed
    output = BytesIO()
    presentation.save(output)
    enriched = _add_mask_opacity_and_effect(output.getvalue(), str(picture.shape_id))
    cloned = clone_source_slides(
        enriched, source_slide_parts=["ppt/slides/slide1.xml"]
    )
    return cloned.package_bytes, str(picture.shape_id), relationship_id


def _shared_image_fixture(tmp_path: Path) -> tuple[bytes, str, str]:
    image_path = tmp_path / "shared-original.png"
    image_path.write_bytes(_png((400, 300), (30, 90, 150, 180)))
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1), Inches(4), Inches(3)
    )
    slide.shapes.add_picture(
        str(image_path), Inches(6), Inches(1), Inches(4), Inches(3)
    )
    relationship_id = picture._pic.blipFill.blip.rEmbed
    output = BytesIO()
    presentation.save(output)
    enriched = _add_mask_opacity_and_effect(output.getvalue(), str(picture.shape_id))
    cloned = clone_source_slides(
        enriched, source_slide_parts=["ppt/slides/slide1.xml"]
    )
    return cloned.package_bytes, str(picture.shape_id), relationship_id


def _cross_slide_shared_image_fixture(tmp_path: Path) -> tuple[bytes, str, str]:
    image_path = tmp_path / "cross-slide-original.png"
    image_path.write_bytes(_png((400, 300), (30, 90, 150, 180)))
    presentation = Presentation()
    first_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = first_slide.shapes.add_picture(
        str(image_path), Inches(2), Inches(1), Inches(6), Inches(4)
    )
    second_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_slide.shapes.add_picture(
        str(image_path), Inches(2), Inches(1), Inches(6), Inches(4)
    )
    relationship_id = picture._pic.blipFill.blip.rEmbed
    output = BytesIO()
    presentation.save(output)
    enriched = _add_mask_opacity_and_effect(output.getvalue(), str(picture.shape_id))
    cloned = clone_source_slides(
        enriched,
        source_slide_parts=[
            "ppt/slides/slide1.xml",
            "ppt/slides/slide2.xml",
        ],
    )
    return cloned.package_bytes, str(picture.shape_id), relationship_id


def _add_mask_opacity_and_effect(package_bytes: bytes, shape_id: str) -> bytes:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}
    root = ET.fromstring(entries["ppt/slides/slide1.xml"])
    picture = next(
        item
        for item in root.findall(f".//{{{PML_NS}}}pic")
        if item.find(f".//{{{PML_NS}}}cNvPr").get("id") == shape_id
    )
    blip = picture.find(f".//{{{DML_NS}}}blip")
    ET.SubElement(blip, f"{{{DML_NS}}}alphaModFix", {"amt": "72000"})
    geometry = picture.find(f".//{{{DML_NS}}}prstGeom")
    geometry.set("prst", "ellipse")
    shape_properties = picture.find(f"{{{PML_NS}}}spPr")
    effects = ET.SubElement(shape_properties, f"{{{DML_NS}}}effectLst")
    ET.SubElement(effects, f"{{{DML_NS}}}outerShdw", {"blurRad": "38100"})
    entries["ppt/slides/slide1.xml"] = ET.tostring(root, xml_declaration=True)
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as package:
        for name, content in entries.items():
            package.writestr(name, content)
    return output.getvalue()


def _slot(shape_id: str, relationship_id: str, *, maximum: float = 1.5) -> OoxmlImageTemplateSlot:
    return OoxmlImageTemplateSlot.model_validate(
        {
            "slotId": "fixture-v1-slide-01-image",
            "semanticRole": "image",
            "contentType": "image",
            "required": True,
            "locator": {
                "slidePart": "ppt/slides/slide1.xml",
                "shapeId": shape_id,
                "placeholderType": None,
                "relationshipId": relationship_id,
            },
            "capacity": {
                "minAspectRatio": 1.2,
                "maxAspectRatio": maximum,
                "cropPolicy": "preserve-frame",
                "alphaRequired": True,
                "maskRequired": True,
            },
            "mutationPolicy": ["image-source"],
            "replacementPolicy": {"overflow": "fail"},
        }
    )


def _picture_state(package_bytes: bytes, shape_id: str, relationship_id: str) -> dict:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        slide = ET.fromstring(package.read("ppt/slides/slide1.xml"))
        relationships = ET.fromstring(package.read("ppt/slides/_rels/slide1.xml.rels"))
        picture = next(
            item
            for item in slide.findall(f".//{{{PML_NS}}}pic")
            if item.find(f".//{{{PML_NS}}}cNvPr").get("id") == shape_id
        )
        relation = next(
            item
            for item in relationships.findall(f"{{{PKG_REL_NS}}}Relationship")
            if item.get("Id") == relationship_id
        )
        media_part = "ppt/media/" + Path(relation.get("Target")).name
        return {
            "xfrm": ET.tostring(picture.find(f".//{{{DML_NS}}}xfrm")),
            "crop": ET.tostring(picture.find(f".//{{{DML_NS}}}srcRect")),
            "mask": ET.tostring(picture.find(f".//{{{DML_NS}}}prstGeom")),
            "alpha": ET.tostring(picture.find(f".//{{{DML_NS}}}alphaModFix")),
            "effect": ET.tostring(picture.find(f".//{{{DML_NS}}}effectLst")),
            "relationshipId": picture.find(f".//{{{DML_NS}}}blip").get(
                f"{{{REL_NS}}}embed"
            ),
            "relationshipTarget": relation.get("Target"),
            "media": package.read(media_part),
        }


def _image_target_usage_count(package_bytes: bytes, relationship_id: str) -> tuple[int, int]:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        slide = ET.fromstring(package.read("ppt/slides/slide1.xml"))
        relationships = ET.fromstring(package.read("ppt/slides/_rels/slide1.xml.rels"))
        relation = next(
            item
            for item in relationships.findall(f"{{{PKG_REL_NS}}}Relationship")
            if item.get("Id") == relationship_id
        )
        target = relation.get("Target")
        target_name = Path(target).name
        relationship_count = 0
        for name in package.namelist():
            if not name.endswith(".rels"):
                continue
            root = ET.fromstring(package.read(name))
            relationship_count += sum(
                item.get("Type", "").endswith("/image")
                and Path(item.get("Target", "")).name == target_name
                for item in root.findall(f"{{{PKG_REL_NS}}}Relationship")
            )
        embed_count = sum(
            blip.get(f"{{{REL_NS}}}embed") == relationship_id
            for blip in slide.findall(f".//{{{DML_NS}}}blip")
        )
        return relationship_count, embed_count


def _rewrite_package_entry(
    package_bytes: bytes,
    entry_name: str,
    content: bytes,
) -> bytes:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        entries = {
            item.filename: package.read(item.filename)
            for item in package.infolist()
            if not item.is_dir()
        }
    entries[entry_name] = content
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as package:
        for name, value in entries.items():
            package.writestr(name, value)
    return output.getvalue()


def _set_media_content_type_override(
    package_bytes: bytes,
    relationship_id: str,
    content_type: str,
) -> bytes:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        relationships = ET.fromstring(
            package.read("ppt/slides/_rels/slide1.xml.rels")
        )
        relation = next(
            item
            for item in relationships.findall(
                f"{{{PKG_REL_NS}}}Relationship"
            )
            if item.get("Id") == relationship_id
        )
        media_part = "ppt/media/" + Path(relation.get("Target", "")).name
        content_types = ET.fromstring(package.read("[Content_Types].xml"))
    for child in list(content_types):
        if (
            child.tag == f"{{{CONTENT_TYPES_NS}}}Override"
            and child.get("PartName") == f"/{media_part}"
        ):
            content_types.remove(child)
    ET.SubElement(
        content_types,
        f"{{{CONTENT_TYPES_NS}}}Override",
        {"PartName": f"/{media_part}", "ContentType": content_type},
    )
    return _rewrite_package_entry(
        package_bytes,
        "[Content_Types].xml",
        ET.tostring(content_types, xml_declaration=True),
    )


def _add_external_image_relationship(
    package_bytes: bytes,
    relationship_id: str,
) -> bytes:
    rels_part = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        relationships = ET.fromstring(package.read(rels_part))
    relation = next(
        item
        for item in relationships.findall(f"{{{PKG_REL_NS}}}Relationship")
        if item.get("Id") == relationship_id
    )
    ET.SubElement(
        relationships,
        f"{{{PKG_REL_NS}}}Relationship",
        {
            "Id": "rId999",
            "Type": relation.get("Type", ""),
            "Target": relation.get("Target", ""),
            "TargetMode": "External",
        },
    )
    return _rewrite_package_entry(
        package_bytes,
        rels_part,
        ET.tostring(relationships, xml_declaration=True),
    )


def _picture_alt(package_bytes: bytes, shape_id: str) -> str | None:
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        slide = ET.fromstring(package.read("ppt/slides/slide1.xml"))
    picture = next(
        item
        for item in slide.findall(f".//{{{PML_NS}}}pic")
        if item.find(f".//{{{PML_NS}}}cNvPr").get("id") == shape_id
    )
    return picture.find(f".//{{{PML_NS}}}cNvPr").get("descr")


def _set_picture_geometry(
    package_bytes: bytes,
    shape_id: str,
    preset: str,
) -> bytes:
    slide_part = "ppt/slides/slide1.xml"
    with zipfile.ZipFile(BytesIO(package_bytes), "r") as package:
        slide = ET.fromstring(package.read(slide_part))
    picture = next(
        item
        for item in slide.findall(f".//{{{PML_NS}}}pic")
        if item.find(f".//{{{PML_NS}}}cNvPr").get("id") == shape_id
    )
    picture.find(f".//{{{DML_NS}}}prstGeom").set("prst", preset)
    return _rewrite_package_entry(
        package_bytes,
        slide_part,
        ET.tostring(slide, xml_declaration=True),
    )


def _sync_reference_image_slot(
    tmp_path: Path,
    package_bytes: bytes,
    shape_id: str,
    relationship_id: str,
    replacement: bytes | None,
    *,
    mime_type: str = "image/png",
    alt: str | None = None,
    maximum_aspect_ratio: float = 1.5,
    alpha_required: bool = True,
    include_capacity: bool = True,
    include_policy: bool = True,
    omit_capacity_field: str | None = None,
    extra_props: dict[str, object] | None = None,
) -> PptxOoxmlSyncResult:
    package_path = tmp_path / "reference-image-sync.pptx"
    package_path.write_bytes(package_bytes)
    imported = generate_pptx_ooxml(package_path, "file_baseline", render=False)
    template_blueprint = imported.template_blueprint
    source = next(
        item
        for slide in template_blueprint["slides"]
        for item in slide["elementSources"]
        if item.get("shapeId") == shape_id
        and item.get("relationshipId") == relationship_id
    )
    source["writable"] = True
    element_id = source["elementId"]
    template_blueprint["referenceTemplateSnapshot"] = {
        "catalogTemplateId": "fixture",
        "catalogTemplateVersion": 2,
        "sourceSha256": "a" * 64,
        "sourceSlideIds": ["fixture-slide-001"],
        "slotAssignmentCount": 1,
    }
    policy = {
        "slotId": "fixture-v2-slide-001-shape-image",
        "elementId": element_id,
        "mutationPolicy": ["image-source"],
        "frameLocked": True,
    }
    if include_capacity:
        image_capacity: dict[str, object] = {
            "minAspectRatio": 1.2,
            "maxAspectRatio": maximum_aspect_ratio,
            "cropPolicy": "preserve-frame",
            "alphaRequired": alpha_required,
            "maskRequired": True,
        }
        if omit_capacity_field is not None:
            image_capacity.pop(omit_capacity_field)
        policy["imageCapacity"] = image_capacity
    template_blueprint["slotEditPolicies"] = [policy] if include_policy else []
    template_slide = next(
        slide
        for slide in template_blueprint["slides"]
        if any(
            item["elementId"] == element_id for item in slide["elementSources"]
        )
    )
    props: dict[str, object] = {}
    if replacement is not None:
        props["src"] = (
            f"data:{mime_type};base64,"
            + base64.b64encode(replacement).decode("ascii")
        )
    if alt is not None:
        props["alt"] = alt
    if extra_props is not None:
        props.update(extra_props)
    result = sync_pptx_ooxml(
        package_path,
        template_blueprint=template_blueprint,
        operations=[
            {
                "type": "update_element_props",
                "slideId": template_slide["slideId"],
                "elementId": element_id,
                "props": props,
            }
        ],
        deck_canvas=imported.canvas,
        synced_deck_version=2,
        render=False,
    )
    return result


def _current_package_bytes(result: PptxOoxmlSyncResult) -> bytes:
    current = next(
        asset for asset in result.assets if asset.asset_id == "current_package"
    )
    return base64.b64decode(current.content_base64)


def test_image_replacement_preserves_frame_crop_mask_rotation_opacity_and_effect(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    assert _image_target_usage_count(source, relationship_id) == (1, 1)
    before = _picture_state(source, shape_id, relationship_id)
    replacement = _png((800, 600), (180, 40, 80, 160))

    result = replace_image_slot(
        source,
        slot=_slot(shape_id, relationship_id),
        image_bytes=replacement,
        mime_type="image/png",
    )
    after = _picture_state(result.package_bytes, shape_id, relationship_id)

    for key in ("xfrm", "crop", "mask", "alpha", "effect"):
        assert after[key] == before[key]
    assert after["relationshipId"] == relationship_id
    assert after["relationshipTarget"] == before["relationshipTarget"]
    assert after["media"] == replacement
    assert result.warning_codes == []
    assert validate_cloned_package(result.package_bytes) == []


def test_reference_editor_sync_preserves_exclusive_image_relationship(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    before = _picture_state(source, shape_id, relationship_id)
    replacement = _png((800, 600), (180, 40, 80, 160))

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        replacement,
    )
    assert result.warnings == []
    assert result.unsupported_operations == []
    assert len(result.applied_operations) == 1
    current = _current_package_bytes(result)
    after = _picture_state(current, shape_id, relationship_id)

    assert after["relationshipId"] == relationship_id
    assert after["relationshipTarget"] == before["relationshipTarget"]
    assert after["media"] == replacement
    assert _image_target_usage_count(current, relationship_id) == (1, 1)
    assert validate_cloned_package(current) == []


def test_reference_editor_sync_updates_alt_without_replacing_media(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    before = _picture_state(source, shape_id, relationship_id)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        None,
        alt="Approved replacement image",
    )

    assert result.warnings == []
    assert result.unsupported_operations == []
    assert len(result.applied_operations) == 1
    current = _current_package_bytes(result)
    after = _picture_state(current, shape_id, relationship_id)
    assert after["media"] == before["media"]
    assert _picture_alt(current, shape_id) == "Approved replacement image"


def test_reference_editor_sync_rejects_crop_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        None,
        extra_props={"crop": {"left": 0.1}},
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_FIELDS_UNSUPPORTED"
    ]
    assert result.applied_operations == []
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_non_xml_alt_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        None,
        alt="bad\u0000alt",
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_FIELDS_UNSUPPORTED"
    ]
    assert result.applied_operations == []
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_shared_image_media_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _shared_image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any("reference image media is shared" in item for item in result.warnings)
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_cross_slide_shared_media_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _cross_slide_shared_image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any("reference image media is shared" in item for item in result.warnings)
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_ignores_external_image_relationship_in_usage_count(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    source = _add_external_image_relationship(source, relationship_id)
    replacement = _png((800, 600), (180, 40, 80, 160))

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        replacement,
    )

    assert result.warnings == []
    assert result.unsupported_operations == []
    assert _picture_state(
        _current_package_bytes(result),
        shape_id,
        relationship_id,
    )["media"] == replacement


def test_reference_editor_sync_rejects_image_format_drift_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _jpeg((800, 600), (180, 40, 80)),
        mime_type="image/jpeg",
        alpha_required=False,
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any(
        "reference image format mismatch" in item for item in result.warnings
    )
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_uses_effective_content_type_override(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    source = _set_media_content_type_override(
        source,
        relationship_id,
        "image/jpeg",
    )

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any(
        "reference image format mismatch" in item for item in result.warnings
    )
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_aspect_ratio_capacity_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((1200, 300), (180, 40, 80, 160)),
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any("capacity aspect ratio exceeded" in item for item in result.warnings)
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_required_alpha_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _jpeg((800, 600), (180, 40, 80)),
        mime_type="image/jpeg",
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any("capacity requires alpha" in item for item in result.warnings)
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_missing_mask_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    source = _set_picture_geometry(source, shape_id, "rect")

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_UPDATE_FAILED"
    ]
    assert result.applied_operations == []
    assert any("capacity requires mask" in item for item in result.warnings)
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_missing_capacity_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
        include_capacity=False,
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_FIELDS_UNSUPPORTED"
    ]
    assert result.applied_operations == []
    assert _current_package_bytes(result) == source


def test_reference_editor_sync_rejects_missing_policy_without_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
        include_policy=False,
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_FIELDS_UNSUPPORTED"
    ]
    assert result.applied_operations == []
    assert _current_package_bytes(result) == source


@pytest.mark.parametrize("missing_field", ["alphaRequired", "maskRequired"])
def test_reference_editor_sync_rejects_incomplete_capacity_without_mutation(
    tmp_path: Path,
    missing_field: str,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    result = _sync_reference_image_slot(
        tmp_path,
        source,
        shape_id,
        relationship_id,
        _png((800, 600), (180, 40, 80, 160)),
        omit_capacity_field=missing_field,
    )

    assert [item.reason_code for item in result.unsupported_operations] == [
        "PROPS_FIELDS_UNSUPPORTED"
    ]
    assert result.applied_operations == []
    assert _current_package_bytes(result) == source


def test_image_replacement_rejects_shared_media_target_without_package_mutation(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _shared_image_fixture(tmp_path)
    relationship_count, embed_count = _image_target_usage_count(
        source, relationship_id
    )
    assert (relationship_count, embed_count) == (1, 2)

    with pytest.raises(SlotCapacityError) as caught:
        replace_image_slot(
            source,
            slot=_slot(shape_id, relationship_id),
            image_bytes=_png((800, 600), (180, 40, 80, 160)),
            mime_type="image/png",
        )

    assert caught.value.code == "OOXML_REFERENCE_IMAGE_MEDIA_SHARED"
    assert caught.value.retryable is False
    assert caught.value.package_bytes == source
    assert caught.value.authored_fallback_created is False


def test_image_replacement_rejects_media_target_reused_by_another_slide(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _cross_slide_shared_image_fixture(tmp_path)
    assert _image_target_usage_count(source, relationship_id) == (2, 1)

    with pytest.raises(SlotCapacityError) as caught:
        replace_image_slot(
            source,
            slot=_slot(shape_id, relationship_id),
            image_bytes=_png((800, 600), (180, 40, 80, 160)),
            mime_type="image/png",
        )

    assert caught.value.code == "OOXML_REFERENCE_IMAGE_MEDIA_SHARED"
    assert caught.value.package_bytes == source


def test_image_aspect_ratio_capacity_fails_without_package_mutation(tmp_path: Path) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)

    with pytest.raises(SlotCapacityError) as caught:
        replace_image_slot(
            source,
            slot=_slot(shape_id, relationship_id, maximum=1.4),
            image_bytes=_png((1200, 300), (10, 20, 30, 255)),
            mime_type="image/png",
        )

    assert caught.value.code == "OOXML_REFERENCE_CAPACITY_IMAGE_ASPECT_RATIO"
    assert caught.value.retryable is False
    assert caught.value.package_bytes == source


def test_image_replacement_uses_effective_content_type_override(
    tmp_path: Path,
) -> None:
    source, shape_id, relationship_id = _image_fixture(tmp_path)
    source = _set_media_content_type_override(
        source,
        relationship_id,
        "image/jpeg",
    )

    with pytest.raises(SlotCapacityError) as caught:
        replace_image_slot(
            source,
            slot=_slot(shape_id, relationship_id),
            image_bytes=_png((800, 600), (180, 40, 80, 160)),
            mime_type="image/png",
        )

    assert caught.value.code == "OOXML_REFERENCE_IMAGE_FORMAT_MISMATCH"
    assert caught.value.package_bytes == source
