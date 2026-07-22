from __future__ import annotations

import hashlib
import posixpath
import zipfile
from io import BytesIO
from pathlib import Path
from pathlib import PurePosixPath
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from app.ai.ooxml_reference_templates.materialize import (
    MaterializationError,
    materialize_reference_package,
)
from app.ai.ooxml_reference_templates.models import (
    OoxmlReferenceTemplateManifest,
    OoxmlTemplateSnapshot,
)


def _sha256(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _source_package(tmp_path: Path) -> Path:
    path = tmp_path / "generated.pptx"
    presentation = Presentation()
    cover = presentation.slides.add_slide(presentation.slide_layouts[5])
    cover.shapes.title.text = "Generated title"
    decoration = cover.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8),
        Inches(6),
        Inches(1),
        Inches(0.25),
    )
    decoration.name = "Locked decoration"
    closing = presentation.slides.add_slide(presentation.slide_layouts[5])
    closing.shapes.title.text = "Closing"
    presentation.save(path)
    return path


def _manifest(path: Path, *, title_shape_id: str = "2") -> OoxmlReferenceTemplateManifest:
    preview_sha = "c" * 64
    return OoxmlReferenceTemplateManifest.model_validate(
        {
            "templateId": "operating-review",
            "version": 1,
            "status": "active",
            "sourceFormat": "pptx",
            "sourceSha256": _sha256(path.read_bytes()),
            "slideCount": 2,
            "canvas": {
                "aspectRatio": "16:9",
                "widthEmu": 12_192_000,
                "heightEmu": 6_858_000,
            },
            "name": "Operating Review",
            "description": "경영 보고와 KPI 중심",
            "preview": {
                "coverPreviewId": "cover",
                "coverPreviewSha256": preview_sha,
                "bodyPreviewId": "body",
                "bodyPreviewSha256": preview_sha,
            },
            "sourceSlides": [
                {
                    "sourceSlideId": "cover-01",
                    "sourceSlidePart": "ppt/slides/slide1.xml",
                    "sourceOrder": 1,
                    "semanticRole": "cover",
                    "relationships": {
                        "layoutPart": "ppt/slideLayouts/slideLayout1.xml",
                        "masterPart": "ppt/slideMasters/slideMaster1.xml",
                        "themePart": "ppt/theme/theme1.xml",
                    },
                    "capacity": {
                        "textSlotCount": 1,
                        "imageSlotCount": 0,
                        "tableSlotCount": 0,
                        "chartSlotCount": 0,
                    },
                    "previewId": "cover",
                    "lockedInventorySha256": "a" * 64,
                    "slots": [
                        {
                            "slotId": "operating-review-v1-slide-01-title",
                            "semanticRole": "title",
                            "contentType": "text",
                            "required": True,
                            "locator": {
                                "slidePart": "ppt/slides/slide1.xml",
                                "shapeId": title_shape_id,
                                "placeholderType": "title",
                                "relationshipId": None,
                            },
                            "capacity": {"maxChars": 80, "maxLines": 2},
                            "mutationPolicy": ["text-content"],
                            "replacementPolicy": {"overflow": "fail"},
                        }
                    ],
                },
                {
                    "sourceSlideId": "closing-02",
                    "sourceSlidePart": "ppt/slides/slide2.xml",
                    "sourceOrder": 2,
                    "semanticRole": "closing",
                    "relationships": {
                        "layoutPart": "ppt/slideLayouts/slideLayout1.xml",
                        "masterPart": "ppt/slideMasters/slideMaster1.xml",
                        "themePart": "ppt/theme/theme1.xml",
                    },
                    "capacity": {
                        "textSlotCount": 0,
                        "imageSlotCount": 0,
                        "tableSlotCount": 0,
                        "chartSlotCount": 0,
                    },
                    "previewId": "body",
                    "lockedInventorySha256": "b" * 64,
                    "slots": [],
                },
            ],
            "provenance": {
                "authorizationStatus": "approved",
                "inventoryVersion": 1,
            },
        }
    )


def _snapshot() -> OoxmlTemplateSnapshot:
    return OoxmlTemplateSnapshot.model_validate(
        {
            "catalogTemplateId": "operating-review",
            "catalogTemplateVersion": 1,
            "sourceSha256": "d" * 64,
            "sourceSlideIds": ["cover-01", "closing-02"],
            "slotAssignmentCount": 1,
        }
    )


def test_materialization_reconciles_each_slot_to_one_imported_element_and_locks_decoration(
    tmp_path: Path,
) -> None:
    package_path = _source_package(tmp_path)
    snapshot = _snapshot().model_copy(
        update={"source_sha256": _sha256(package_path.read_bytes())}
    )

    materialized = materialize_reference_package(
        package_path,
        baseline_file_id="file_reference_baseline",
        current_file_id="file_reference_current",
        manifest=_manifest(package_path),
        snapshot=snapshot,
        render=False,
    )

    blueprint = materialized.template_blueprint
    assert blueprint["sourcePackageFileId"] == "file_reference_baseline"
    assert blueprint["currentPackageFileId"] == "file_reference_current"
    assert blueprint["referenceTemplateSnapshot"] == snapshot.model_dump(
        by_alias=True
    )
    assert blueprint["slotEditPolicies"] == [
        {
            "slotId": "operating-review-v1-slide-01-title",
            "elementId": blueprint["slotEditPolicies"][0]["elementId"],
            "mutationPolicy": ["text-content"],
            "frameLocked": True,
        }
    ]

    slot_element_id = blueprint["slotEditPolicies"][0]["elementId"]
    matching_sources = [
        source
        for slide in blueprint["slides"]
        for source in slide["elementSources"]
        if source["elementId"] == slot_element_id
    ]
    assert len(matching_sources) == 1
    assert matching_sources[0] | {
        "slidePart": "ppt/slides/slide1.xml",
        "shapeId": "2",
        "writable": True,
    } == matching_sources[0]

    elements = [
        element
        for slide in materialized.blueprint["slides"]
        for element in slide["elements"]
    ]
    assert next(
        element for element in elements if element["elementId"] == slot_element_id
    )["locked"] is False
    assert all(
        element["locked"] is True
        for element in elements
        if element["elementId"] != slot_element_id
    )


def test_materialization_fails_closed_when_slot_locator_has_no_unique_element(
    tmp_path: Path,
) -> None:
    package_path = _source_package(tmp_path)
    snapshot = _snapshot().model_copy(
        update={"source_sha256": _sha256(package_path.read_bytes())}
    )

    with pytest.raises(MaterializationError, match="SLOT_LOCATOR_UNRESOLVED"):
        materialize_reference_package(
            package_path,
            baseline_file_id="file_reference_baseline",
            current_file_id="file_reference_current",
            manifest=_manifest(package_path, title_shape_id="9999"),
            snapshot=snapshot,
            render=False,
        )


def test_materialization_issues_chart_capability_only_after_locator_proof(
    tmp_path: Path,
) -> None:
    package_path = tmp_path / "chart-reference.pptx"
    presentation = Presentation()
    chart_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    data = ChartData()
    data.categories = ["North", "South"]
    data.add_series("Revenue", (10, 20))
    chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(4),
        data,
    )
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(package_path)

    with zipfile.ZipFile(BytesIO(package_path.read_bytes()), "r") as package:
        slide_part = "ppt/slides/slide1.xml"
        slide = ET.fromstring(package.read(slide_part))
        frame = next(node for node in slide.iter() if node.tag.endswith("graphicFrame"))
        shape = next(node for node in frame.iter() if node.tag.endswith("cNvPr"))
        chart_ref = next(node for node in frame.iter() if node.tag.endswith("chart"))
        relationship_id = next(
            value for key, value in chart_ref.attrib.items() if key.endswith("}id")
        )
        path = PurePosixPath(slide_part)
        rels_part = str(path.parent / "_rels" / f"{path.name}.rels")
        relationships = ET.fromstring(package.read(rels_part))
        chart_relationship = next(
            node
            for node in relationships
            if node.attrib.get("Id") == relationship_id
        )
        chart_part = posixpath.normpath(
            posixpath.join(
                posixpath.dirname(slide_part),
                chart_relationship.attrib["Target"],
            )
        )
        chart_path = PurePosixPath(chart_part)
        chart_rels_part = str(
            chart_path.parent / "_rels" / f"{chart_path.name}.rels"
        )
        chart_relationships = ET.fromstring(package.read(chart_rels_part))
        workbook_relationship = next(
            node
            for node in chart_relationships
            if node.attrib.get("Type", "").endswith("/package")
        )
        workbook_part = posixpath.normpath(
            posixpath.join(
                posixpath.dirname(chart_part),
                workbook_relationship.attrib["Target"],
            )
        )
        workbook_fingerprint = _sha256(package.read(workbook_part))

    manifest_payload = _manifest(package_path).model_dump(by_alias=True)
    manifest_payload["sourceSlides"][0]["capacity"] = {
        "textSlotCount": 0,
        "imageSlotCount": 0,
        "tableSlotCount": 0,
        "chartSlotCount": 1,
    }
    manifest_payload["sourceSlides"][0]["slots"] = [
        {
            "slotId": "operating-review-v1-slide-01-chart",
            "semanticRole": "chart",
            "contentType": "chart",
            "required": True,
            "locator": {
                "slidePart": slide_part,
                "shapeId": shape.attrib["id"],
                "placeholderType": None,
                "relationshipId": relationship_id,
            },
            "capacity": {
                "chartType": "column",
                "maxCategories": 2,
                "maxSeries": 1,
                "workbookUpdatePolicy": "atomic",
                "workbookFingerprint": workbook_fingerprint,
            },
            "mutationPolicy": ["chart-data"],
            "replacementPolicy": {"overflow": "fail"},
        }
    ]
    manifest = OoxmlReferenceTemplateManifest.model_validate(manifest_payload)
    snapshot = _snapshot().model_copy(
        update={"source_sha256": _sha256(package_path.read_bytes())}
    )

    materialized = materialize_reference_package(
        package_path,
        baseline_file_id="file_reference_baseline",
        current_file_id="file_reference_current",
        manifest=manifest,
        snapshot=snapshot,
        render=False,
    )

    source = next(
        source
        for slide in materialized.template_blueprint["slides"]
        for source in slide["elementSources"]
        if source.get("elementType") == "chart"
    )
    assert source["sourceType"] == "chart"
    assert source["relationshipId"] == relationship_id
    assert source["ooxmlEditCapabilities"]["chartData"] is True
    assert source["ooxmlEditCapabilities"]["frame"] is False
    assert source["chartDataLocator"]["workbookFingerprint"] == (
        workbook_fingerprint
    )
