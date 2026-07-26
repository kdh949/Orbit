from __future__ import annotations

from functools import lru_cache

import copy

import difflib


import posixpath


import zipfile


from io import BytesIO


from typing import Any

from xml.etree import ElementTree as ET


from pptx import Presentation


from app.ai.pptx_ooxml.common import (
    CONTENT_TYPES_NS,
    DML_NS,
    NOTES_MASTER_CONTENT_TYPE,
    NOTES_MASTER_REL_TYPE,
    NOTES_SLIDE_CONTENT_TYPE,
    NOTES_SLIDE_REL_TYPE,
    PKG_REL_NS,
    PptxOoxmlUnsupportedReasonCode,
    SLIDE_REL_TYPE,
    THEME_CONTENT_TYPE,
    THEME_REL_TYPE,
    XML_SPACE,
)


def update_speaker_notes_body(
    operation: dict[str, Any],
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    template_blueprint: dict[str, Any],
    source_package: zipfile.ZipFile,
) -> PptxOoxmlUnsupportedReasonCode | None:
    from app.ai.pptx_ooxml.rendering import (
        xml_bytes,
    )
    from app.ai.pptx_ooxml.routing import (
        direct_local_children,
        first_local_child,
        is_safe_notes_part,
    )

    slide_id = str(operation.get("slideId", ""))
    source_slide_part = str(operation.get("sourceSlidePart", ""))
    matching_slides = [
        slide
        for slide in template_blueprint.get("slides", [])
        if isinstance(slide, dict)
        and slide.get("slideId") == slide_id
        and (
            not source_slide_part
            or str(slide.get("sourceSlidePart", "")) == source_slide_part
        )
    ]
    if len(matching_slides) != 1:
        return "NOTES_BODY_LOCATOR_UNSAFE"
    notes_page = matching_slides[0].get("notesPage")
    if not isinstance(notes_page, dict):
        return "NOTES_BODY_LOCATOR_UNSAFE"
    speaker_notes = operation.get("speakerNotes")
    if not isinstance(speaker_notes, str):
        return "NOTES_BODY_UPDATE_FAILED"
    if notes_page.get("status") == "absent":
        if not speaker_notes:
            return None
        return create_speaker_notes_page(
            matching_slides[0],
            notes_page,
            speaker_notes,
            package_entries,
            added_entries,
            template_blueprint,
            source_package,
        )
    if notes_page.get("bodyWritable") is not True:
        return "NOTES_BODY_NOT_WRITABLE"

    notes_part = str(notes_page.get("sourceNotesPart", ""))
    body_shape_id = str(notes_page.get("bodyShapeId", ""))
    if not is_safe_notes_part(notes_part) or not body_shape_id:
        return "NOTES_BODY_LOCATOR_UNSAFE"
    notes_xml = package_entries.get(notes_part)
    if notes_xml is None:
        return "NOTES_PART_MISSING"
    try:
        root = ET.fromstring(notes_xml)
    except ET.ParseError:
        return "NOTES_BODY_UPDATE_FAILED"
    common_slide = first_local_child(root, "cSld")
    shape_tree = (
        first_local_child(common_slide, "spTree") if common_slide is not None else None
    )
    if shape_tree is None:
        return "NOTES_BODY_LOCATOR_UNSAFE"
    body_shapes = [
        shape
        for shape in direct_local_children(shape_tree, "sp")
        if notes_placeholder_type(shape) == "body"
    ]
    matching_body_shapes = [
        shape for shape in body_shapes if notes_shape_id(shape) == body_shape_id
    ]
    if len(body_shapes) != 1 or len(matching_body_shapes) != 1:
        return "NOTES_BODY_LOCATOR_UNSAFE"
    text_body = first_local_child(matching_body_shapes[0], "txBody")
    if text_body is None:
        return "NOTES_BODY_LOCATOR_UNSAFE"
    if notes_text_body_text(text_body) == speaker_notes:
        return None
    replace_notes_text_body(text_body, speaker_notes)
    package_entries[notes_part] = xml_bytes(root)
    return None


def create_speaker_notes_page(
    slide: dict[str, Any],
    notes_page: dict[str, Any],
    speaker_notes: str,
    package_entries: dict[str, bytes],
    added_entries: dict[str, bytes],
    template_blueprint: dict[str, Any],
    source_package: zipfile.ZipFile,
) -> PptxOoxmlUnsupportedReasonCode | None:
    from app.ai.pptx_ooxml.operations import (
        resolve_relationship_part,
    )
    from app.ai.pptx_ooxml.rendering import (
        int_value,
        xml_bytes,
    )
    from app.ai.pptx_ooxml.routing import (
        first_local_child,
        is_safe_notes_master_part,
        is_safe_slide_part,
        is_safe_theme_part,
        rels_part_for_slide_part,
    )
    from app.ai.pptx_ooxml.shapes import (
        next_relationship_id,
    )

    slide_part = str(slide.get("sourceSlidePart", ""))
    slide_rels_part = rels_part_for_slide_part(slide_part)
    slide_rels_xml = package_entries.get(slide_rels_part)
    content_types_xml = package_entries.get("[Content_Types].xml")
    presentation_rels_xml = package_entries.get("ppt/_rels/presentation.xml.rels")
    if (
        not is_safe_slide_part(slide_part)
        or slide_part not in source_package.namelist()
        or slide_rels_xml is None
        or content_types_xml is None
        or presentation_rels_xml is None
    ):
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    notes_width_emu = int_value(notes_page.get("notesWidthEmu"), 0)
    notes_height_emu = int_value(notes_page.get("notesHeightEmu"), 0)
    if notes_width_emu <= 0 or notes_height_emu <= 0:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"

    try:
        slide_rels_root = ET.fromstring(slide_rels_xml)
        content_types_root = ET.fromstring(content_types_xml)
        presentation_rels_root = ET.fromstring(presentation_rels_xml)
    except ET.ParseError:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    if relationship_type_nodes(slide_rels_root, NOTES_SLIDE_REL_TYPE):
        return "NOTES_MASTER_CAPABILITY_UNSAFE"

    template = minimal_notes_package_template()
    if template is None:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    source_names = (
        set(source_package.namelist()) | set(package_entries) | set(added_entries)
    )
    actual_master_parts = {
        name for name in source_names if is_safe_notes_master_part(name)
    }
    blueprint_master_parts = {
        str(candidate_notes_page.get("sourceNotesMasterPart", ""))
        for candidate_slide in template_blueprint.get("slides", [])
        if isinstance(candidate_slide, dict)
        and isinstance(
            (candidate_notes_page := candidate_slide.get("notesPage")),
            dict,
        )
        and candidate_notes_page.get("sourceNotesMasterPart")
    }
    if (
        len(actual_master_parts) > 1
        or len(blueprint_master_parts) > 1
        or blueprint_master_parts - actual_master_parts
    ):
        return "NOTES_MASTER_CAPABILITY_UNSAFE"

    existing_presentation_master_relationships = relationship_type_nodes(
        presentation_rels_root,
        NOTES_MASTER_REL_TYPE,
    )

    # A notes master may have been created by an earlier update_speaker_notes
    # operation in this same sync, in which case it lives in added_entries (and
    # its theme/rels too) rather than package_entries. Read from both so a
    # second notes-less slide can reuse the freshly created master instead of
    # failing NOTES_MASTER_CAPABILITY_UNSAFE.
    def pending_part(part: str) -> bytes | None:
        if part in package_entries:
            return package_entries[part]
        return added_entries.get(part)

    if actual_master_parts:
        notes_master_part = next(iter(actual_master_parts))
        if blueprint_master_parts != {notes_master_part}:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        notes_master_xml = pending_part(notes_master_part)
        if notes_master_xml is None:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        try:
            ET.fromstring(notes_master_xml)
        except ET.ParseError:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        notes_master_rels_part = rels_part_for_slide_part(notes_master_part)
        notes_master_rels_xml = pending_part(notes_master_rels_part)
        if notes_master_rels_xml is None:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        try:
            notes_master_rels_root = ET.fromstring(notes_master_rels_xml)
        except ET.ParseError:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        theme_relationships = relationship_type_nodes(
            notes_master_rels_root,
            THEME_REL_TYPE,
        )
        theme_part = (
            resolve_relationship_part(
                notes_master_part,
                str(theme_relationships[0].get("Target", "")),
            )
            if len(theme_relationships) == 1
            and relationship_is_internal(theme_relationships[0])
            else ""
        )
        theme_xml = pending_part(theme_part) if theme_part else None
        if theme_xml is None:
            try:
                theme_xml = source_package.read(theme_part)
            except KeyError:
                return "NOTES_MASTER_CAPABILITY_UNSAFE"
        try:
            ET.fromstring(theme_xml)
        except ET.ParseError:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        if (
            len(existing_presentation_master_relationships) != 1
            or not relationship_is_internal(
                existing_presentation_master_relationships[0]
            )
            or resolve_relationship_part(
                "ppt/presentation.xml",
                str(existing_presentation_master_relationships[0].get("Target", "")),
            )
            != notes_master_part
            or not is_safe_theme_part(theme_part)
            or theme_part not in source_names
            or not content_type_override_matches(
                content_types_root,
                notes_master_part,
                NOTES_MASTER_CONTENT_TYPE,
            )
            or not content_type_override_matches(
                content_types_root,
                theme_part,
                THEME_CONTENT_TYPE,
            )
        ):
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
    else:
        if blueprint_master_parts or existing_presentation_master_relationships:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        notes_master_part = next_indexed_part(
            "ppt/notesMasters/notesMaster",
            ".xml",
            source_names,
        )
        theme_part = next_indexed_part(
            "ppt/theme/theme",
            ".xml",
            source_names,
        )
        if not notes_master_part or not theme_part:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        try:
            notes_master_rels_root = ET.fromstring(template["notesMasterRels"])
        except ET.ParseError:
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        theme_relationships = relationship_type_nodes(
            notes_master_rels_root,
            THEME_REL_TYPE,
        )
        if len(theme_relationships) != 1 or not relationship_is_internal(
            theme_relationships[0]
        ):
            return "NOTES_MASTER_CAPABILITY_UNSAFE"
        theme_relationships[0].set(
            "Target",
            posixpath.relpath(theme_part, posixpath.dirname(notes_master_part)),
        )
        added_entries[notes_master_part] = template["notesMaster"]
        added_entries[rels_part_for_slide_part(notes_master_part)] = xml_bytes(
            notes_master_rels_root
        )
        added_entries[theme_part] = template["theme"]
        ET.SubElement(
            presentation_rels_root,
            f"{{{PKG_REL_NS}}}Relationship",
            {
                "Id": next_relationship_id(presentation_rels_root),
                "Type": NOTES_MASTER_REL_TYPE,
                "Target": posixpath.relpath(notes_master_part, "ppt"),
            },
        )
        ensure_content_type_override(
            content_types_root,
            notes_master_part,
            NOTES_MASTER_CONTENT_TYPE,
        )
        ensure_content_type_override(
            content_types_root,
            theme_part,
            THEME_CONTENT_TYPE,
        )
        source_names.update(
            {
                notes_master_part,
                rels_part_for_slide_part(notes_master_part),
                theme_part,
            }
        )

    notes_part = next_indexed_part(
        "ppt/notesSlides/notesSlide",
        ".xml",
        source_names,
    )
    if not notes_part:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    try:
        notes_root = ET.fromstring(template["notesSlide"])
        notes_rels_root = ET.fromstring(template["notesSlideRels"])
    except ET.ParseError:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    body_shapes = notes_body_shapes(notes_root)
    if len(body_shapes) != 1:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    body_shape_id = notes_shape_id(body_shapes[0])
    text_body = first_local_child(body_shapes[0], "txBody")
    if not body_shape_id or text_body is None:
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    notes_master_relationships = relationship_type_nodes(
        notes_rels_root,
        NOTES_MASTER_REL_TYPE,
    )
    slide_relationships = relationship_type_nodes(notes_rels_root, SLIDE_REL_TYPE)
    if (
        len(notes_master_relationships) != 1
        or len(slide_relationships) != 1
        or not relationship_is_internal(notes_master_relationships[0])
        or not relationship_is_internal(slide_relationships[0])
    ):
        return "NOTES_MASTER_CAPABILITY_UNSAFE"
    notes_master_relationships[0].set(
        "Target",
        posixpath.relpath(notes_master_part, posixpath.dirname(notes_part)),
    )
    slide_relationships[0].set(
        "Target",
        posixpath.relpath(slide_part, posixpath.dirname(notes_part)),
    )
    replace_notes_text_body(text_body, speaker_notes)

    ET.SubElement(
        slide_rels_root,
        f"{{{PKG_REL_NS}}}Relationship",
        {
            "Id": next_relationship_id(slide_rels_root),
            "Type": NOTES_SLIDE_REL_TYPE,
            "Target": posixpath.relpath(notes_part, posixpath.dirname(slide_part)),
        },
    )
    ensure_content_type_override(
        content_types_root,
        notes_part,
        NOTES_SLIDE_CONTENT_TYPE,
    )
    added_entries[notes_part] = xml_bytes(notes_root)
    added_entries[rels_part_for_slide_part(notes_part)] = xml_bytes(notes_rels_root)
    package_entries[slide_rels_part] = xml_bytes(slide_rels_root)
    package_entries["ppt/_rels/presentation.xml.rels"] = xml_bytes(
        presentation_rels_root
    )
    package_entries["[Content_Types].xml"] = xml_bytes(content_types_root)
    slide["notesPage"] = {
        "status": "preserved",
        "sourceNotesPart": notes_part,
        "sourceNotesMasterPart": notes_master_part,
        "bodyShapeId": body_shape_id,
        "bodyWritable": True,
        "notesWidthEmu": notes_width_emu,
        "notesHeightEmu": notes_height_emu,
        "hasNonBodyContent": False,
    }
    return None


@lru_cache(maxsize=1)
def minimal_notes_package_template() -> dict[str, bytes] | None:
    from app.ai.pptx_ooxml.operations import (
        resolve_relationship_part,
    )

    try:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.notes_slide.notes_text_frame.text = ""
        buffer = BytesIO()
        presentation.save(buffer)
        with zipfile.ZipFile(BytesIO(buffer.getvalue()), "r") as package:
            notes_master_rels = package.read(
                "ppt/notesMasters/_rels/notesMaster1.xml.rels"
            )
            notes_master_rels_root = ET.fromstring(notes_master_rels)
            theme_relationships = relationship_type_nodes(
                notes_master_rels_root,
                THEME_REL_TYPE,
            )
            if len(theme_relationships) != 1 or not relationship_is_internal(
                theme_relationships[0]
            ):
                return None
            theme_part = resolve_relationship_part(
                "ppt/notesMasters/notesMaster1.xml",
                str(theme_relationships[0].get("Target", "")),
            )
            template = {
                "notesSlide": package.read("ppt/notesSlides/notesSlide1.xml"),
                "notesSlideRels": package.read(
                    "ppt/notesSlides/_rels/notesSlide1.xml.rels"
                ),
                "notesMaster": package.read("ppt/notesMasters/notesMaster1.xml"),
                "notesMasterRels": notes_master_rels,
                "theme": package.read(theme_part),
            }
        if len(notes_body_shapes(ET.fromstring(template["notesSlide"]))) != 1:
            return None
        ET.fromstring(template["notesMaster"])
        ET.fromstring(template["theme"])
        return template
    except (KeyError, OSError, ET.ParseError, zipfile.BadZipFile):
        return None


def notes_body_shapes(root: ET.Element[Any]) -> list[ET.Element[Any]]:
    from app.ai.pptx_ooxml.routing import (
        direct_local_children,
        first_local_child,
    )

    common_slide = first_local_child(root, "cSld")
    shape_tree = (
        first_local_child(common_slide, "spTree") if common_slide is not None else None
    )
    if shape_tree is None:
        return []
    return [
        shape
        for shape in direct_local_children(shape_tree, "sp")
        if notes_placeholder_type(shape) == "body"
    ]


def relationship_type_nodes(
    root: ET.Element[Any],
    relationship_type: str,
) -> list[ET.Element[Any]]:
    return [
        relationship
        for relationship in list(root)
        if relationship.get("Type") == relationship_type
    ]


def relationship_is_internal(relationship: ET.Element[Any]) -> bool:
    return relationship.get("TargetMode") in {None, "Internal"}


def content_type_override_matches(
    root: ET.Element[Any],
    part: str,
    content_type: str,
) -> bool:
    matches = [
        item
        for item in list(root)
        if item.tag == f"{{{CONTENT_TYPES_NS}}}Override"
        and item.get("PartName") == f"/{part}"
    ]
    return len(matches) == 1 and matches[0].get("ContentType") == content_type


def ensure_content_type_override(
    root: ET.Element[Any],
    part: str,
    content_type: str,
) -> None:
    if content_type_override_matches(root, part, content_type):
        return
    ET.SubElement(
        root,
        f"{{{CONTENT_TYPES_NS}}}Override",
        {"PartName": f"/{part}", "ContentType": content_type},
    )


def next_indexed_part(prefix: str, suffix: str, names: set[str]) -> str:
    for index in range(1, 10_001):
        candidate = f"{prefix}{index}{suffix}"
        if candidate not in names:
            return candidate
    return ""


def notes_placeholder_type(shape: ET.Element[Any]) -> str:
    from app.ai.pptx_ooxml.routing import (
        first_local_child,
    )

    non_visual = first_local_child(shape, "nvSpPr")
    non_visual_properties = (
        first_local_child(non_visual, "nvPr") if non_visual is not None else None
    )
    placeholder = (
        first_local_child(non_visual_properties, "ph")
        if non_visual_properties is not None
        else None
    )
    return str(placeholder.get("type", "")) if placeholder is not None else ""


def notes_shape_id(shape: ET.Element[Any]) -> str:
    from app.ai.pptx_ooxml.routing import (
        first_local_child,
    )

    non_visual = first_local_child(shape, "nvSpPr")
    properties = (
        first_local_child(non_visual, "cNvPr") if non_visual is not None else None
    )
    return str(properties.get("id", "")) if properties is not None else ""


def notes_text_body_text(text_body: ET.Element[Any]) -> str:
    from app.ai.pptx_ooxml.routing import (
        direct_local_children,
    )

    return "\n".join(
        notes_paragraph_text(paragraph)
        for paragraph in direct_local_children(text_body, "p")
    )


def notes_paragraph_text(paragraph: ET.Element[Any]) -> str:
    from app.ai.pptx_ooxml.routing import (
        local_name,
    )

    fragments: list[str] = []
    for item in paragraph.iter():
        name = local_name(item)
        if name == "t":
            fragments.append(str(item.text or ""))
        elif name == "br":
            fragments.append("\n")
        elif name == "tab":
            fragments.append("\t")
    return "".join(fragments)


def replace_notes_text_body(text_body: ET.Element[Any], speaker_notes: str) -> None:
    from app.ai.pptx_ooxml.routing import (
        direct_local_children,
    )

    existing = direct_local_children(text_body, "p")
    existing_text = [notes_paragraph_text(paragraph) for paragraph in existing]
    desired_text = speaker_notes.split("\n")
    template_indexes: dict[int, int] = {}
    matcher = difflib.SequenceMatcher(a=existing_text, b=desired_text, autojunk=False)
    for tag, old_start, old_end, new_start, new_end in matcher.get_opcodes():
        if tag == "equal":
            for offset in range(new_end - new_start):
                template_indexes[new_start + offset] = old_start + offset
        elif tag == "replace":
            for offset in range(new_end - new_start):
                if old_start < old_end:
                    template_indexes[new_start + offset] = min(
                        old_start + offset,
                        old_end - 1,
                    )

    new_paragraphs: list[ET.Element[Any]] = []
    for index, text in enumerate(desired_text):
        template_index = template_indexes.get(index)
        if template_index is None and existing:
            template_index = min(index, len(existing) - 1)
        template = existing[template_index] if template_index is not None else None
        if (
            template is not None
            and template_index is not None
            and existing_text[template_index] == text
        ):
            new_paragraphs.append(copy.deepcopy(template))
        else:
            new_paragraphs.append(notes_paragraph_from_template(template, text))

    for paragraph in existing:
        text_body.remove(paragraph)
    text_body.extend(new_paragraphs)


def notes_paragraph_from_template(
    template: ET.Element[Any] | None,
    text: str,
) -> ET.Element[Any]:
    from app.ai.pptx_ooxml.routing import (
        first_local_child,
        local_name,
    )

    paragraph = (
        copy.deepcopy(template)
        if template is not None
        else ET.Element(f"{{{DML_NS}}}p")
    )
    style = None
    if template is not None:
        for child in list(paragraph):
            if local_name(child) not in {"pPr", "endParaRPr"}:
                paragraph.remove(child)
        for child in list(template):
            if local_name(child) in {"r", "fld"}:
                style = first_local_child(child, "rPr")
                if style is not None:
                    break
    if not text:
        return paragraph

    run = ET.Element(f"{{{DML_NS}}}r")
    if style is not None:
        run.append(copy.deepcopy(style))
    text_element = ET.SubElement(run, f"{{{DML_NS}}}t")
    text_element.text = text
    if text[:1].isspace() or text[-1:].isspace():
        text_element.set(XML_SPACE, "preserve")
    end_index = next(
        (
            index
            for index, child in enumerate(list(paragraph))
            if local_name(child) == "endParaRPr"
        ),
        len(paragraph),
    )
    paragraph.insert(end_index, run)
    return paragraph
